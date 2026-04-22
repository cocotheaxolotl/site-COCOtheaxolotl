import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

path = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
print(f"Exists: {os.path.exists(path)}")
prs = Presentation(path)
print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inches")
print()

for i, slide in enumerate(prs.slides):
    texts = []
    images = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        try:
            if shape.image:
                images.append(f"{shape.image.content_type}, {len(shape.image.blob)//1024}KB")
        except:
            pass

    text_summary = " | ".join(texts[:3]) if texts else "(no text)"
    img_summary = ", ".join(images) if images else "(no images)"
    print(f"Slide {i+1:2d}: {len(images)} img, {len(texts)} txt | {text_summary[:80]} | {img_summary[:60]}")
