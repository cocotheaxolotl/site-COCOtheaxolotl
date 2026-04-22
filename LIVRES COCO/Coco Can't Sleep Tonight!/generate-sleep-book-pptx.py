"""
Generate "Coco ne dort pas ce soir !" — French spread layout
Format: 8.5 x 11 trim + 0.125" bleed = 8.75 x 11.25 per page
Layout: text page LEFT, full-page illustration RIGHT
Target: 30 pages (1 spread per animal)
KDP bleed-ready
"""

import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

# ── Page dimensions with KDP bleed ──
PAGE_W = Inches(8.75)
PAGE_H = Inches(11.25)
BLEED  = Inches(0.125)
SAFE_MARGIN = Inches(0.25)
GUTTER = Inches(0.25)
TEXT_WIDTH  = PAGE_W - SAFE_MARGIN - GUTTER

# ── Colors ──
SOFT_BLUE    = RGBColor(0x3D, 0x3D, 0x8E)
DARK_PINK    = RGBColor(0xD6, 0x00, 0x6E)
BLACK        = RGBColor(0x33, 0x33, 0x33)
GREY         = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_BG     = RGBColor(0xF8, 0xF8, 0xFF)
FACT_BG      = RGBColor(0xE8, 0xEA, 0xF6)
FACT_BORDER  = RGBColor(0x7B, 0x7F, 0xC4)
FACT_TEXT     = RGBColor(0x3D, 0x3D, 0x8E)
NAVY         = RGBColor(0x00, 0x00, 0x80)
PLAC_BG      = RGBColor(0xF5, 0xF5, 0xF5)
PLAC_BORDER  = RGBColor(0xDD, 0xDD, 0xDD)

FONT = "Berlin Sans FB Demi"
a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def set_rpr(paragraph, size_pt, bold, color, font_name=FONT):
    pPr = paragraph._element.find(f"{a_ns}pPr")
    if pPr is None:
        pPr = etree.SubElement(paragraph._element, f"{a_ns}pPr")
        paragraph._element.insert(0, pPr)
    defRPr = etree.SubElement(pPr, f"{a_ns}defRPr")
    defRPr.set("sz", str(int(size_pt * 100)))
    defRPr.set("b", "1" if bold else "0")
    sf = etree.SubElement(defRPr, f"{a_ns}solidFill")
    sc = etree.SubElement(sf, f"{a_ns}srgbClr")
    sc.set("val", str(color))
    la = etree.SubElement(defRPr, f"{a_ns}latin")
    la.set("typeface", font_name)


def _make_run(parent, text, sz, bold, color, italic=False, font_name=FONT):
    """Create a <a:r> run — formatting inherited from defRPr, modifiable dans PowerPoint."""
    r = etree.SubElement(parent, f"{a_ns}r")
    etree.SubElement(r, f"{a_ns}t").text = text


def tbox(slide, left, top, width, height, text, sz, bold, color,
         align=PP_ALIGN.CENTER, vcenter=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(f"{a_ns}bodyPr")
    bp.set("wrap", "square")
    if vcenter:
        bp.set("anchor", "ctr")
    else:
        etree.SubElement(bp, f"{a_ns}spAutoFit")
    p = tf.paragraphs[0]
    p.alignment = align
    set_rpr(p, sz, bold, color)

    for li, line in enumerate(text.split("\n")):
        if li > 0:
            etree.SubElement(p._element, f"{a_ns}br")
        _make_run(p._element, line, sz, bold, color)
    return txBox


def fact_box(slide, left, top, width, height, title, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = FACT_BG
    shape.line.color.rgb = FACT_BORDER
    shape.line.width = Pt(2)
    fe = shape._element.find(f".//{a_ns}solidFill")
    if fe is not None:
        ce = fe.find(f"{a_ns}srgbClr")
        if ce is not None:
            etree.SubElement(ce, f"{a_ns}alpha").set("val", "90000")
    tf = shape.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(f"{a_ns}bodyPr")
    bp.set("wrap", "square")
    bp.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_rpr(p, 24, False, NAVY)
    # Bold title
    r = etree.SubElement(p._element, f"{a_ns}r")
    rP = etree.SubElement(r, f"{a_ns}rPr")
    rP.set("sz", "2400"); rP.set("b", "1")
    sf = etree.SubElement(rP, f"{a_ns}solidFill")
    sc = etree.SubElement(sf, f"{a_ns}srgbClr"); sc.set("val", str(NAVY))
    la = etree.SubElement(rP, f"{a_ns}latin"); la.set("typeface", FONT)
    etree.SubElement(r, f"{a_ns}t").text = title
    # Break + text
    etree.SubElement(p._element, f"{a_ns}br")
    r2 = etree.SubElement(p._element, f"{a_ns}r")
    rP2 = etree.SubElement(r2, f"{a_ns}rPr")
    rP2.set("sz", "2400"); rP2.set("b", "0")
    sf2 = etree.SubElement(rP2, f"{a_ns}solidFill")
    sc2 = etree.SubElement(sf2, f"{a_ns}srgbClr"); sc2.set("val", str(NAVY))
    la2 = etree.SubElement(rP2, f"{a_ns}latin"); la2.set("typeface", FONT)
    etree.SubElement(r2, f"{a_ns}t").text = text


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def img_slide(prs, label="[ Illustration ]", img_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if img_path and Path(img_path).exists():
        # Insert actual image — full page with bleed
        slide.shapes.add_picture(str(img_path), 0, 0, PAGE_W, PAGE_H)
    else:
        # Placeholder rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            BLEED, BLEED, PAGE_W - 2*BLEED, PAGE_H - 2*BLEED)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PLAC_BG
        shape.line.color.rgb = PLAC_BORDER
        shape.line.width = Pt(1)
        shape.line.dash_style = 4
        tf = shape.text_frame; tf.word_wrap = True
        bp = tf._txBody.find(f"{a_ns}bodyPr"); bp.set("anchor", "ctr")
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_rpr(p, 14, False, PLAC_BORDER)
        r = etree.SubElement(p._element, f"{a_ns}r")
        etree.SubElement(r, f"{a_ns}t").text = label
    return slide


# ═══════════════════════════════════════
IMG_DIR = Path(__file__).parent
BG_IMAGE = IMG_DIR / "fond_nuit_étoilee.png"


def txt_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if BG_IMAGE.exists():
        slide.shapes.add_picture(str(BG_IMAGE), 0, 0, PAGE_W, PAGE_H)
    else:
        set_bg(slide, LIGHT_BG)
    return slide

# Image file mapping (prompt # → filename)
IMAGES = {
    'cover':        IMG_DIR / "couverture en FR.png",
    'paupieres':    IMG_DIR / "COCO NE DORT PAS + PAUPIÈRES.png",
    'eclisse':      IMG_DIR / "COCO S'ÉCLIPSE.png",
    'hibou':        IMG_DIR / "le-hibou.png",
    'dauphin':      IMG_DIR / "le_dauphin.png",
    'loutre':       IMG_DIR / "les-loutres.png",
    'koala':        IMG_DIR / "koala.png",
    'chauve_souris':IMG_DIR / "chauve-souris.png",
    'chat':         IMG_DIR / "le-chat.png",
    'cheval':       IMG_DIR / "le-cheval.png",
    'flamant':      IMG_DIR / "flamant-rose.png",
    'paresseux':    IMG_DIR / "paresseux.png",
    'maman':        IMG_DIR / "MAMAN AXOLOTL.png",
    'borde':        IMG_DIR / "COCO BORDÉ + RÉGÉNÉRATION.png",
    'bonne_nuit':   IMG_DIR / "BONNE NUIT, COCO.png",
    'back_cover':   IMG_DIR / "COUVERTURE ARRIÈRE.png",
}


def build():
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    # ══════════════════════════════════
    # SLIDE 1 — COVER
    # ══════════════════════════════════
    s = img_slide(prs, "[ COUVERTURE ]", IMAGES['cover'])
    tbox(s, SAFE_MARGIN, Inches(7.0), TEXT_WIDTH, Inches(1.5),
         "Coco ne dort pas ce soir !", 40, True, DARK_PINK)
    tbox(s, SAFE_MARGIN, Inches(8.5), TEXT_WIDTH, Inches(0.7),
         "Une aventure de Coco l\u2019Axolotl", 20, False, SOFT_BLUE)
    tbox(s, SAFE_MARGIN, Inches(9.3), TEXT_WIDTH, Inches(0.5),
         "Dr. Anita NIRVENA  \u2022  Illustrations : Loopinky", 13, False, GREY)

    # ══════════════════════════════════
    # SLIDE 2 — COPYRIGHT
    # ══════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tbox(s, GUTTER, Inches(3.0), TEXT_WIDTH, Inches(6),
        "Coco ne dort pas ce soir !\nUne aventure de Coco l\u2019Axolotl\n\n"
        "Par : Dr. Anita NIRVENA\nIllustrations : Loopinky\n\n"
        "\u00a9 2026 Dr. Anita NIRVENA\nTous droits r\u00e9serv\u00e9s.\n\n"
        "www.cocotheaxolotl.org\nPremi\u00e8re \u00e9dition, 2026",
        12, False, GREY)

    # ══════════════════════════════════
    # SLIDES 3-4 — COCO CAN'T SLEEP
    # ══════════════════════════════════
    s = txt_page(prs)
    tbox(s, GUTTER, Inches(1.5), TEXT_WIDTH, Inches(8),
        "Ce soir, c\u2019\u00e9tait l\u2019heure de dormir. Mais Coco n\u2019avait pas du tout sommeil.\n\n"
        "\u2014Maman, je n\u2019arrive pas \u00e0 dormir !\n\n"
        "\u2014Ferme les yeux et...\n"
        "\u2014Mais Maman, je n\u2019ai pas de paupi\u00e8res !\n\n"
        "Maman ria.\n"
        "\u2014C\u2019est vrai ! Alors reste tranquille et le sommeil viendra.",
        28, True, BLACK, vcenter=True)
    img_slide(prs, "[ Paupières ]", IMAGES['paupieres'])

    # ══════════════════════════════════
    # SLIDES 5-6 — SNEAKS OUT
    # ══════════════════════════════════
    s = txt_page(prs)
    tbox(s, GUTTER, Inches(1.5), TEXT_WIDTH, Inches(8),
        "Mais Coco ne pouvait pas rester tranquille. Pas ce soir.\n\n"
        "Quand Maman et Papa se furent endormis, Coco se faufila doucement hors de la rivi\u00e8re.\n\n"
        "\u2014Si je ne dors pas, je vais d\u00e9couvrir ce que font les autres animaux la nuit !",
        28, True, BLACK, vcenter=True)
    img_slide(prs, "[ S'éclipse ]", IMAGES['eclisse'])

    # ═══════════════════════════════════
    # SLIDES 7-24 — 9 ANIMALS (1 spread each)
    # ═══════════════════════════════════

    animals = [
        # (emoji, name, story_text, fact_text, illust_label, img_key)

        ("\U0001f989", "Le Hibou",
         "Coco rencontra un hibou perch\u00e9 sur une branche.\n\n"
         "\u2014Bonsoir, Hibou ! Pourquoi tu ne dors pas ?\n\n"
         "\u2014Hou hou ! Je suis nocturne ! Je travaille la nuit et je dors le jour. Mais m\u00eame moi, j\u2019ai besoin de dormir. Tout le monde a besoin de dormir !",
         "Les hiboux dorment le jour, cach\u00e9s dans les arbres. Ils peuvent tourner la t\u00eate presque tout autour, comme pour regarder derri\u00e8re eux !",
         "Hibou sur branche", 'hibou'),

        ("\U0001f42c", "Le Dauphin",
         "Au bord de l\u2019oc\u00e9an, un dauphin nageait tranquillement.\n\n"
         "\u2014Tu ne dors pas non plus ?\n\n"
         "\u2014Oh si ! Je n\u2019endors qu\u2019une moiti\u00e9 de mon cerveau \u00e0 la fois. L\u2019autre reste \u00e9veill\u00e9e pour nager et respirer !\n\n"
         "\u2014Tu dors ET tu nages en m\u00eame temps ?!\n\n"
         "\u2014Le sommeil est si important que j\u2019ai trouv\u00e9 un moyen de ne jamais m\u2019en passer !",
         "Les dauphins dorment avec un seul \u0153il ferm\u00e9 ! Le c\u00f4t\u00e9 du cerveau qui dort change toutes les deux heures.",
         "Dauphin", 'dauphin'),

        ("\U0001f9a6", "La Loutre",
         "Deux loutres flottaient sur le dos, patte dans la patte.\n\n"
         "\u2014Pourquoi vous vous tenez la main ?\n\n"
         "\u2014Pour ne pas s\u2019\u00e9loigner en dormant ! Se sentir en s\u00e9curit\u00e9, c\u2019est tr\u00e8s important pour bien dormir.\n\n"
         "Coco sourit.\n"
         "\u2014C\u2019est comme quand je me blottis contre ma peluche !",
         "Les loutres de mer se tiennent vraiment par la main en dormant ! Elles s\u2019enroulent aussi dans des algues pour ne pas d\u00e9river.",
         "Loutres", 'loutre'),

        ("\U0001f428", "Le Koala",
         "Dans un eucalyptus, un koala dormait profond\u00e9ment.\n\n"
         "\u2014Koala ! Tu dors ?\n"
         "\u2014Zzz... oui... 22 heures par jour...\n\n"
         "\u201422 HEURES ?! Mais pourquoi autant ?\n\n"
         "\u2014Mon corps a besoin d\u2019\u00e9nergie pour dig\u00e9rer mes feuilles... Le sommeil, \u00e7a recharge ton corps en \u00e9nergie ! Et il se rendormit aussit\u00f4t.",
         "Les koalas dorment jusqu\u2019\u00e0 22 heures par jour ! Les feuilles d\u2019eucalyptus sont tr\u00e8s dures \u00e0 dig\u00e9rer et donnent tr\u00e8s peu d\u2019\u00e9nergie.",
         "Koala", 'koala'),

        ("\U0001f987", "La Chauve-Souris",
         "Sous un pont, des chauves-souris \u00e9taient suspendues la t\u00eate en bas.\n\n"
         "\u2014Vous dormez \u00e0 l\u2019envers ?!\n\n"
         "\u2014Bien s\u00fbr ! Et pendant qu\u2019on dort, notre cerveau trie tout ce qu\u2019on a appris \u2014 o\u00f9 sont les insectes, quels chemins \u00e9viter...\n\n"
         "\u2014C\u2019est pour \u00e7a qu\u2019on se sent plus malin apr\u00e8s une bonne nuit !",
         "Pendant le sommeil, le cerveau trie et range les souvenirs. C\u2019est pour \u00e7a qu\u2019on retient mieux ses le\u00e7ons apr\u00e8s avoir bien dormi !",
         "Chauve-souris", 'chauve_souris'),

        ("\U0001f431", "Le Chat",
         "Pr\u00e8s d\u2019une maison, un chat \u00e9tait roul\u00e9 en boule sur un muret.\n\n"
         "\u2014Tu fais un gros dodo ?\n\n"
         "\u2014Je fais plein de petits dodos. On appelle \u00e7a des siestes !\n\n"
         "\u2014Pour les petits axolotls, c\u2019est mieux de faire un long sommeil la nuit. Pendant le sommeil profond, ton corps se r\u00e9pare !",
         "Les chats dorment entre 12 et 16 heures par jour ! C\u2019est de l\u00e0 que vient le mot anglais \"catnap\".",
         "Chat", 'chat'),

        ("\U0001f40e", "Le Cheval",
         "Dans un pr\u00e9, un cheval se tenait immobile sous la lune.\n\n"
         "\u2014Tu... tu dors debout ?\n\n"
         "\u2014Oui ! Comme \u00e7a je suis pr\u00eat \u00e0 courir. Mais pour r\u00eaver, je dois m\u2019allonger.\n\n"
         "\u2014Le sommeil a diff\u00e9rentes \u00e9tapes : d\u2019abord l\u00e9ger, puis profond, puis les r\u00eaves ! Et \u00e7a recommence toute la nuit.",
         "Les chevaux ne dorment profond\u00e9ment que 2 \u00e0 3 heures par jour \u2014 et seulement allong\u00e9s !",
         "Cheval", 'cheval'),

        ("\U0001f9a9", "Le Flamant Rose",
         "Au bord d\u2019un \u00e9tang, un flamant rose dormait sur une seule patte !\n\n"
         "\u2014Comment tu ne tombes pas ?!\n\n"
         "\u2014L\u2019\u00e9quilibre, \u00e7a demande un corps bien repos\u00e9 ! Quand je suis fatigu\u00e9...\n\n"
         "Le flamant vacilla.\n"
         "\u2014...voil\u00e0 ce qui arrive !\n\n"
         "Coco ria, mais commen\u00e7ait \u00e0 comprendre.",
         "Les flamants dorment souvent sur une seule patte ! Les scientifiques pensent que c\u2019est pour garder l\u2019autre au chaud.",
         "Flamant", 'flamant'),

        ("\U0001f9a5", "Le Paresseux",
         "Coco commen\u00e7ait \u00e0 b\u00e2iller. Un paresseux pendait d\u2019une branche.\n\n"
         "\u2014Je commence... \u00e0 \u00eatre... fatigu\u00e9...\n\n"
         "\u2014Bienvenuuuue... au cluuuub... Quand on ne dort pas assez... on devient lent... grognon... et tout est plus difficile...\n\n"
         "\u2014Allez... va te coucher... petit axolotl... zzz...",
         "Les paresseux dorment environ 15 heures par jour ! Sans assez de sommeil, notre corps se sent lourd et notre cerveau a du mal \u00e0 se concentrer.",
         "Paresseux", 'paresseux'),
    ]

    for i, (emoji, name, story, fact, illust, img_key) in enumerate(animals):
        s = txt_page(prs)
        # Counter
        tbox(s, PAGE_W - Inches(1.5), Inches(0.3), Inches(1.0), Inches(0.3),
             f"{i+1} / 9", 10, False, GREY, align=PP_ALIGN.RIGHT)
        # Animal header
        tbox(s, GUTTER, Inches(0.7), TEXT_WIDTH, Inches(0.7),
             f"{emoji}  {name}", 28, True, SOFT_BLUE)
        # Story text
        tbox(s, GUTTER, Inches(1.6), TEXT_WIDTH, Inches(5.6),
             story, 28, True, BLACK, vcenter=True)
        # Fact box
        fact_box(s, SAFE_MARGIN, Inches(8.0), PAGE_W - 2*SAFE_MARGIN, Inches(2.5),
                 "Le savais-tu ?", fact)
        # Illustration
        img_slide(prs, f"[ {illust} ]", IMAGES[img_key])

    # ══════════════════════════════════
    # SLIDES 25-26 — MOMMY + AXOLOTL SLEEP
    # ══════════════════════════════════
    s = txt_page(prs)
    tbox(s, GUTTER, Inches(0.7), TEXT_WIDTH, Inches(0.7),
         "\U0001f30a  Maman Axolotl", 28, True, SOFT_BLUE)
    tbox(s, GUTTER, Inches(1.6), TEXT_WIDTH, Inches(5.6),
        "\u2014Coco ?\n\n"
        "C\u2019\u00e9tait Maman, au bord de la rivi\u00e8re.\n\n"
        "\u2014Maman ! J\u2019ai d\u00e9couvert comment tous les animaux dorment ! Tout le monde en a besoin. Moi aussi !\n\n"
        "\u2014Et sais-tu comment dorment les axolotls ? On se pose au fond de l\u2019eau, bien tranquilles. M\u00eame les yeux ouverts... on dort profond\u00e9ment !",
        28, True, BLACK, vcenter=True)
    fact_box(s, SAFE_MARGIN, Inches(8.0), PAGE_W - 2*SAFE_MARGIN, Inches(2.5),
             "Le savais-tu ?",
             "Les axolotls n\u2019ont pas de paupi\u00e8res ! Ils dorment les yeux grands ouverts, pos\u00e9s au fond de l\u2019eau.")
    img_slide(prs, "[ Maman ]", IMAGES['maman'])

    # ══════════════════════════════════
    # SLIDES 27-28 — TUCKED IN + REGENERATION
    # ══════════════════════════════════
    s = txt_page(prs)
    tbox(s, GUTTER, Inches(1.5), TEXT_WIDTH, Inches(5.8),
        "Maman porta Coco jusqu\u2019\u00e0 son coin douillet au fond de la rivi\u00e8re.\n\n"
        "\u2014Le sommeil, \u00e7a recharge l\u2019\u00e9nergie, \u00e7a range le cerveau, \u00e7a r\u00e9pare le corps... pas vrai ?\n\n"
        "\u2014Tout \u00e7a \u00e0 la fois, mon Coco. Et nous, les axolotls, on peut m\u00eame faire repousser nos pattes ! Mais pour \u00e7a, il faut bien dormir.",
        28, True, BLACK, vcenter=True)
    fact_box(s, SAFE_MARGIN, Inches(8.0), PAGE_W - 2*SAFE_MARGIN, Inches(2.5),
             "Le savais-tu ?",
             "Les axolotls ont un super-pouvoir : ils peuvent r\u00e9g\u00e9n\u00e9rer leurs pattes, leur queue et m\u00eame des parties de leur cerveau !")
    img_slide(prs, "[ Bordé ]", IMAGES['borde'])

    # ══════════════════════════════════
    # SLIDES 29-30 — GOODNIGHT
    # ══════════════════════════════════
    s = txt_page(prs)
    tbox(s, GUTTER, Inches(1.5), TEXT_WIDTH, Inches(5.5),
        "Coco se blottit dans son petit nid d\u2019eau douce avec sa peluche.\n\n"
        "Il pensa \u00e0 tous les animaux qui dormaient en ce moment \u2014 chacun \u00e0 sa fa\u00e7on, mais tous parce que le sommeil est magique.\n\n"
        "Ses branchies redevinrent toutes roses et touffues.\n\n"
        "Et Coco s\u2019endormit.",
        28, True, BLACK, vcenter=True)
    tbox(s, GUTTER, Inches(8.5), TEXT_WIDTH, Inches(1.5),
         "Bonne nuit, Coco.", 40, True, SOFT_BLUE)
    img_slide(prs, "[ Bonne nuit ]", IMAGES['bonne_nuit'])

    return prs


if __name__ == '__main__':
    import re as _re
    out_dir = Path(__file__).parent
    prs = build()
    base = "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads"
    # Auto-versioning: find highest existing V number and increment
    existing = list(out_dir.glob(f"{base}-V*.pptx"))
    if existing:
        nums = [int(m.group(1)) for f in existing
                if (m := _re.search(r'-V(\d+)\.pptx$', f.name))]
        next_v = max(nums) + 1 if nums else 1
    else:
        next_v = 1
    filename = f"{base}-V{next_v}.pptx"
    path = out_dir / filename
    prs.save(str(path))
    print(f"OK  {filename}")
    print(f"    Slides: {len(prs.slides)}")
    print(f"    Format: 8.75 x 11.25 (trim 8.5x11 + bleed)")
    print(f"    Layout: text LEFT, illustration RIGHT")
