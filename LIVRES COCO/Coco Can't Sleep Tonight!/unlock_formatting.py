"""
Retire les verrous de formatage (rPr) des runs dans un PPTX existant.
Ne touche a RIEN d'autre : ni images, ni positions, ni tailles, ni texte.
Sauvegarde en nouvelle version.
"""
from pathlib import Path
from pptx import Presentation
from lxml import etree

a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

src = Path(__file__).parent / "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads-V5.pptx"
dst = Path(__file__).parent / "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads-V6.pptx"

prs = Presentation(str(src))

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para._element.findall(f"{a_ns}r"):
                rPr = run.find(f"{a_ns}rPr")
                if rPr is not None:
                    run.remove(rPr)
                    count += 1

prs.save(str(dst))
print(f"OK  {dst.name}")
print(f"    {count} verrous rPr retires")
print(f"    Images, positions, tailles : intacts")
