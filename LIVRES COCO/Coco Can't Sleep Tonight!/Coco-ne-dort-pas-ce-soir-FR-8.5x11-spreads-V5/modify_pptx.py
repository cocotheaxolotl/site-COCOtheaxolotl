"""
Modify the interior PPTX:
1. Insert a glossary slide at position 2 (after CHER PARENT)
2. Add pretty page numbers to story pages
3. Add emojis to titles that are missing them
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

INPUT = 'Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads-V5.pptx'
OUTPUT = 'Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads-V5-modified.pptx'

prs = Presentation(INPUT)
slides = prs.slides

# ============================================================
# STEP 3: Add emojis to titles missing them (do this BEFORE inserting slide)
# ============================================================
title_emojis = {
    2: ("L\u2019heure de dormir", "\U0001f319 L\u2019heure de dormir"),
    4: ("L\u2019aventure commence", "\U0001f31f L\u2019aventure commence"),
    26: ("Le super-pouvoir", "\u2728 Le super-pouvoir"),
    28: ("Bonne nuit Coco !", "\U0001f4a4 Bonne nuit Coco !"),
}

for slide_idx, (old_text, new_text) in title_emojis.items():
    slide = slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if full == old_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = new_text
                            print(f"  Slide {slide_idx+1}: '{old_text}' -> '{new_text}'")
                            break
                    break
                break

print("Emojis added to titles.")

# ============================================================
# STEP 2: Add page numbers to story pages
# ============================================================
# Story pages are slides 3-30 (index 2-29), text on odd indices (2,4,6,...28)
# Images on even indices (3,5,7,...29)
# Page numbering: start from 1 on slide index 2

PAGE_NUM_COLOR = RGBColor(0x69, 0x5A, 0x4F)  # warm brown

page_number = 1
for slide_idx in range(2, 30):  # slides 3-30
    slide = slides[slide_idx]

    # Outer corner: odd pages (1,3,5...) = right, even pages (2,4,6...) = left
    if page_number % 2 == 1:  # recto -> right
        left = Inches(7.2)
        align = PP_ALIGN.RIGHT
    else:  # verso -> left
        left = Inches(0.3)
        align = PP_ALIGN.LEFT

    top = Inches(10.25)
    width = Inches(1.0)
    height = Inches(0.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    run.text = f"\u2014 {page_number} \u2014"
    run.font.name = "Berlin Sans FB Demi"
    run.font.size = Pt(14)
    run.font.color.rgb = PAGE_NUM_COLOR
    run.font.bold = True

    page_number += 1

print(f"Page numbers added (1-{page_number-1}).")

# ============================================================
# STEP 1: Insert TABLE OF CONTENTS slide at position 2
# ============================================================

blank_layout = prs.slide_layouts[6]  # blank
toc_slide = prs.slides.add_slide(blank_layout)

# Background
background = toc_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xEC)  # warm cream

# Title
title_box = toc_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.5), Inches(0.9))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\U0001f4d6 Table des mati\u00e8res"
run.font.name = "Berlin Sans FB Demi"
run.font.size = Pt(32)
run.font.color.rgb = RGBColor(0x2D, 0x1B, 0x4E)
run.font.bold = True

# TOC entries: (emoji+title, page_number)
toc_entries = [
    ("\U0001f319 L\u2019heure de dormir", 1),
    ("\U0001f31f L\u2019aventure commence", 3),
    ("\U0001f989 Le Hibou", 5),
    ("\U0001f42c Le Dauphin", 7),
    ("\U0001f9a6 La Loutre", 9),
    ("\U0001f428 Le Koala", 11),
    ("\U0001f987 La Chauve-Souris", 13),
    ("\U0001f431 Le Chat", 15),
    ("\U0001f40e Le Cheval", 17),
    ("\U0001f9a9 Le Flamant Rose", 19),
    ("\U0001f9a5 Le Paresseux", 21),
    ("\U0001f30a Maman Axolotl", 23),
    ("\u2728 Le super-pouvoir", 25),
    ("\U0001f4a4 Bonne nuit Coco !", 27),
]

# Content box
content_box = toc_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.9), Inches(8.5))
tf = content_box.text_frame
tf.word_wrap = True

# Alternating colors for entries
TOC_COLORS = [
    RGBColor(0x2D, 0x1B, 0x4E),  # deep purple
    RGBColor(0xA4, 0x79, 0x83),  # mauve/pink
]

for idx, (title, page) in enumerate(toc_entries):
    p_entry = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
    p_entry.space_before = Pt(6)
    p_entry.space_after = Pt(6)

    color = TOC_COLORS[idx % 2]

    # Title run
    run_title = p_entry.add_run()
    run_title.text = title
    run_title.font.name = "Berlin Sans FB Demi"
    run_title.font.size = Pt(16)
    run_title.font.color.rgb = color
    run_title.font.bold = True

    # Dots + page number
    dots = " " + "\u00b7" * 20 + " "
    run_dots = p_entry.add_run()
    run_dots.text = dots
    run_dots.font.name = "Berlin Sans FB Demi"
    run_dots.font.size = Pt(10)
    run_dots.font.color.rgb = RGBColor(0xC0, 0xB0, 0xA0)

    run_page = p_entry.add_run()
    run_page.text = str(page)
    run_page.font.name = "Berlin Sans FB Demi"
    run_page.font.size = Pt(16)
    run_page.font.color.rgb = color
    run_page.font.bold = True

print("Table of contents slide created.")

# Move TOC slide (last) to position 2 (index 1)
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)
toc_elem = slide_ids[-1]
slide_list.remove(toc_elem)
slide_list.insert(1, toc_elem)

print("TOC moved to position 2.")

# ============================================================
# Save
# ============================================================
prs.save(OUTPUT)
print(f"\nSaved to: {OUTPUT}")
print("Done!")
