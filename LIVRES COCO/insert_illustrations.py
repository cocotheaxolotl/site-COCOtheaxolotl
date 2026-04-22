"""Insert real illustrations into PowerPoint book files.
Replaces placeholder rectangles with actual images.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_DIR = Path(__file__).parent
COLORED_DIR = SCRIPT_DIR / "whose-egg-colored"

# Slide index → image filename (without path)
# Slide 0 = cover, Slide 1 = intro, Slides 2-30 = animals, Slide 31 = ending, Slide 32 = back cover
SLIDE_TO_IMAGE = {
    # 0: cover - not yet available
    1: "coco découvre la coquille.png",       # intro
    2: "axolotl-kitten-colored.png",          # kitten
    3: "axolotl_puffy-colored.png",           # puppy
    4: "axolotl-rabbit-colored.png",          # rabbit
    5: "axolotl-unicorn-colored.png",         # unicorn
    6: "axolotl-panda-colored.png",           # panda
    7: "axolotl-dolphin-colored.png",         # dolphin
    8: "axolotl-foal-colored.png",            # foal
    9: "axolotl-butterfly-colored.png",       # butterfly
    10: "axolotl-dinosaur-colored.png",       # dinosaur
    11: "axolotl-lion-colored.png",           # lion
    12: "axolotl-bear-colored.png",           # bear
    13: "axolotl-elephant-colored.png",       # elephant
    14: "axolotl-giraffe-colored.png",        # giraffe
    15: "axolotl-penguin-colored.png",        # penguin
    16: "axolotl-koala-colored.png",          # koala
    17: "axolotl-fox-colored.png",            # fox
    18: "axolotl-owl-colored.png",            # owl
    19: "axolotl-turtle-colored.png",         # turtle
    20: "axolotl-frog-colored.png",           # frog
    21: "axolotl-highland-colored.png",       # highland cow
    22: "axolotl-flamingo-colored.png",       # flamingo
    # 23-30: not yet available (sloth → dragon)
    # 31: ending - not yet available
    # 32: back cover - not yet available
}

PPTX_FILES = [
    "Whose-Egg-Is-This-EN.pptx",
    "Whose-Egg-Is-This-FR.pptx",
    "Whose-Egg-Is-This-ES.pptx",
    "Whose-Egg-Is-This-EN-8.5x11.pptx",
    "Whose-Egg-Is-This-FR-8.5x11.pptx",
    "Whose-Egg-Is-This-ES-8.5x11.pptx",
]


def replace_placeholder_with_image(slide, placeholder_name, image_path, prs):
    """Replace a rounded rectangle placeholder with an image."""
    # Find the placeholder shape
    target = None
    for shape in slide.shapes:
        if shape.name == placeholder_name and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            target = shape
            break

    if not target:
        # Try any auto_shape that contains "[ " text (placeholder indicator)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                if "[ " in shape.text_frame.text or "[" in shape.text_frame.text:
                    target = shape
                    break

    if not target:
        return False

    # Get position and size of the placeholder
    left = target.left
    top = target.top
    width = target.width
    height = target.height

    # Remove the placeholder shape
    sp = target._element
    sp.getparent().remove(sp)

    # Add the image in its place, maintaining aspect ratio
    # Images are 1024x1536 (portrait 2:3), fit within the placeholder box
    img_ratio = 1024 / 1536  # ~0.667
    box_ratio = width / height

    if img_ratio > box_ratio:
        # Image is wider relative to box - fit to width
        new_width = width
        new_height = int(width / img_ratio)
        new_left = left
        new_top = top + (height - new_height) // 2
    else:
        # Image is taller relative to box - fit to height
        new_height = height
        new_width = int(height * img_ratio)
        new_left = left + (width - new_width) // 2
        new_top = top

    slide.shapes.add_picture(str(image_path), new_left, new_top, new_width, new_height)
    return True


def process_pptx(pptx_path):
    """Process a single PowerPoint file."""
    prs = Presentation(str(pptx_path))
    inserted = 0
    skipped = 0

    for slide_idx, image_name in SLIDE_TO_IMAGE.items():
        if slide_idx >= len(prs.slides):
            continue

        image_path = COLORED_DIR / image_name
        if not image_path.exists():
            print(f"  WARNING: {image_name} not found, skipping slide {slide_idx}")
            skipped += 1
            continue

        slide = prs.slides[slide_idx]

        # Try common placeholder names
        replaced = False
        for name in ["Rounded Rectangle 1", "Rounded Rectangle 2"]:
            if replace_placeholder_with_image(slide, name, image_path, prs):
                replaced = True
                break

        if not replaced:
            # Try any auto_shape placeholder
            replaced = replace_placeholder_with_image(slide, None, image_path, prs)

        if replaced:
            inserted += 1
        else:
            print(f"  WARNING: No placeholder found on slide {slide_idx}")
            skipped += 1

    prs.save(str(pptx_path))
    return inserted, skipped


def main():
    print("Inserting illustrations into PowerPoint files...\n")

    for filename in PPTX_FILES:
        filepath = SCRIPT_DIR / filename
        if not filepath.exists():
            print(f"SKIP: {filename} not found")
            continue

        print(f"Processing {filename}...")
        inserted, skipped = process_pptx(filepath)
        print(f"  OK: {inserted} images inserted, {skipped} skipped\n")

    print("Done!")


if __name__ == "__main__":
    main()
