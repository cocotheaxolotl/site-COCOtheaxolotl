"""
Generate "Whose Egg Is This?" book as PowerPoint (EN / FR / ES)
Format: 8.5 x 11 inches — KDP coloring book standard (trim 8.25 x 11)
Each slide = one page of the book
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Dimensions ── 8.5 x 11 (KDP trim = 8.25 x 11)
W = Inches(8.5)
H = Inches(11)

# ── Colors ──
PINK      = RGBColor(0xFF, 0x69, 0xB4)
DARK_PINK = RGBColor(0xD6, 0x00, 0x6E)
LIGHT_BG  = RGBColor(0xFF, 0xF0, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x33, 0x33, 0x33)
GREY      = RGBColor(0x99, 0x99, 0x99)
FACT_BG   = RGBColor(0xFF, 0xF0, 0xF5)
FACT_BORDER = RGBColor(0xFF, 0xB6, 0xC1)
PLAC_BORDER = RGBColor(0xFF, 0xB6, 0xC1)

FONT = 'Baloo 2'
FONT_FALLBACK = 'Comic Sans MS'

def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=BLACK, align=PP_ALIGN.CENTER, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_placeholder(slide, left, top, width, height, label="[ Illustration ]"):
    """Dashed-border placeholder for illustration"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF9, 0xFB)
    shape.line.color.rgb = PLAC_BORDER
    shape.line.width = Pt(2)
    shape.line.dash_style = 4  # dash
    # Label inside
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = PLAC_BORDER
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(40)
    return shape

def add_fun_fact(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = FACT_BG
    shape.line.color.rgb = FACT_BORDER
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x8B, 0x00, 0x4A)
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return shape

# ── Book content (all 3 languages) ──
# Each animal: (emoji, name_en, name_fr, name_es,
#               illust_en, illust_fr, illust_es,
#               story_en, story_fr, story_es,
#               fact_en, fact_fr, fact_es)

animals = [
    ("🐱", "Baby Kitten", "Bébé Chaton", "Gatito",
     "Coco meets baby kitten\nplaying with yarn",
     "Coco rencontre un chaton\njouant avec une pelote",
     "Coco encuentra un gatito\njugando con un ovillo",
     'Coco met a tiny kitten playing with a ball of yarn.\n"Is this your eggshell, little kitten?"\nThe kitten purred and shook its head.\n"No, silly! Kittens don\'t hatch from eggs! But keep looking!"',
     'Coco rencontra un petit chaton qui jouait avec une pelote de laine.\n« Est-ce ta coquille, petit chaton ? »\nLe chaton ronronna et secoua la tête.\n« Mais non ! Les chatons ne sortent pas des œufs ! Mais continue à chercher ! »',
     'Coco encontró un gatito jugando con un ovillo de lana.\n«¿Es tu cascarón, gatito?»\nEl gatito ronroneó y negó con la cabeza.\n«¡No, tontita! ¡Los gatitos no nacen de huevos! ¡Pero sigue buscando!»',
     "Did you know? Kittens are born with their eyes closed. They can't see for the first 7-10 days!",
     "Le savais-tu ? Les chatons naissent les yeux fermés. Ils ne voient rien pendant les 7 à 10 premiers jours !",
     "¿Sabías que...? ¡Los gatitos nacen con los ojos cerrados y no pueden ver durante los primeros 7 a 10 días!"),

    ("🐶", "Baby Puppy", "Bébé Chiot", "Perrito",
     "Coco meets baby puppy\nwagging its tail",
     "Coco rencontre un chiot\nremuant la queue",
     "Coco encuentra un perrito\nmoviendo la cola",
     'Next, Coco found a puppy wagging its tail.\n"Is this your eggshell, little puppy?"\nThe puppy barked happily.\n"Woof! No way! Puppies don\'t come from eggs! Try asking someone else!"',
     'Ensuite, Coco trouva un chiot qui remuait la queue.\n« Est-ce ta coquille, petit chiot ? »\nLe chiot aboya joyeusement.\n« Wouf ! Pas du tout ! Les chiots ne viennent pas des œufs ! Essaie de demander à quelqu\'un d\'autre ! »',
     'Después, Coco encontró un perrito moviendo la cola.\n«¿Es tu cascarón, perrito?»\nEl perrito ladró feliz.\n«¡Guau! ¡Para nada! ¡Los perritos no vienen de huevos! ¡Pregúntale a alguien más!»',
     "Did you know? Puppies are born deaf! They can't hear anything for the first two weeks.",
     "Le savais-tu ? Les chiots naissent sourds ! Ils n'entendent rien pendant les deux premières semaines.",
     "¿Sabías que...? ¡Los perritos nacen sordos! No pueden oír nada durante las dos primeras semanas."),

    ("🐰", "Baby Bunny", "Bébé Lapin", "Conejito",
     "Coco meets baby bunny\nnibbling a carrot",
     "Coco rencontre un lapin\ngrignotant une carotte",
     "Coco encuentra un conejito\nmordisqueando una zanahoria",
     'Coco hopped over to a fluffy bunny nibbling a carrot.\n"Is this your eggshell, little bunny?"\nThe bunny wiggled its nose.\n"Hee hee! No! Bunnies don\'t hatch from eggs! But I do love Easter eggs!"',
     'Coco sautilla jusqu\'à un petit lapin tout doux qui grignotait une carotte.\n« Est-ce ta coquille, petit lapin ? »\nLe lapin remua son petit nez.\n« Hi hi ! Non ! Les lapins ne sortent pas des œufs ! Mais j\'adore les œufs de Pâques ! »',
     'Coco saltó hacia un conejito esponjoso que mordisqueaba una zanahoria.\n«¿Es tu cascarón, conejito?»\nEl conejito movió su naricita.\n«¡Ji ji! ¡No! ¡Los conejos no nacen de huevos! ¡Pero me encantan los huevos de Pascua!»',
     "Did you know? Baby bunnies are called kittens too! A group of baby bunnies is called a litter.",
     "Le savais-tu ? Les bébés lapins s'appellent des « lapereaux » ! Une portée peut compter jusqu'à 12 petits.",
     "¿Sabías que...? ¡Los bebés conejo también se llaman gazapos! Una camada puede tener hasta 12 conejitos."),

    ("🦄", "Baby Unicorn", "Bébé Licorne", "Unicornio bebé",
     "Coco meets baby unicorn\nwith sparkly horn",
     "Coco rencontre une licorne\navec corne scintillante",
     "Coco encuentra un unicornio\ncon cuerno brillante",
     'Coco spotted a baby unicorn with a sparkly little horn.\n"Is this your eggshell, little unicorn?"\nThe unicorn shook its rainbow mane.\n"Oh no! Unicorns are born with a sprinkle of magic, not from eggs!"',
     'Coco aperçut un bébé licorne avec une petite corne scintillante.\n« Est-ce ta coquille, petite licorne ? »\nLa licorne secoua sa crinière arc-en-ciel.\n« Oh non ! Les licornes naissent avec une pincée de magie, pas dans des œufs ! »',
     'Coco vio un bebé unicornio con un cuernito brillante.\n«¿Es tu cascarón, unicornio?»\nEl unicornio sacudió su melena arcoíris.\n«¡Oh no! ¡Los unicornios nacen con una pizca de magia, no de huevos!»',
     "Did you know? The unicorn is the national animal of Scotland!",
     "Le savais-tu ? La licorne est l'animal national de l'Écosse !",
     "¿Sabías que...? ¡El unicornio es el animal nacional de Escocia!"),

    ("🐼", "Baby Panda", "Bébé Panda", "Pandita",
     "Coco meets baby panda\nmunching bamboo",
     "Coco rencontre un panda\ncroquant du bambou",
     "Coco encuentra un pandita\nmasticando bambú",
     'A round baby panda was munching on bamboo.\n"Is this your eggshell, little panda?"\nThe panda giggled.\n"No no! Pandas don\'t come from eggs! But this bamboo is yummy — want some?"',
     'Un bébé panda tout rond croquait du bambou.\n« Est-ce ta coquille, petit panda ? »\nLe panda rigola.\n« Non non ! Les pandas ne viennent pas des œufs ! Mais ce bambou est trop bon — tu en veux ? »',
     'Un bebé panda redondito masticaba bambú.\n«¿Es tu cascarón, pandita?»\nEl panda se rió.\n«¡No no! ¡Los pandas no vienen de huevos! Pero este bambú está riquísimo, ¿quieres un poco?»',
     "Did you know? Baby pandas are born pink and tiny — about the size of a stick of butter!",
     "Le savais-tu ? Les bébés pandas naissent roses et minuscules — environ la taille d'une plaquette de beurre !",
     "¿Sabías que...? ¡Los bebés panda nacen rosados y diminutos, del tamaño de una barra de mantequilla!"),

    ("🐬", "Baby Dolphin", "Bébé Dauphin", "Delfincito",
     "Coco meets baby dolphin\nsplashing in waves",
     "Coco rencontre un dauphin\nfaisant des éclaboussures",
     "Coco encuentra un delfincito\nchapoteando en las olas",
     'Coco swam to a baby dolphin splashing in the waves.\n"Is this your eggshell, little dolphin?"\nThe dolphin did a flip.\n"Haha! Dolphins are born in the ocean, not in eggs! Keep searching, friend!"',
     'Coco nagea jusqu\'à un bébé dauphin qui faisait des éclaboussures.\n« Est-ce ta coquille, petit dauphin ? »\nLe dauphin fit un saut.\n« Ha ha ! Les dauphins naissent dans l\'océan, pas dans des œufs ! Continue à chercher ! »',
     'Coco nadó hasta un bebé delfín que chapoteaba en las olas.\n«¿Es tu cascarón, delfincito?»\nEl delfín dio una pirueta.\n«¡Jaja! ¡Los delfines nacen en el océano, no en huevos! ¡Sigue buscando!»',
     "Did you know? Baby dolphins can swim from the moment they're born!",
     "Le savais-tu ? Les bébés dauphins savent nager dès la naissance !",
     "¿Sabías que...? ¡Los delfines bebé pueden nadar desde que nacen!"),

    ("🐴", "Baby Foal", "Bébé Poulain", "Potrito",
     "Coco meets wobbly baby horse\nlearning to stand",
     "Coco rencontre un poulain\napprenant à se lever",
     "Coco encuentra un potrito\naprendiendo a pararse",
     'A wobbly baby horse was learning to stand.\n"Is this your eggshell, little foal?"\nThe foal stumbled and laughed.\n"Neigh! Horses don\'t hatch from eggs! We stand up right after we\'re born!"',
     'Un bébé cheval tout chancelant apprenait à se tenir debout.\n« Est-ce ta coquille, petit poulain ? »\nLe poulain trébucha et rigola.\n« Hiii ! Les chevaux ne sortent pas des œufs ! On se lève juste après être nés ! »',
     'Un potrito tambaleante aprendía a ponerse de pie.\n«¿Es tu cascarón, potrito?»\nEl potrito tropezó y se rió.\n«¡Iiih! ¡Los caballos no nacen de huevos! ¡Nos ponemos de pie justo al nacer!»',
     "Did you know? Baby horses can walk and run just a few hours after birth!",
     "Le savais-tu ? Les poulains peuvent marcher et courir quelques heures après la naissance !",
     "¿Sabías que...? ¡Los potros pueden caminar y correr pocas horas después de nacer!"),

    ("🦋", "Baby Butterfly", "Bébé Papillon", "Mariposita",
     "Coco meets baby butterfly\nstretching new wings",
     "Coco rencontre un papillon\ndéployant ses ailes",
     "Coco encuentra una mariposa\nestirando sus alas",
     'A baby butterfly was stretching its new wings.\n"Is this your eggshell, little butterfly?"\nThe butterfly fluttered.\n"Well... I did come from an egg! But it was a tiny, tiny egg on a leaf. That big shell isn\'t mine!"',
     'Un bébé papillon déployait ses nouvelles ailes.\n« Est-ce ta coquille, petit papillon ? »\nLe papillon voleta.\n« Eh bien... je suis sorti d\'un œuf ! Mais c\'était un tout tout petit œuf sur une feuille. Cette grosse coquille n\'est pas la mienne ! »',
     'Una mariposa bebé estiraba sus alas nuevas.\n«¿Es tu cascarón, mariposita?»\nLa mariposa revoloteó.\n«Bueno... ¡yo sí salí de un huevo! Pero era un huevito chiquito en una hoja. ¡Ese cascarón grande no es mío!»',
     "Did you know? Before becoming a butterfly, they're caterpillars! The change is called metamorphosis.",
     "Le savais-tu ? Avant de devenir papillon, ils sont d'abord des chenilles ! Ce changement s'appelle la métamorphose.",
     "¿Sabías que...? ¡Antes de ser mariposa, son orugas! Este cambio se llama metamorfosis."),

    ("🦖", "Baby Dinosaur", "Bébé Dinosaure", "Bebé Dinosaurio",
     "Coco meets baby T-Rex\nplaying with a rock",
     "Coco rencontre un bébé T-Rex\njouant avec un caillou",
     "Coco encuentra un bebé T-Rex\njugando con una piedra",
     'Coco stumbled upon a baby T-Rex playing with a rock.\n"Is this your eggshell, little dinosaur?"\nThe dino roared softly.\n"RAWR! Dinosaurs did hatch from eggs, long ago! But that shell is too small for me!"',
     'Coco tomba sur un bébé T-Rex qui jouait avec un caillou.\n« Est-ce ta coquille, petit dinosaure ? »\nLe dino rugit doucement.\n« GRRRR ! Les dinosaures sortaient bien des œufs, il y a très longtemps ! Mais cette coquille est trop petite pour moi ! »',
     'Coco se encontró con un bebé T-Rex jugando con una piedra.\n«¿Es tu cascarón, dinosaurio?»\nEl dino rugió suavecito.\n«¡RAWR! ¡Los dinosaurios sí nacían de huevos, hace muuucho tiempo! ¡Pero ese cascarón es muy pequeño para mí!»',
     "Did you know? The biggest dinosaur eggs ever found were as big as a basketball!",
     "Le savais-tu ? Les plus gros œufs de dinosaure jamais trouvés faisaient la taille d'un ballon de basket !",
     "¿Sabías que...? ¡Los huevos de dinosaurio más grandes jamás encontrados eran del tamaño de un balón de baloncesto!"),

    ("🦁", "Baby Lion", "Bébé Lion", "Leoncito",
     "Coco meets baby lion\npracticing its roar",
     "Coco rencontre un lionceau\ns'entraînant à rugir",
     "Coco encuentra un leoncito\npracticando su rugido",
     'A brave little lion cub was practicing its roar.\n"Is this your eggshell, little lion?"\nThe cub tried to roar but it came out like a squeak.\n"Mew! No! Lions are not born from eggs! But I\'m the king of this jungle!"',
     'Un petit lionceau courageux s\'entraînait à rugir.\n« Est-ce ta coquille, petit lion ? »\nLe lionceau essaya de rugir mais ça fit un petit couinement.\n« Miiou ! Non ! Les lions ne naissent pas dans des œufs ! Mais je suis le roi de cette jungle ! »',
     'Un valiente cachorro de león practicaba su rugido.\n«¿Es tu cascarón, leoncito?»\nEl cachorro intentó rugir pero le salió un maullido.\n«¡Miu! ¡No! ¡Los leones no nacen de huevos! ¡Pero soy el rey de esta selva!»',
     "Did you know? Lion cubs are born with spots! The spots fade away as they grow up.",
     "Le savais-tu ? Les lionceaux naissent avec des taches ! Elles disparaissent en grandissant.",
     "¿Sabías que...? ¡Los cachorros de león nacen con manchas! Desaparecen cuando crecen."),

    ("🐻", "Baby Bear", "Bébé Ours", "Osito",
     "Coco meets baby bear\nlicking honey",
     "Coco rencontre un ourson\nléchant du miel",
     "Coco encuentra un osito\nlamiendo miel",
     'A fuzzy bear cub was licking honey from its paws.\n"Is this your eggshell, little bear?"\nThe cub licked its lips.\n"Nope! Bears are born in cozy caves, not eggs! Want some honey?"',
     'Un petit ourson tout poilu léchait du miel sur ses pattes.\n« Est-ce ta coquille, petit ours ? »\nL\'ourson se lécha les babines.\n« Non ! Les ours naissent dans des grottes bien douillettes, pas dans des œufs ! Tu veux du miel ? »',
     'Un osito peludo lamía miel de sus patitas.\n«¿Es tu cascarón, osito?»\nEl osito se relamió.\n«¡No! ¡Los osos nacen en cuevas calentitas, no en huevos! ¿Quieres un poco de miel?»',
     "Did you know? Bear cubs are born during winter while mama bear is hibernating!",
     "Le savais-tu ? Les oursons naissent en hiver pendant que maman ours hiberne !",
     "¿Sabías que...? ¡Los oseznos nacen en invierno mientras mamá osa hiberna!"),

    ("🐘", "Baby Elephant", "Bébé Éléphant", "Elefantito",
     "Coco meets baby elephant\nspraying water",
     "Coco rencontre un éléphanteau\ns'aspergeant d'eau",
     "Coco encuentra un elefantito\nrociándose agua",
     'A baby elephant was spraying water with its trunk.\n"Is this your eggshell, little elephant?"\nThe elephant trumpeted and laughed.\n"PFRRRT! No way! I was way too big to fit in any egg!"',
     'Un bébé éléphant s\'aspergeait d\'eau avec sa trompe.\n« Est-ce ta coquille, petit éléphant ? »\nL\'éléphant barrit et rigola.\n« PFRRRT ! Pas possible ! J\'étais bien trop grand pour rentrer dans un œuf ! »',
     'Un elefantito se rociaba agua con su trompa.\n«¿Es tu cascarón, elefantito?»\nEl elefante trompeteó y se rió.\n«¡PFRRRT! ¡Imposible! ¡Yo era demasiado grande para caber en un huevo!»',
     "Did you know? Baby elephants weigh about 100 kg at birth — that's like a big refrigerator!",
     "Le savais-tu ? Les bébés éléphants pèsent environ 100 kg à la naissance — c'est comme un gros réfrigérateur !",
     "¿Sabías que...? ¡Los elefantes bebé pesan unos 100 kg al nacer, como un refrigerador grande!"),

    ("🦒", "Baby Giraffe", "Bébé Girafe", "Jirafita",
     "Coco meets baby giraffe\nreaching for a leaf",
     "Coco rencontre un girafon\nessayant d'attraper une feuille",
     "Coco encuentra una jirafita\nalcanzando una hoja",
     'A tall baby giraffe was trying to reach a leaf.\n"Is this your eggshell, little giraffe?"\nThe giraffe bent its long neck down.\n"No! Giraffes don\'t come from eggs. But even as a baby, I\'m already so tall!"',
     'Un bébé girafe essayait d\'attraper une feuille.\n« Est-ce ta coquille, petite girafe ? »\nLa girafe baissa son long cou.\n« Non ! Les girafes ne viennent pas des œufs. Mais même bébé, je suis déjà tellement grande ! »',
     'Una jirafita intentaba alcanzar una hoja.\n«¿Es tu cascarón, jirafita?»\nLa jirafa agachó su largo cuello.\n«¡No! Las jirafas no vienen de huevos. ¡Pero hasta de bebé ya soy altísima!»',
     "Did you know? Baby giraffes are 1.80 m tall at birth — taller than most adult humans!",
     "Le savais-tu ? Les bébés girafes mesurent 1m80 à la naissance — plus grand que la plupart des adultes !",
     "¿Sabías que...? ¡Las jirafas bebé miden 1.80 m al nacer, más alto que la mayoría de los adultos!"),

    ("🐧", "Baby Penguin", "Bébé Pingouin", "Pingüinito",
     "Coco meets baby penguin\nsliding on its belly",
     "Coco rencontre un bébé pingouin\nglissant sur le ventre",
     "Coco encuentra un pingüinito\ndeslizándose",
     'A baby penguin was sliding on its belly on the ice.\n"Is this your eggshell, little penguin?"\nThe penguin clapped its flippers.\n"Yes, I did hatch from an egg! But my egg was kept warm by my daddy\'s feet! That shell isn\'t mine."',
     'Un bébé pingouin glissait sur le ventre sur la glace.\n« Est-ce ta coquille, petit pingouin ? »\nLe pingouin tapa des ailerons.\n« Oui, moi je suis bien sorti d\'un œuf ! Mais mon œuf, c\'est papa qui le gardait au chaud sur ses pieds ! Cette coquille n\'est pas la mienne. »',
     'Un pingüinito se deslizaba sobre su barriguita en el hielo.\n«¿Es tu cascarón, pingüinito?»\nEl pingüino aplaudió con sus aletas.\n«¡Sí, yo sí nací de un huevo! ¡Pero mi huevo lo cuidaba mi papá en sus pies! Ese cascarón no es mío.»',
     "Did you know? Emperor penguin daddies keep the egg warm on their feet for 2 months without eating!",
     "Le savais-tu ? Les papas manchots empereurs gardent l'œuf au chaud sur leurs pieds pendant 2 mois sans manger !",
     "¿Sabías que...? ¡Los papás pingüino emperador mantienen el huevo caliente en sus pies durante 2 meses sin comer!"),

    ("🐨", "Baby Koala", "Bébé Koala", "Koalita",
     "Coco meets sleepy baby koala\non mama's back",
     "Coco rencontre un koala endormi\nsur le dos de maman",
     "Coco encuentra un koalita dormilón\nen la espalda de mamá",
     'A sleepy baby koala was clinging to its mama\'s back.\n"Is this your eggshell, little koala?"\nThe koala yawned.\n"Zzz... no... koalas ride in mama\'s pouch... not eggs... zzz..."',
     'Un bébé koala tout endormi s\'accrochait au dos de sa maman.\n« Est-ce ta coquille, petit koala ? »\nLe koala bâilla.\n« Zzz... non... les koalas se baladent dans la poche de maman... pas dans des œufs... zzz... »',
     'Un koalita dormilón se aferraba a la espalda de su mamá.\n«¿Es tu cascarón, koalita?»\nEl koala bostezó.\n«Zzz... no... los koalas viajan en la bolsa de mamá... no en huevos... zzz...»',
     "Did you know? Koalas sleep up to 22 hours a day! That's almost the whole day!",
     "Le savais-tu ? Les koalas dorment jusqu'à 22 heures par jour ! C'est presque toute la journée !",
     "¿Sabías que...? ¡Los koalas duermen hasta 22 horas al día! ¡Eso es casi todo el día!"),

    ("🦊", "Baby Fox", "Bébé Renard", "Zorrito",
     "Coco meets baby fox\nchasing its tail",
     "Coco rencontre un renardeau\ncourant après sa queue",
     "Coco encuentra un zorrito\npersiguiendo su cola",
     'A playful baby fox was chasing its own tail.\n"Is this your eggshell, little fox?"\nThe fox spun around.\n"Nope! I was born in a den underground! This tail is more fun than any egg!"',
     'Un bébé renard joueur courait après sa propre queue.\n« Est-ce ta coquille, petit renard ? »\nLe renard tourna sur lui-même.\n« Non ! Je suis né dans un terrier sous la terre ! Ma queue est bien plus amusante qu\'un œuf ! »',
     'Un zorrito juguetón perseguía su propia cola.\n«¿Es tu cascarón, zorrito?»\nEl zorro dio vueltas.\n«¡No! ¡Yo nací en una madriguera bajo tierra! ¡Mi cola es más divertida que cualquier huevo!»',
     "Did you know? Baby foxes are called kits! They love to play-fight with their brothers and sisters.",
     "Le savais-tu ? Les bébés renards s'appellent des renardeaux ! Ils adorent se bagarrer pour jouer.",
     "¿Sabías que...? ¡Los bebés zorro se llaman cachorros! Les encanta jugar a pelear con sus hermanos."),

    ("🦉", "Baby Owl", "Bébé Hibou", "Buhíto",
     "Coco meets baby owl\nblinking huge eyes",
     "Coco rencontre un bébé hibou\nclignant ses grands yeux",
     "Coco encuentra un buhíto\nparpadeando",
     'A fluffy owlet was blinking its huge eyes.\n"Is this your eggshell, little owl?"\nThe owl hooted softly.\n"Hoo hoo! Yes, owls hatch from eggs! But mine was in a nest high up in a tree. That one\'s not mine!"',
     'Un bébé hibou tout duveteux clignait ses grands yeux.\n« Est-ce ta coquille, petit hibou ? »\nLe hibou hulula doucement.\n« Hou hou ! Oui, les hiboux sortent des œufs ! Mais le mien était dans un nid tout en haut d\'un arbre. Celui-là n\'est pas le mien ! »',
     'Un buhíto esponjoso parpadeaba con sus enormes ojos.\n«¿Es tu cascarón, buhíto?»\nEl búho ululó suavecito.\n«¡Uh uh! ¡Sí, los búhos nacen de huevos! Pero el mío estaba en un nido bien alto en un árbol. ¡Ese no es mío!»',
     "Did you know? Owls can turn their heads almost all the way around — up to 270 degrees!",
     "Le savais-tu ? Les hiboux peuvent tourner la tête presque complètement — jusqu'à 270 degrés !",
     "¿Sabías que...? ¡Los búhos pueden girar la cabeza casi por completo, hasta 270 grados!"),

    ("🐢", "Baby Turtle", "Bébé Tortue", "Tortuguita",
     "Coco meets baby turtle\ncrawling toward the sea",
     "Coco rencontre une tortue\nrampant vers la mer",
     "Coco encuentra una tortuguita\ncaminando al mar",
     'A tiny baby turtle was slowly crawling toward the sea.\n"Is this your eggshell, little turtle?"\nThe turtle smiled.\n"Yes! I hatched from an egg on the beach! But my shell was buried in the sand. That one isn\'t mine!"',
     'Une toute petite tortue rampait lentement vers la mer.\n« Est-ce ta coquille, petite tortue ? »\nLa tortue sourit.\n« Oui ! Je suis sortie d\'un œuf sur la plage ! Mais ma coquille était enterrée dans le sable. Celle-là n\'est pas la mienne ! »',
     'Una tortuguita caminaba lentamente hacia el mar.\n«¿Es tu cascarón, tortuguita?»\nLa tortuga sonrió.\n«¡Sí! ¡Yo nací de un huevo en la playa! Pero mi cascarón estaba enterrado en la arena. ¡Ese no es mío!»',
     "Did you know? Sea turtles return to the same beach where they were born to lay their own eggs!",
     "Le savais-tu ? Les tortues de mer reviennent sur la même plage où elles sont nées pour pondre !",
     "¿Sabías que...? ¡Las tortugas marinas regresan a la misma playa donde nacieron para poner sus huevos!"),

    ("🐸", "Baby Frog", "Bébé Grenouille", "Ranita",
     "Coco meets baby frog\non a lily pad",
     "Coco rencontre une grenouille\nsur un nénuphar",
     "Coco encuentra una ranita\nen un nenúfar",
     'A little frog was sitting on a lily pad.\n"Is this your eggshell, little frog?"\nThe frog jumped up.\n"RIBBIT! Frogs don\'t come from eggs like that! We start as tiny dots in jelly in the water! I was a tadpole first!"',
     'Une petite grenouille était assise sur un nénuphar.\n« Est-ce ta coquille, petite grenouille ? »\nLa grenouille bondit.\n« COÂÂÂ ! Les grenouilles ne sortent pas d\'un œuf comme ça ! On commence comme des petits points dans de la gelée dans l\'eau ! J\'étais un têtard d\'abord ! »',
     'Una ranita estaba sentada en una hoja de nenúfar.\n«¿Es tu cascarón, ranita?»\nLa rana saltó.\n«¡CROAC! ¡Las ranas no nacen de huevos así! ¡Empezamos como puntitos en gelatina en el agua! ¡Primero fui un renacuajo!»',
     "Did you know? Frogs go through metamorphosis — from egg to tadpole to frog!",
     "Le savais-tu ? Les grenouilles font une métamorphose — d'œuf à têtard puis grenouille !",
     "¿Sabías que...? ¡Las ranas pasan por una metamorfosis: de huevo a renacuajo y luego rana!"),

    ("🐮", "Baby Highland Cow", "Bébé Vache Highland", "Vaquita Highland",
     "Coco meets baby Highland cow\npeeking through fluffy bangs",
     "Coco rencontre une vache Highland\nregardant à travers sa frange",
     "Coco encuentra una vaquita Highland\nasomando entre su flequillo",
     'A fluffy baby Highland cow peeked through its long bangs.\n"Is this your eggshell, little cow?"\nThe calf swished its fluffy hair.\n"Moo! No way! I just have the best hair of all the babies here!"',
     'Un bébé vache Highland regarda à travers sa longue frange.\n« Est-ce ta coquille, petite vache ? »\nLe veau secoua ses poils tout doux.\n« Meuh ! Pas du tout ! J\'ai juste la plus belle coiffure de tous les bébés ici ! »',
     'Una vaquita Highland asomó entre su largo flequillo.\n«¿Es tu cascarón, vaquita?»\nEl ternero agitó su pelo esponjoso.\n«¡Muu! ¡Para nada! ¡Solo tengo el mejor peinado de todos los bebés de aquí!»',
     "Did you know? Highland cows come from Scotland and their long hair keeps them warm in cold weather!",
     "Le savais-tu ? Les vaches Highland viennent d'Écosse et leurs longs poils les protègent du froid !",
     "¿Sabías que...? ¡Las vacas Highland vienen de Escocia y su pelo largo las mantiene calientes en el frío!"),

    ("🦩", "Baby Flamingo", "Bébé Flamant Rose", "Flaminguito",
     "Coco meets baby flamingo\nstanding on one leg",
     "Coco rencontre un flamant rose\nsur une patte",
     "Coco encuentra un flaminguito\nen una pata",
     'A baby flamingo was standing on one leg (or at least trying to).\n"Is this your eggshell, little flamingo?"\nThe flamingo wobbled.\n"Yes, I hatched from an egg! But mine was on a big mud nest. And look — I\'m grey, not pink yet!"',
     'Un bébé flamant rose se tenait sur une patte (ou essayait).\n« Est-ce ta coquille, petit flamant ? »\nLe flamant vacilla.\n« Oui, je suis sorti d\'un œuf ! Mais le mien était sur un gros nid de boue. Et regarde — je suis gris, pas encore rose ! »',
     'Un flaminguito se sostenía en una pata (o lo intentaba).\n«¿Es tu cascarón, flaminguito?»\nEl flamenco se tambaleó.\n«¡Sí, yo nací de un huevo! Pero el mío estaba en un nido grande de barro. ¡Y mira, soy gris, todavía no soy rosa!»',
     "Did you know? Flamingos are born grey! They turn pink from eating tiny shrimp.",
     "Le savais-tu ? Les flamants naissent gris ! Ils deviennent roses en mangeant des petites crevettes.",
     "¿Sabías que...? ¡Los flamencos nacen grises! Se vuelven rosas por comer camarones pequeñitos."),

    ("🦥", "Baby Sloth", "Bébé Paresseux", "Perezosito",
     "Coco meets baby sloth\nhanging upside down",
     "Coco rencontre un paresseux\nsuspendu à l'envers",
     "Coco encuentra un perezosito\ncolgando boca abajo",
     'A baby sloth was hanging upside down from a branch, veeeery slowly.\n"Is this... your eggshell... little sloth?"\nThe sloth blinked.\n"Nooooo... sloths... don\'t... come... from... eggs... we just... hang... around..."',
     'Un bébé paresseux pendait à l\'envers sur une branche, trèèèès lentement.\n« Est-ce que c\'est... ta coquille... petit paresseux ? »\nLe paresseux cligna des yeux.\n« Nooooon... les paresseux... ne viennent... pas... des œufs... on fait juste... rien du tout... »',
     'Un perezosito colgaba boca abajo de una rama, muuuy lentamente.\n«¿Es... tu cascarón... perezosito?»\nEl perezoso pestañeó.\n«Nooooo... los perezosos... no nacen... de huevos... solo... colgamos... por aquí...»',
     "Did you know? Sloths are so slow that algae grows on their fur, making them look green!",
     "Le savais-tu ? Les paresseux sont si lents que des algues poussent sur leur fourrure !",
     "¿Sabías que...? ¡Los perezosos son tan lentos que les crecen algas en el pelo!"),

    ("🦦", "Baby Otter", "Bébé Loutre", "Nutria bebé",
     "Coco meets baby otter\nfloating with a shell",
     "Coco rencontre une loutre\nflottant avec un coquillage",
     "Coco encuentra una nutria\nflotando con una conchita",
     'A baby otter was floating on its back holding a little shell.\n"Is this your eggshell, little otter?"\nThe otter giggled.\n"No! Otters are born in the water! But I do love collecting shells! Can I keep it?"',
     'Un bébé loutre flottait sur le dos en tenant un petit coquillage.\n« Est-ce ta coquille, petite loutre ? »\nLa loutre gloussa.\n« Non ! Les loutres naissent dans l\'eau ! Mais j\'adore collectionner les coquilles ! Je peux la garder ? »',
     'Una nutria bebé flotaba boca arriba sujetando una conchita.\n«¿Es tu cascarón, nutria?»\nLa nutria se rió.\n«¡No! ¡Las nutrias nacen en el agua! Pero me encanta coleccionar conchas. ¿Me la puedo quedar?»',
     "Did you know? Sea otters hold hands while sleeping so they don't drift apart!",
     "Le savais-tu ? Les loutres de mer se tiennent par la main en dormant pour ne pas dériver !",
     "¿Sabías que...? ¡Las nutrias marinas se toman de las manos al dormir para no separarse!"),

    ("🦔", "Baby Hedgehog", "Bébé Hérisson", "Erizito",
     "Coco meets baby hedgehog\ncurled up in a ball",
     "Coco rencontre un hérisson\nroulé en boule",
     "Coco encuentra un erizito\nhecho bolita",
     'A tiny baby hedgehog was curled up in a ball.\n"Is this your eggshell, little hedgehog?"\nThe hedgehog uncurled slowly.\n"No! I was born soft and tiny. My spikes grew later! Be careful though — I\'m a little prickly!"',
     'Un tout petit hérisson était roulé en boule.\n« Est-ce ta coquille, petit hérisson ? »\nLe hérisson se déroula lentement.\n« Non ! Je suis né tout doux et tout petit. Mes piquants ont poussé après ! Attention quand même — je pique un peu ! »',
     'Un erizito estaba hecho bolita.\n«¿Es tu cascarón, erizito?»\nEl erizo se desenrolló lentamente.\n«¡No! Yo nací blandito y chiquito. ¡Mis púas crecieron después! ¡Pero cuidado, pico un poquito!»',
     "Did you know? Baby hedgehogs are born with soft spines hidden under their skin!",
     "Le savais-tu ? Les bébés hérissons naissent avec des piquants mous cachés sous la peau !",
     "¿Sabías que...? ¡Los erizos bebé nacen con púas blandas escondidas bajo la piel!"),

    ("🦙", "Baby Llama", "Bébé Lama", "Llamita",
     "Coco meets baby llama\nprancing proudly",
     "Coco rencontre un lama\ntrottinant fièrement",
     "Coco encuentra una llamita\ntrotando orgullosa",
     'A fluffy baby llama was prancing around proudly.\n"Is this your eggshell, little llama?"\nThe llama tossed its fluffy head.\n"HA! Llamas don\'t hatch from eggs! We stand up and run right after we\'re born! Pretty cool, right?"',
     'Un bébé lama tout duveteux trottinait fièrement.\n« Est-ce ta coquille, petit lama ? »\nLe lama secoua sa tête touffue.\n« HA ! Les lamas ne sortent pas des œufs ! On se lève et on court juste après être nés ! Trop cool, non ? »',
     'Una llamita esponjosa trotaba orgullosa.\n«¿Es tu cascarón, llamita?»\nLa llama sacudió su cabeza peluda.\n«¡JA! ¡Las llamas no nacen de huevos! ¡Nos paramos y corremos justo al nacer! ¿Genial, no?»',
     "Did you know? Llamas hum to talk to each other! A mama llama hums to her baby all the time.",
     "Le savais-tu ? Les lamas fredonnent pour se parler ! Maman lama fredonne à son bébé tout le temps.",
     "¿Sabías que...? ¡Las llamas tararean para comunicarse! Mamá llama le tararea a su bebé todo el tiempo."),

    ("🐥", "Baby Duckling", "Bébé Caneton", "Patito",
     "Coco meets baby duckling\npaddling in a puddle",
     "Coco rencontre un caneton\nbarbotant dans une flaque",
     "Coco encuentra un patito\nchapoteando en un charco",
     'A yellow duckling was paddling in a puddle.\n"Is this your eggshell, little duck?"\nThe duckling quacked.\n"QUACK! Yes, I came from an egg! But mine was in a nest by the pond! That shell isn\'t mine!"',
     'Un caneton jaune barbotait dans une flaque.\n« Est-ce ta coquille, petit canard ? »\nLe caneton caqueta.\n« COIN ! Oui, je suis sorti d\'un œuf ! Mais le mien était dans un nid au bord de la mare ! Celui-là n\'est pas le mien ! »',
     'Un patito amarillo chapoteaba en un charco.\n«¿Es tu cascarón, patito?»\nEl patito graznó.\n«¡CUAC! ¡Sí, yo salí de un huevo! Pero el mío estaba en un nido junto al estanque! ¡Ese no es mío!»',
     "Did you know? Ducklings can swim within hours of hatching! They follow the first thing they see.",
     "Le savais-tu ? Les canetons savent nager quelques heures après leur naissance !",
     "¿Sabías que...? ¡Los patitos pueden nadar pocas horas después de nacer!"),

    ("🐤", "Baby Chick", "Bébé Poussin", "Pollito",
     "Coco meets baby chick\npecking at seeds",
     "Coco rencontre un poussin\npicorant des graines",
     "Coco encuentra un pollito\npicoteando semillas",
     'A round little chick was pecking at seeds.\n"Is this your eggshell, little chick?"\nThe chick cheeped.\n"PEEP PEEP! I hatched from an egg too! But mine was warm under mama hen! That one looks different from mine!"',
     'Un petit poussin tout rond picorait des graines.\n« Est-ce ta coquille, petit poussin ? »\nLe poussin piailla.\n« PIOU PIOU ! Je suis sorti d\'un œuf aussi ! Mais le mien était bien au chaud sous maman poule ! Celui-là est différent du mien ! »',
     'Un pollito redondito picoteaba semillas.\n«¿Es tu cascarón, pollito?»\nEl pollito piaba.\n«¡PÍO PÍO! ¡Yo también nací de un huevo! Pero el mío estaba calentito bajo mamá gallina! ¡Ese se ve diferente al mío!»',
     "Did you know? Baby chicks start peeping inside the egg before they hatch!",
     "Le savais-tu ? Les poussins piaillent dans l'œuf avant d'éclore !",
     "¿Sabías que...? ¡Los pollitos pían dentro del huevo antes de nacer!"),

    ("🐑", "Baby Lamb", "Bébé Agneau", "Corderito",
     "Coco meets baby lamb\njumping in a meadow",
     "Coco rencontre un agneau\nsautant dans un pré",
     "Coco encuentra un corderito\nsaltando en un prado",
     'A woolly little lamb was jumping in a meadow.\n"Is this your eggshell, little lamb?"\nThe lamb bounced.\n"BAA! No! Lambs don\'t come from eggs! We come with fluffy wool and springy legs!"',
     'Un petit agneau laineux sautait dans un pré.\n« Est-ce ta coquille, petit agneau ? »\nL\'agneau bondit.\n« BÊÊ ! Non ! Les agneaux ne viennent pas des œufs ! On arrive avec de la laine toute douce et des pattes qui rebondissent ! »',
     'Un corderito lanudo saltaba en un prado.\n«¿Es tu cascarón, corderito?»\nEl cordero brincó.\n«¡BEE! ¡No! ¡Los corderos no nacen de huevos! ¡Llegamos con lana suavecita y patas saltarinas!»',
     "Did you know? Baby lambs can recognize their mama's voice from birth!",
     "Le savais-tu ? Les agneaux reconnaissent la voix de leur maman dès la naissance !",
     "¿Sabías que...? ¡Los corderos reconocen la voz de su mamá desde que nacen!"),

    ("🐉", "Baby Dragon", "Bébé Dragon", "Dragoncito",
     "Coco meets baby dragon\nblowing tiny smoke rings",
     "Coco rencontre un bébé dragon\nsoufflant des ronds de fumée",
     "Coco encuentra un dragoncito\nsoplando anillos de humo",
     'A baby dragon was blowing tiny smoke rings.\n"Is this your eggshell, little dragon?"\nThe dragon puffed.\n"Hmm! Dragons DO hatch from eggs! Big, sparkly eggs! But that shell isn\'t sparkly enough to be mine!"',
     'Un bébé dragon soufflait de petits ronds de fumée.\n« Est-ce ta coquille, petit dragon ? »\nLe dragon souffla.\n« Hmm ! Les dragons SORTENT bien des œufs ! De gros œufs tout brillants ! Mais cette coquille ne brille pas assez pour être la mienne ! »',
     'Un dragoncito soplaba pequeños anillos de humo.\n«¿Es tu cascarón, dragoncito?»\nEl dragón sopló.\n«¡Hmm! ¡Los dragones SÍ nacen de huevos! ¡Huevos grandes y brillantes! ¡Pero ese cascarón no brilla lo suficiente para ser mío!»',
     "Did you know? In many stories, dragon eggs need fire to hatch!",
     "Le savais-tu ? Dans beaucoup d'histoires, les œufs de dragon ont besoin de feu pour éclore !",
     "¿Sabías que...? ¡En muchas historias, los huevos de dragón necesitan fuego para eclosionar!"),
]

# ── Titles & texts per language ──
langs = {
    'en': {
        'title': "Whose Egg Is This?",
        'subtitle': "A Coco the Axolotl Adventure",
        'author': "Coco the Axolotl",
        'copyright': (
            "Whose Egg Is This?\n"
            "A Coco the Axolotl Adventure\n\n"
            "Written and illustrated by Coco the Axolotl\n\n"
            "\u00a9 2025 Coco the Axolotl. All rights reserved.\n\n"
            "No part of this publication may be reproduced, stored in a\n"
            "retrieval system, or transmitted in any form or by any means,\n"
            "electronic, mechanical, photocopying, recording, or otherwise,\n"
            "without the prior written permission of the copyright owner.\n\n"
            "www.cocotheaxolotl.org\n\n"
            "First edition, 2025"
        ),
        'intro': 'One sunny morning, little Coco found a cracked eggshell by the river.\n\n"Oh! Whose egg is this?" Coco wondered.\n"I must find who it belongs to!"\n\nAnd so, Coco set off on a big adventure to meet all the baby animals.',
        'intro_illust': "Coco by the river,\nfinding the cracked eggshell",
        'ending': 'Coco was tired. So many baby animals, but still no owner for the eggshell. Walking slowly back home along the river, Coco carried the shell on top like a little hat.\n\nThen — Mama Axolotl was standing by the water, and in her arms was the tiniest, cutest baby axolotl ever!\n\n"Mama! Whose egg was this?"\n\nMama smiled softly. "That\'s your little sister\'s eggshell, Coco. She hatched this morning! Say hello to Lili!"\n\nBaby Lili looked up at Coco with big sparkly eyes and giggled.\n\nCoco laughed and kept the shell on top. "I\'m keeping it! It\'s my special hat now!"\n\nAnd from that day on, Coco wore the eggshell like a tiny explorer helmet — showing baby Lili all the wonderful animal friends along the way.',
        'ending_illust': "Coco with eggshell on head (Caliméro style),\nMama Axolotl holding baby Lili,\nall baby animal friends waving",
        'the_end': "The End!",
        'back_title': "Did you enjoy the adventure?",
        'back_text': "Visit Coco online for free coloring pages,\ngames, and more fun!\n\ncocotheaxolotl.org\n\nFind Coco on mugs, t-shirts, stickers\nand more at our shop!",
    },
    'fr': {
        'title': "À qui est cet œuf ?",
        'subtitle': "Une aventure de Coco l'Axolotl",
        'author': "Coco the Axolotl",
        'copyright': (
            "\u00c0 qui est cet \u0153uf ?\n"
            "Une aventure de Coco l'Axolotl\n\n"
            "\u00c9crit et illustr\u00e9 par Coco the Axolotl\n\n"
            "\u00a9 2025 Coco the Axolotl. Tous droits r\u00e9serv\u00e9s.\n\n"
            "Aucune partie de cette publication ne peut \u00eatre reproduite,\n"
            "stock\u00e9e dans un syst\u00e8me de r\u00e9cup\u00e9ration, ou transmise\n"
            "sous quelque forme ou par quelque moyen que ce soit,\n"
            "\u00e9lectronique, m\u00e9canique, photocopie, enregistrement ou autre,\n"
            "sans l'autorisation \u00e9crite pr\u00e9alable du titulaire des droits.\n\n"
            "www.cocotheaxolotl.org\n\n"
            "Premi\u00e8re \u00e9dition, 2025"
        ),
        'intro': 'Par un beau matin ensoleillé, Coco trouva une coquille d\'œuf cassée au bord de la rivière.\n\n« Oh ! À qui est cet œuf ? » se demanda Coco.\n« Je dois trouver à qui il appartient ! »\n\nEt Coco partit pour une grande aventure à la rencontre de tous les bébés animaux.',
        'intro_illust': "Coco au bord de la rivière,\ntrouvant la coquille cassée",
        'ending': 'Coco était fatigué. Tant de bébés animaux, mais toujours pas de propriétaire pour la coquille. En rentrant doucement à la maison le long de la rivière, Coco portait la coquille sur la tête comme un petit chapeau.\n\nPuis — Maman Axolotl était là, au bord de l\'eau, et dans ses bras, il y avait le plus petit, le plus mignon bébé axolotl du monde !\n\n« Maman ! À qui est cet œuf ? »\n\nMaman sourit tendrement. « C\'est la coquille de ta petite sœur, Coco. Elle est née ce matin ! Dis bonjour à Lili ! »\n\nBébé Lili leva les yeux vers Coco avec de grands yeux tout brillants et fit un petit rire.\n\nCoco rigola et garda la coquille sur la tête. « Je la garde ! C\'est mon chapeau spécial maintenant ! »\n\nEt depuis ce jour-là, Coco porta la coquille comme un petit casque d\'explorateur — en présentant à bébé Lili tous les merveilleux amis animaux rencontrés en chemin.',
        'ending_illust': "Coco avec la coquille sur la tête (style Caliméro),\nMaman Axolotl tenant bébé Lili,\ntous les bébés animaux font coucou",
        'the_end': "Fin !",
        'back_title': "Tu as aimé l'aventure ?",
        'back_text': "Retrouve Coco en ligne pour des coloriages\ngratuits, des jeux et encore plus de fun !\n\ncocotheaxolotl.org\n\nRetrouve Coco sur des mugs, t-shirts,\nstickers et plus dans notre boutique !",
    },
    'es': {
        'title': "¿De quién es este huevo?",
        'subtitle': "Una aventura de Coco el Ajolote",
        'author': "Coco the Axolotl",
        'copyright': (
            "\u00bfDe qui\u00e9n es este huevo?\n"
            "Una aventura de Coco el Ajolote\n\n"
            "Escrito e ilustrado por Coco the Axolotl\n\n"
            "\u00a9 2025 Coco the Axolotl. Todos los derechos reservados.\n\n"
            "Ninguna parte de esta publicaci\u00f3n puede ser reproducida,\n"
            "almacenada en un sistema de recuperaci\u00f3n, o transmitida\n"
            "de ninguna forma ni por ning\u00fan medio, electr\u00f3nico,\n"
            "mec\u00e1nico, fotocopia, grabaci\u00f3n u otro,\n"
            "sin el permiso previo por escrito del titular de los derechos.\n\n"
            "www.cocotheaxolotl.org\n\n"
            "Primera edici\u00f3n, 2025"
        ),
        'intro': 'Una soleada mañana, Coco encontró un cascarón roto junto al río.\n\n«¡Oh! ¿De quién es este huevo?» se preguntó Coco.\n«¡Tengo que encontrar a quién pertenece!»\n\nY así, Coco partió en una gran aventura para conocer a todos los bebés animales.',
        'intro_illust': "Coco junto al río,\nencontrando el cascarón roto",
        'ending': 'Coco estaba cansado. Tantos bebés animales, pero el cascarón seguía sin dueño. Caminando lentamente de vuelta a casa junto al río, Coco llevaba el cascarón en la cabeza como un sombrerito.\n\nEntonces — Mamá Ajolote estaba junto al agua, ¡y en sus brazos había el más pequeño y tierno bebé ajolote jamás visto!\n\n«¡Mamá! ¿De quién es este huevo?»\n\nMamá sonrió con ternura. «Es el cascarón de tu hermanita, Coco. ¡Nació esta mañana! ¡Saluda a Lili!»\n\nLa bebé Lili miró a Coco con sus grandes ojos brillantes y se rió.\n\nCoco se rió y se dejó el cascarón en la cabeza. «¡Me lo quedo! ¡Ahora es mi sombrero especial!»\n\nY desde ese día, Coco llevó el cascarón como un pequeño casco de explorador — mientras le presentaba a su hermanita Lili todos los maravillosos amigos animales que había conocido.',
        'ending_illust': "Coco con el cascarón en la cabeza (estilo Calimero),\nMamá Ajolote con bebé Lili,\ntodos los bebés animales saludando",
        'the_end': "¡Fin!",
        'back_title': "¿Te gustó la aventura?",
        'back_text': "¡Visita a Coco en línea para dibujos gratis,\njuegos y más diversión!\n\ncocotheaxolotl.org\n\n¡Encuentra a Coco en tazas, camisetas,\nstickers y más en nuestra tienda!",
    },
}

def build_book(lang_code, texts, animals_data):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ── COVER ──
    slide = prs.slides.add_slide(blank_layout)
    # Pink gradient background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    # Cover illustration placeholder
    add_placeholder(slide, Inches(1), Inches(0.5), Inches(6.5), Inches(6),
                    f"[ COVER ART ]\n{texts['intro_illust']}")

    # Title
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(7.5), Inches(1.4),
                texts['title'], size=40, bold=True, color=DARK_PINK)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(8.2), Inches(7.5), Inches(0.7),
                texts['subtitle'], size=20, color=PINK)

    # Author
    add_textbox(slide, Inches(0.5), Inches(9.2), Inches(7.5), Inches(0.5),
                texts['author'], size=14, color=GREY)

    # ── COPYRIGHT PAGE ──
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(6.5), Inches(5),
                texts['copyright'], size=12, color=GREY, align=PP_ALIGN.CENTER)

    # ── INTRO ──
    slide = prs.slides.add_slide(blank_layout)
    add_placeholder(slide, Inches(1), Inches(0.4), Inches(6.5), Inches(5.5),
                    f"[ Illustration ]\n{texts['intro_illust']}")
    add_textbox(slide, Inches(0.6), Inches(6.2), Inches(7.3), Inches(4),
                texts['intro'], size=17, color=BLACK, align=PP_ALIGN.CENTER)

    # ── ANIMAL PAGES ──
    lang_idx = {'en': 0, 'fr': 1, 'es': 2}[lang_code]

    for i, a in enumerate(animals_data):
        emoji = a[0]
        name = a[1 + lang_idx]
        illust = a[4 + lang_idx]
        story = a[7 + lang_idx]
        fact = a[10 + lang_idx]

        slide = prs.slides.add_slide(blank_layout)

        # Page number
        add_textbox(slide, Inches(6.8), Inches(0.2), Inches(1.2), Inches(0.4),
                    f"{i+1} / 29", size=10, color=PLAC_BORDER, align=PP_ALIGN.RIGHT)

        # Illustration placeholder (bigger for portrait format)
        add_placeholder(slide, Inches(1), Inches(0.4), Inches(6.5), Inches(5),
                        f"[ Illustration ]\n{illust}")

        # Animal name
        add_textbox(slide, Inches(0.5), Inches(5.6), Inches(7.5), Inches(0.6),
                    f"{emoji} {name}", size=24, bold=True, color=DARK_PINK)

        # Story text
        add_textbox(slide, Inches(0.6), Inches(6.2), Inches(7.3), Inches(3),
                    story, size=14, color=BLACK, align=PP_ALIGN.CENTER)

        # Fun fact
        add_fun_fact(slide, Inches(1), Inches(9.3), Inches(6.5), Inches(1.1), fact)

    # ── ENDING ──
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    add_placeholder(slide, Inches(1), Inches(0.3), Inches(6.5), Inches(4.5),
                    f"[ FINAL ILLUSTRATION ]\n{texts['ending_illust']}")

    add_textbox(slide, Inches(0.5), Inches(5), Inches(7.5), Inches(4.8),
                texts['ending'], size=13, color=BLACK, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(10), Inches(7.5), Inches(0.6),
                texts['the_end'], size=30, bold=True, color=DARK_PINK)

    # ── BACK COVER ──
    slide = prs.slides.add_slide(blank_layout)
    add_placeholder(slide, Inches(2.25), Inches(0.5), Inches(4), Inches(4),
                    "[ Coco waving goodbye\nwith eggshell hat ]")

    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(7.5), Inches(0.6),
                texts['back_title'], size=22, bold=True, color=DARK_PINK)

    add_textbox(slide, Inches(0.5), Inches(5.6), Inches(7.5), Inches(2.8),
                texts['back_text'], size=15, color=BLACK, align=PP_ALIGN.CENTER)

    # QR note
    add_textbox(slide, Inches(2.5), Inches(8.8), Inches(3.5), Inches(0.5),
                "[ QR Code → cocotheaxolotl.org ]", size=12, color=PLAC_BORDER)

    add_textbox(slide, Inches(0.5), Inches(10.2), Inches(7.5), Inches(0.4),
                "© Coco the Axolotl • cocotheaxolotl.org", size=10, color=GREY)

    return prs


# ── Generate 3 files ──
out_dir = os.path.dirname(os.path.abspath(__file__))

for code in ['en', 'fr', 'es']:
    prs = build_book(code, langs[code], animals)
    filename = f"Whose-Egg-Is-This-{code.upper()}-8.5x11.pptx"
    path = os.path.join(out_dir, filename)
    prs.save(path)
    print(f"OK  {filename}")

print("\nDone! 3 PowerPoint files generated.")
