"""Build coloring book PowerPoints (B&W outlines) in 8.5x11 for EN/FR/ES.
Each animal gets a full-page B&W illustration with the animal name at the bottom."""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import re
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

SCRIPT_DIR = Path(__file__).parent
COLORING_DIR = SCRIPT_DIR / "whose-egg-coloring"

# Slide index (in source pptx) → coloring image filename
SLIDE_TO_COLORING = {
    1: "coco-eggshell_coloring .png",
    2: "axolotl-kitten-coloring.png",
    3: "axolotl_puffy-coloring.png",
    4: "axolotl-rabbit-coloring.png",
    5: "axolotl-unicorn-coloring.png",
    6: "axolotl-panda-coloring.png",
    7: "axolotl-dolphin-coloring.png",
    8: "axolotl-foal-coloring.png",
    9: "axolotl-butterfly-coloring.png",
    10: "axolotl-dinosaur-coloring.png",
    11: "axolotl-lion-coloring.png",
    12: "axolotl-bear-coloring.png",
    13: "axolotl-elephant-coloring.png",
    14: "axolotl-giraffe-coloring.png",
    15: "axolotl-penguin-coloring.png",
    16: "axolotl-koala-coloring.png",
    17: "axolotl-fox-coloring.png",
    18: "axolotl-owl-coloring.png",
    19: "axolotl-turtle-coloring.png",
    20: "axolotl-frog-coloring.png",
    21: "axolotl-highland-coloring.png",
    22: "axolotl-flamingo-coloring.png",
    23: "axolotl-sloth-coloring.png",
    24: "axolotl-otter-coloring.png",
    25: "axolotl-hedgehog-coloring.png",
    26: "axolotl-lama-coloring.png",
    27: "axolotl-duck-coloring.png",
    28: "axolotl-chick-coloring.png",
    29: "axolotl-lamb-coloring.png",
    30: "axolotl-dragon-coloring.png",
    31: "axolotl-baby-coloring.png",
}

FONT_NAME = "Baloo 2"
PINK_DARK = RGBColor(0xD6, 0x00, 0x6E)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)

# Copyright page text per language
COPYRIGHT_TEXTS = {
    'en': (
        "Whose Egg Is This? \u2014 Coloring Book\n"
        "A Coco the Axolotl Adventure\n\n"
        "Written and illustrated by Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. All rights reserved.\n\n"
        "No part of this publication may be reproduced, stored in a\n"
        "retrieval system, or transmitted in any form or by any means,\n"
        "electronic, mechanical, photocopying, recording, or otherwise,\n"
        "without the prior written permission of the copyright owner.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "First edition, 2025"
    ),
    'fr': (
        "\u00c0 qui est cet \u0153uf ? \u2014 Livre de Coloriage\n"
        "Une aventure de Coco l'Axolotl\n\n"
        "\u00c9crit et illustr\u00e9 par Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. Tous droits r\u00e9serv\u00e9s.\n\n"
        "Aucune partie de cette publication ne peut \u00eatre reproduite,\n"
        "stock\u00e9e dans un syst\u00e8me de r\u00e9cup\u00e9ration, ou transmise\n"
        "sous quelque forme ou par quelque moyen que ce soit,\n"
        "\u00e9lectronique, m\u00e9canique, photocopie, enregistrement ou autre,\n"
        "sans l'autorisation \u00e9crite pr\u00e9alable du titulaire des droits.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "Premi\u00e8re \u00e9dition, 2025"
    ),
    'es': (
        "\u00bfDe qui\u00e9n es este huevo? \u2014 Libro para Colorear\n"
        "Una aventura de Coco el Ajolote\n\n"
        "Escrito e ilustrado por Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. Todos los derechos reservados.\n\n"
        "Ninguna parte de esta publicaci\u00f3n puede ser reproducida,\n"
        "almacenada en un sistema de recuperaci\u00f3n, o transmitida\n"
        "de ninguna forma ni por ning\u00fan medio, electr\u00f3nico,\n"
        "mec\u00e1nico, fotocopia, grabaci\u00f3n u otro,\n"
        "sin el permiso previo por escrito del titular de los derechos.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "Primera edici\u00f3n, 2025"
    ),
}

PAGE_W = Inches(8.5)
PAGE_H = Inches(11)


def extract_texts(pptx_path):
    """Extract text from source PPTX."""
    prs = Presentation(str(pptx_path))
    slides_data = []
    for i, slide in enumerate(prs.slides):
        data = {"index": i, "texts": {}}
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = ""
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    if pi > 0:
                        full_text += "\n"
                    full_text += para.text
                full_text = full_text.replace("\x0b", "\n")
                full_text = re.sub(r'[\x00-\x08\x0c\x0e-\x1f]', '', full_text)
                data["texts"][shape.name] = full_text.strip()
        slides_data.append(data)
    return slides_data


def add_text_box(slide, left, top, width, height, text, size_pt, bold, color,
                 alignment=PP_ALIGN.CENTER):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    bodyPr = tf._txBody.find(f"{a}bodyPr")
    bodyPr.set("wrap", "square")
    etree.SubElement(bodyPr, f"{a}spAutoFit")

    p = tf.paragraphs[0]
    p.alignment = alignment

    # Set default run properties
    pPr = p._element.find(f"{a}pPr")
    if pPr is None:
        pPr = etree.SubElement(p._element, f"{a}pPr")
        p._element.insert(0, pPr)
    defRPr = etree.SubElement(pPr, f"{a}defRPr")
    defRPr.set("sz", str(int(size_pt * 100)))
    defRPr.set("b", "1" if bold else "0")
    solidFill = etree.SubElement(defRPr, f"{a}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{a}srgbClr")
    srgbClr.set("val", str(color))
    latin = etree.SubElement(defRPr, f"{a}latin")
    latin.set("typeface", FONT_NAME)

    lines = text.split("\n")
    for li, line in enumerate(lines):
        if li > 0:
            etree.SubElement(p._element, f"{a}br")
        r = etree.SubElement(p._element, f"{a}r")
        t = etree.SubElement(r, f"{a}t")
        t.text = line

    return txBox


def add_coloring_image(slide, image_path):
    """Add a B&W coloring image to a slide, fitted to page with margin for name."""
    if not image_path or not Path(image_path).exists():
        return None

    img = PILImage.open(image_path)
    img_w, img_h = img.size
    img_ratio = img_w / img_h

    # Leave space at bottom for animal name
    avail_h = Inches(9.5)  # 11 - 0.5 top - 1.0 bottom for name
    avail_w = Inches(7.5)  # 8.5 - 0.5 margins each side
    avail_ratio = avail_w / avail_h

    if img_ratio > avail_ratio:
        # Wider than available: fit width
        new_width = avail_w
        new_height = int(avail_w / img_ratio)
    else:
        # Taller: fit height
        new_height = avail_h
        new_width = int(avail_h * img_ratio)

    left = (PAGE_W - new_width) // 2
    top = Inches(0.3)

    slide.shapes.add_picture(str(image_path), left, top, new_width, new_height)


def build_coloring_book(slides_data, output_path):
    """Build a coloring book PPTX with B&W illustrations.
    Source is a spread-format PPTX (64 or 66 slides):
      64 slides: 0=cover, 1=intro text, 2=intro image,
        3,5,7...59 = animal text pages, 61=ending text, 63=back cover
      66 slides: 0=cover, 1=copyright, 2=intro text, 3=intro image,
        4,6,8...60 = animal text pages, 62=ending text, 65=back cover
    """
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    margin_x = Inches(0.5)
    text_width = PAGE_W - 2 * margin_x

    # Detect format: 66 slides = has copyright page (offset +1)
    has_copyright = len(slides_data) == 66
    offset = 1 if has_copyright else 0

    # --- COVER PAGE ---
    # Detect language from intro text (at index 1 or 2 depending on copyright page)
    intro_idx = 2 if has_copyright else 1
    intro_texts = str(slides_data[intro_idx]["texts"]) if len(slides_data) > intro_idx else ""
    coloring_label = "Coloring Book"
    if "coquille" in intro_texts or "matin" in intro_texts:
        coloring_label = "Livre de Coloriage"
    elif "mañana" in intro_texts or "huevo" in intro_texts or "cascarón" in intro_texts:
        coloring_label = "Libro para Colorear"

    # Get title from cover slide
    cover_texts = slides_data[0]["texts"]
    # Find the title (contains "Egg"/"œuf"/"huevo")
    title = ""
    for key, val in cover_texts.items():
        if any(w in val.lower() for w in ["egg", "oeuf", "œuf", "huevo"]):
            title = val
            break
    if not title:
        title = "Whose Egg Is This?"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(2.0), text_width, Inches(2.0),
                 title, 44, True, PINK_DARK)
    add_text_box(slide, margin_x, Inches(4.5), text_width, Inches(1.0),
                 coloring_label, 30, True, TEXT_DARK)

    # --- COPYRIGHT PAGE ---
    # Detect language for copyright text
    lang = 'en'
    if "Livre de Coloriage" in coloring_label:
        lang = 'fr'
    elif "Libro para Colorear" in coloring_label:
        lang = 'es'
    copyright_text = COPYRIGHT_TEXTS.get(lang, COPYRIGHT_TEXTS['en'])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(3.5), text_width, Inches(5),
                 copyright_text, 12, False, LIGHT_GRAY)

    # --- INTRO COLORING PAGE ---
    intro_coloring = SLIDE_TO_COLORING.get(1)
    if intro_coloring:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_coloring_image(slide, COLORING_DIR / intro_coloring)

    # --- 29 ANIMALS: each gets one coloring page ---
    for animal_num in range(29):
        # In spread format: text page = (3 or 4) + animal_num * 2
        text_slide_idx = (3 + offset) + animal_num * 2
        old_idx = animal_num + 2  # for SLIDE_TO_COLORING mapping

        # Get animal name from the text page
        animal_name = f"Animal {animal_num + 1}"
        if text_slide_idx < len(slides_data):
            animal_texts = slides_data[text_slide_idx]["texts"]
            for key, val in animal_texts.items():
                # Animal name is short and has emoji
                if len(val) < 50 and any(ord(c) > 0x1F00 for c in val):
                    animal_name = val
                    break

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # B&W coloring image
        coloring_name = SLIDE_TO_COLORING.get(old_idx)
        if coloring_name:
            add_coloring_image(slide, COLORING_DIR / coloring_name)

        # Animal name at bottom
        add_text_box(slide, margin_x, Inches(10.0), text_width, Inches(0.7),
                     animal_name, 24, True, PINK_DARK)

    # --- ENDING COLORING PAGE ---
    ending_coloring = SLIDE_TO_COLORING.get(31)
    if ending_coloring:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_coloring_image(slide, COLORING_DIR / ending_coloring)

    # --- BACK COVER ---
    back_idx = len(slides_data) - 1  # last slide
    back_data = slides_data[back_idx]["texts"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for key, val in back_data.items():
        if "enjoy" in val.lower() or "aventure" in val.lower() or "disfrutaste" in val.lower():
            add_text_box(slide, margin_x, Inches(3.0), text_width, Inches(1.0),
                         val, 22, True, PINK_DARK)
        elif "cocotheaxolotl" in val.lower() or "visit" in val.lower():
            add_text_box(slide, margin_x, Inches(4.5), text_width, Inches(3.0),
                         val, 16, False, TEXT_DARK)
        elif "©" in val:
            add_text_box(slide, margin_x, Inches(9.5), text_width, Inches(0.5),
                         val, 10, False, LIGHT_GRAY)

    prs.save(str(output_path))
    return len(prs.slides)


def main():
    sources = {
        "EN": "Whose-Egg-Is-This-EN-spreads.pptx",
        "FR": "Whose-Egg-Is-This-FR-spreads.pptx",
        "ES": "Whose-Egg-Is-This-ES-spreads.pptx",
    }

    # Extract text from spread files (which have the corrected text)
    all_data = {}
    for lang, source_name in sources.items():
        source_path = SCRIPT_DIR / source_name
        if not source_path.exists():
            print(f"SKIP: {source_name} not found")
            continue
        print(f"Extracting text from {source_name}...")
        all_data[lang] = extract_texts(source_path)
        print(f"  Found {len(all_data[lang])} slides")

    outputs = [
        ("EN", "Whose-Egg-Is-This-EN-coloring-8.5x11.pptx"),
        ("FR", "Whose-Egg-Is-This-FR-coloring-8.5x11.pptx"),
        ("ES", "Whose-Egg-Is-This-ES-coloring-8.5x11.pptx"),
    ]

    for lang, output_name in outputs:
        if lang not in all_data:
            continue
        output_path = SCRIPT_DIR / output_name
        print(f"Building {output_name}...")
        num_slides = build_coloring_book(all_data[lang], output_path)
        print(f"  OK: {num_slides} slides created\n")

    print("Done!")


if __name__ == "__main__":
    main()
