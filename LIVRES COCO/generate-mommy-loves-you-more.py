"""
Generate "Mommy Loves You More" — A Coco the Axolotl Picture Book
Format: 8.5 x 8.5 inches (standard square picture book — KDP compatible)
Layout: Illustration FULL PAGE + text overlaid on top (real picture book style)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

W = Inches(8.5)
H = Inches(8.5)

PINK        = RGBColor(0xFF, 0x69, 0xB4)
DARK_PINK   = RGBColor(0xD6, 0x00, 0x6E)
SOFT_PINK   = RGBColor(0xFF, 0xE4, 0xF0)
WARM_GOLD   = RGBColor(0xFF, 0xD7, 0x80)
NIGHT_BLUE  = RGBColor(0x1A, 0x1A, 0x5E)
DEEP_NIGHT  = RGBColor(0x0D, 0x0D, 0x2E)
SOFT_BLUE   = RGBColor(0xE8, 0xF4, 0xFF)
WARM_MORNING= RGBColor(0xFF, 0xF8, 0xE1)
SOFT_GREEN  = RGBColor(0xE8, 0xF5, 0xE9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GREY        = RGBColor(0x99, 0x99, 0x99)
PLAC_BORDER = RGBColor(0xFF, 0xB6, 0xC1)

FONT = 'Baloo 2'

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def full_page_placeholder(slide, brief):
    """Illustration fills the entire slide"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.1), Inches(0.1),
                                   W - Inches(0.2), H - Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xFA)
    shape.line.color.rgb = PLAC_BORDER
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 4
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = brief
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xCC, 0x88, 0xAA)
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return shape

def text_overlay(slide, left, top, width, height, text,
                 size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                 italic=False, bg_color=None):
    """Text box overlaid on the illustration"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return txBox

def page_num(slide, num, color=GREY):
    text_overlay(slide, Inches(3.9), Inches(8.05), Inches(0.7), Inches(0.35),
                 str(num), size=12, color=color, align=PP_ALIGN.CENTER)

def story_slide(prs, page_no, illust_brief,
                narration, narration_top=Inches(6.0),
                narration_size=22, narration_color=WHITE,
                narration_italic=False, narration_align=PP_ALIGN.LEFT,
                narration_bg=None,
                speech=None, speech_top=None,
                speech_size=24, speech_color=WHITE, speech_bold=True):
    """
    Full-page illustration + text overlay
    narration_top: vertical position of text (default = bottom third)
    """
    s = new_slide(prs)
    full_page_placeholder(s, illust_brief)

    # Narration text
    text_overlay(s,
        left=Inches(0.45), top=narration_top,
        width=Inches(7.6), height=Inches(8.5) - narration_top - Inches(0.4),
        text=narration,
        size=narration_size, color=narration_color,
        italic=narration_italic, align=narration_align,
        bg_color=narration_bg)

    # Optional speech line (dialogue, bigger)
    if speech:
        st = speech_top if speech_top else narration_top + Inches(0.9)
        text_overlay(s,
            left=Inches(0.45), top=st,
            width=Inches(7.6), height=Inches(1.5),
            text=speech,
            size=speech_size, color=speech_color,
            bold=speech_bold, align=PP_ALIGN.LEFT)

    page_num(s, page_no, color=narration_color)
    return s


# ══════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ──────────────────────────────────────────
# SLIDE 1 — COVER
# ──────────────────────────────────────────
s = new_slide(prs)
full_page_placeholder(s,
    "[ COVER ILLUSTRATION ]\n\n"
    "Coco (small pink axolotl, eggshell hat, sneakers)\n"
    "reaching up with arms wide open toward Mama.\n"
    "Mama reaching down with arms even wider.\n"
    "Warm golden stars all around. Night sky.\n"
    "Magical, joyful, full of love.")

text_overlay(s, Inches(0.3), Inches(0.35), Inches(7.9), Inches(1.3),
             "Mommy Loves You More",
             size=42, bold=True, color=DARK_PINK, align=PP_ALIGN.CENTER)

text_overlay(s, Inches(0.3), Inches(1.6), Inches(7.9), Inches(0.55),
             "A Coco the Axolotl Adventure",
             size=18, color=PINK, align=PP_ALIGN.CENTER, italic=True)

text_overlay(s, Inches(0.3), Inches(7.6), Inches(7.9), Inches(0.55),
             "By Dr. Anita NIRVENA   |   Illustrations: Loopinky",
             size=13, color=DARK_PINK, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 2 — COPYRIGHT
# ──────────────────────────────────────────
s = new_slide(prs)
bg(s, WHITE)
text_overlay(s, Inches(1.5), Inches(2.5), Inches(5.5), Inches(4.5),
    "Mommy Loves You More\n\n"
    "By Dr. Anita NIRVENA\n"
    "Copyright 2026 by Dr. Anita NIRVENA\n"
    "All rights reserved.\n\n"
    "Illustrations: Loopinky\n\n"
    "No part of this book may be reproduced\n"
    "without written permission.\n\n"
    "First Edition, 2026",
    size=13, color=GREY, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 3 — DEDICATION
# ──────────────────────────────────────────
s = new_slide(prs)
bg(s, SOFT_PINK)
text_overlay(s, Inches(1.0), Inches(3.0), Inches(6.5), Inches(2.5),
    "For every little heart that loves without measure\n"
    "-- and every mama who loves you more.",
    size=22, color=DARK_PINK, align=PP_ALIGN.CENTER, italic=True)

# ──────────────────────────────────────────
# PAGE 1 — Bedtime / Mama tucks Coco in
# ──────────────────────────────────────────
story_slide(prs, page_no=1,
    illust_brief=(
        "[ ILLUSTRATION p.1 ]\n"
        "Mama axolotl gently tucking little Coco\n"
        "into a cozy bed. Soft moonlight through\n"
        "window. Both smiling warmly.\n"
        "Cozy bedroom, soft blues and pinks."
    ),
    narration="Every night, Mommy tucked Coco in --\njust right.\nNot too tight. Not too loose. Just right.",
    narration_top=Inches(6.2),
    narration_size=21, narration_color=NIGHT_BLUE
)

# ──────────────────────────────────────────
# PAGE 2 — First "I love you" exchange
# ──────────────────────────────────────────
story_slide(prs, page_no=2,
    illust_brief=(
        "[ ILLUSTRATION p.2 ]\n"
        "Close-up: Coco looking up at Mama\n"
        "with big loving eyes. Mama leaning in\n"
        "close, whispering with a warm smile.\n"
        "Soft moonlight. Very intimate."
    ),
    narration='Coco looked up and said:',
    narration_top=Inches(5.5),
    narration_size=20, narration_color=NIGHT_BLUE,
    speech='"I love you, Mommy."\n\nMommy smiled: "I love you MORE."',
    speech_top=Inches(6.3),
    speech_size=22, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 3 — Coco wonders
# ──────────────────────────────────────────
story_slide(prs, page_no=3,
    illust_brief=(
        "[ ILLUSTRATION p.3 ]\n"
        "Coco sitting up in bed, finger on chin,\n"
        "thinking deeply. Tiny thought bubbles\n"
        "floating around Coco's head.\n"
        "Playful, curious expression."
    ),
    narration="Coco thought about this.",
    narration_top=Inches(5.8),
    narration_size=22, narration_color=DARK_PINK,
    speech='"More? How much more?"',
    speech_top=Inches(6.6),
    speech_size=24, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 4 — Morning hug
# ──────────────────────────────────────────
story_slide(prs, page_no=4,
    illust_brief=(
        "[ ILLUSTRATION p.4 ]\n"
        "Coco hugging Mama tight from behind\n"
        "in sunny kitchen. Arms wrapped around her.\n"
        "Breakfast on table. Warm golden morning light.\n"
        "Happy and cozy."
    ),
    narration="The next morning, Coco hugged Mommy tight.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=BLACK,
    speech='"I love you as big as my hug!"',
    speech_top=Inches(6.6),
    speech_size=22, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 5 — Mama hugs back even tighter
# ──────────────────────────────────────────
story_slide(prs, page_no=5,
    illust_brief=(
        "[ ILLUSTRATION p.5 ]\n"
        "Mama hugging Coco back even tighter,\n"
        "both laughing. Mama's arms wrap completely\n"
        "around little Coco. Pure joy, warm kitchen."
    ),
    narration="Mommy hugged back -- even tighter.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=BLACK,
    speech='"I love you MORE than that."',
    speech_top=Inches(6.6),
    speech_size=22, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 6 — Park: arms wide
# ──────────────────────────────────────────
story_slide(prs, page_no=6,
    illust_brief=(
        "[ ILLUSTRATION p.6 ]\n"
        "Coco in sunny park, arms stretched as wide\n"
        "as possible, big grin. Bright greens and yellows.\n"
        "Joyful, energetic pose."
    ),
    narration="At the park, Coco stretched arms out wide --\nas wide as they could go.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=BLACK,
    speech='"I love you THIS much!"',
    speech_top=Inches(6.7),
    speech_size=24, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 7 — Mama stretches even wider
# ──────────────────────────────────────────
story_slide(prs, page_no=7,
    illust_brief=(
        "[ ILLUSTRATION p.7 ]\n"
        "Mama stretching arms dramatically wider\n"
        "than Coco, laughing. Arms seem to go\n"
        "beyond the page edges. Funny and playful."
    ),
    narration="Mommy stretched her arms out wider.\nMuch, much wider.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=BLACK,
    speech='"I love you MORE than that."',
    speech_top=Inches(6.7),
    speech_size=22, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 8 — Stars
# ──────────────────────────────────────────
story_slide(prs, page_no=8,
    illust_brief=(
        "[ ILLUSTRATION p.8 ]\n"
        "Coco and Mama sitting outside at night,\n"
        "Coco pointing up at a sky full of stars.\n"
        "Both looking up in wonder. Deep blue, moonlight."
    ),
    narration="That night, Coco pointed to the sky.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=WHITE,
    narration_bg=NIGHT_BLUE,
    speech='"I love you all the way to the stars!"',
    speech_top=Inches(6.6),
    speech_size=21, speech_color=WARM_GOLD
)

# ──────────────────────────────────────────
# PAGE 9 — Mama points past stars
# ──────────────────────────────────────────
story_slide(prs, page_no=9,
    illust_brief=(
        "[ ILLUSTRATION p.9 ]\n"
        "Mama pointing dramatically past the stars,\n"
        "past the moon, into infinite deep space.\n"
        "Both in awe. Deep blues and silvers."
    ),
    narration="Mommy pointed past the stars.\nPast the moon. Past everything.",
    narration_top=Inches(5.6),
    narration_size=20, narration_color=WHITE,
    narration_bg=NIGHT_BLUE,
    speech='"I love you MORE than all of that."',
    speech_top=Inches(6.7),
    speech_size=21, speech_color=WARM_GOLD
)

# ──────────────────────────────────────────
# PAGE 10 — Coco asks the question
# ──────────────────────────────────────────
story_slide(prs, page_no=10,
    illust_brief=(
        "[ ILLUSTRATION p.10 ]\n"
        "Coco looking up at Mama with genuinely\n"
        "curious, innocent eyes under night sky.\n"
        "Quiet, intimate moment. Close-up. Starlight."
    ),
    narration="Coco was quiet for a moment.",
    narration_top=Inches(5.8),
    narration_size=20, narration_color=DARK_PINK,
    speech='"How can you always love me more?"',
    speech_top=Inches(6.6),
    speech_size=21, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 11 — Mama pulls Coco close
# ──────────────────────────────────────────
story_slide(prs, page_no=11,
    illust_brief=(
        "[ ILLUSTRATION p.11 ]\n"
        "Mama sitting down, pulling Coco gently\n"
        "into her arms. Both quiet, stars above.\n"
        "Tender, still moment. Warm golden glow."
    ),
    narration="Mommy sat down.\nPulled Coco close.\n\nAnd said nothing for a little while.",
    narration_top=Inches(5.6),
    narration_size=20, narration_color=WHITE,
    narration_italic=True,
    narration_bg=NIGHT_BLUE
)

# ──────────────────────────────────────────
# PAGE 12 — Before you were born
# ──────────────────────────────────────────
story_slide(prs, page_no=12,
    illust_brief=(
        "[ ILLUSTRATION p.12 — MOST EMOTIONAL ]\n"
        "Dreamlike: Mama holding tiny baby Coco\n"
        "for the very first time. Soft glowing stars\n"
        "and golden light. Magical, tender, beautiful."
    ),
    narration='"Because before you were born,"',
    narration_top=Inches(5.5),
    narration_size=22, narration_color=WARM_GOLD,
    narration_bg=DEEP_NIGHT,
    speech='"I already loved you."',
    speech_top=Inches(6.4),
    speech_size=26, speech_color=WARM_GOLD
)

# ──────────────────────────────────────────
# PAGE 13 — Before your first smile
# ──────────────────────────────────────────
story_slide(prs, page_no=13,
    illust_brief=(
        "[ ILLUSTRATION p.13 ]\n"
        "Three small glowing vignettes like stars:\n"
        "1. Baby Coco's first smile\n"
        "2. Baby Coco saying 'Mama!'\n"
        "3. Baby Coco's first hug. Warm golden glow."
    ),
    narration=(
        '"Before your first smile.\n'
        'Before your first word.\n'
        'Before your first hug.\n\n'
        'My love didn\'t start when I met you.\n'
        'It started before."'
    ),
    narration_top=Inches(3.8),
    narration_size=19, narration_color=WHITE,
    narration_italic=True,
    narration_bg=DEEP_NIGHT
)

# ──────────────────────────────────────────
# PAGE 14 — I will always love you more
# ──────────────────────────────────────────
story_slide(prs, page_no=14,
    illust_brief=(
        "[ ILLUSTRATION p.14 ]\n"
        "Mama looking directly at Coco, eye to eye,\n"
        "holding Coco's small face in her hands.\n"
        "Pure, infinite love. Warm golden close-up."
    ),
    narration=(
        '"So no matter how much you love me...\n\n'
        'I will always love you first.\n\n'
        'And I will always love you more."'
    ),
    narration_top=Inches(4.5),
    narration_size=20, narration_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 15 — Coco thinks a long time
# ──────────────────────────────────────────
story_slide(prs, page_no=15,
    illust_brief=(
        "[ ILLUSTRATION p.15 ]\n"
        "Coco sitting under starry sky, eyes wide,\n"
        "thinking hard. Stars twinkling around.\n"
        "Tiny smile starting to form."
    ),
    narration="Coco thought about this for a long time.\n\nA very long time.",
    narration_top=Inches(6.0),
    narration_size=22, narration_color=DARK_PINK,
    narration_italic=True
)

# ──────────────────────────────────────────
# PAGE 16 — The biggest smile
# ──────────────────────────────────────────
story_slide(prs, page_no=16,
    illust_brief=(
        "[ ILLUSTRATION p.16 ]\n"
        "Coco with the biggest, most radiant smile.\n"
        "Eyes sparkling. One bright star shining above.\n"
        "Coco looks like they just discovered a secret."
    ),
    narration="Then Coco smiled.\n\nThe. Biggest. Smile.",
    narration_top=Inches(6.0),
    narration_size=26, narration_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 17 — Coco wins!
# ──────────────────────────────────────────
story_slide(prs, page_no=17,
    illust_brief=(
        "[ ILLUSTRATION p.17 ]\n"
        "Coco pointing triumphantly at Mama, laughing.\n"
        "Eyes full of mischief and love.\n"
        "Mama looks delighted and won over."
    ),
    narration='"Well then..."',
    narration_top=Inches(5.6),
    narration_size=22, narration_color=DARK_PINK,
    speech='"...I love you MORE, Mommy!"',
    speech_top=Inches(6.5),
    speech_size=26, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 18 — Okay. You win.
# ──────────────────────────────────────────
story_slide(prs, page_no=18,
    illust_brief=(
        "[ ILLUSTRATION p.18 ]\n"
        "Mama laughing out loud, head back, tears of joy.\n"
        "Coco jumping and dancing with delight.\n"
        "Both glowing with happiness."
    ),
    narration="Mommy laughed out loud.",
    narration_top=Inches(5.8),
    narration_size=22, narration_color=BLACK,
    speech='"Okay. You win."',
    speech_top=Inches(6.7),
    speech_size=28, speech_color=DARK_PINK
)

# ──────────────────────────────────────────
# PAGE 19-20 — Big hug + final line
# ──────────────────────────────────────────
story_slide(prs, page_no=19,
    illust_brief=(
        "[ ILLUSTRATION p.19-20 — FINALE ]\n"
        "Mama and Coco in the biggest hug,\n"
        "both laughing and happy.\n"
        "Stars all around. Warm golden light between them.\n"
        "Beautiful, peaceful night. Most beautiful image."
    ),
    narration="And they hugged.",
    narration_top=Inches(5.0),
    narration_size=28, narration_color=WHITE,
    narration_italic=True,
    narration_bg=NIGHT_BLUE,
    speech="Nobody loses\nwhen love is the game.",
    speech_top=Inches(6.2),
    speech_size=24, speech_color=WARM_GOLD
)

# ──────────────────────────────────────────
# BACK COVER
# ──────────────────────────────────────────
s = new_slide(prs)
full_page_placeholder(s,
    "[ BACK COVER ILLUSTRATION ]\n\n"
    "Mama and Coco, small and big, walking home\n"
    "hand in hand under a sky full of stars.\n"
    "Warm light ahead. Peaceful and beautiful.")

text_overlay(s, Inches(0.5), Inches(0.4), Inches(7.5), Inches(1.0),
             "Mommy Loves You More",
             size=30, bold=True, color=DARK_PINK, align=PP_ALIGN.CENTER)

text_overlay(s, Inches(0.8), Inches(6.2), Inches(6.9), Inches(2.0),
    "How much does Mommy love you?\n\n"
    "More than hugs. More than stars. More than everything.\n\n"
    "She loved you before you even knew her name.",
    size=15, color=DARK_PINK, align=PP_ALIGN.CENTER, italic=True)

text_overlay(s, Inches(0.5), Inches(8.1), Inches(7.5), Inches(0.3),
    "Discover more Coco the Axolotl adventures!  |  cocotheaxolotl.org",
    size=10, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "Mommy-Loves-You-More-v2.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"   {len(prs.slides)} slides")
