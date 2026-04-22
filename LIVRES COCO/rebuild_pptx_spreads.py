"""Rebuild PowerPoint books with spread layout: text on left page, full-page illustration on right page.
Text pages have the coloring version as a transparent watermark background."""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import re
import tempfile
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

SCRIPT_DIR = Path(__file__).parent
COLORED_DIR = SCRIPT_DIR / "whose-egg-colored"
COLORING_DIR = SCRIPT_DIR / "whose-egg-coloring"

# Slide index (in OLD pptx) → colored image filename
SLIDE_TO_IMAGE = {
    0: "axolotl -cover-colored.png",  # cover
    1: "coco découvre la coquille.png",
    2: "axolotl-kitten-colored.png",
    3: "axolotl_puffy-colored.png",
    4: "axolotl-rabbit-colored.png",
    5: "axolotl-unicorn-colored.png",
    6: "axolotl-panda-colored.png",
    7: "axolotl-dolphin-colored.png",
    8: "axolotl-foal-colored.png",
    9: "axolotl-butterfly-colored.png",
    10: "axolotl-dinosaur-colored.png",
    11: "axolotl-lion-colored.png",
    12: "axolotl-bear-colored.png",
    13: "axolotl-elephant-colored.png",
    14: "axolotl-giraffe-colored.png",
    15: "axolotl-penguin-colored.png",
    16: "axolotl-koala-colored.png",
    17: "axolotl-fox-colored.png",
    18: "axolotl-owl-colored.png",
    19: "axolotl-turtle-colored.png",
    20: "axolotl-frog-colored.png",
    21: "axolotl-highland-colored.png",
    22: "axolotl-flamingo-colored.png",
    23: "axolotl-sloth-colored.png",
    24: "axolotl-otter-colored.png",
    25: "axolotl-hedgehog-colored.png",
    26: "axolotl-lama-colored.png",
    27: "axolotl-duck-colored.png",
    28: "axolotl-chick-colored.png",
    29: "axolotl-lamb-colored.png",
    30: "axolotl-dragon-colored.png",
    31: "axolotl-baby-colored.png",  # ending
    32: "axolotl-goodbye-colored.png",  # goodbye
    33: None,  # back cover
}

# Slide index (in OLD pptx) → coloring image filename
SLIDE_TO_COLORING = {
    1: "coco-eggshell_coloring .png",
    2: "axolotl-kitten-coloring.png",
    3: "axolotl_puffy-coloring.png",
    4: "axolotl-rabbit-coloring.png",
    5: "axolotl-unicorn-coloring.png",
    6: "axolotl-panda-coloring.png",
    7: "axolotl-dolphin-coloring.png",
    8: "axolotl-foal-coloring.png",
    9: "axolotl-butterfly-coloring.png",
    10: "axolotl-dinosaur-coloring.png",
    11: "axolotl-lion-coloring.png",
    12: "axolotl-bear-coloring.png",
    13: "axolotl-elephant-coloring.png",
    14: "axolotl-giraffe-coloring.png",
    15: "axolotl-penguin-coloring.png",
    16: "axolotl-koala-coloring.png",
    17: "axolotl-fox-coloring.png",
    18: "axolotl-owl-coloring.png",
    19: "axolotl-turtle-coloring.png",
    20: "axolotl-frog-coloring.png",
    21: "axolotl-highland-coloring.png",
    22: "axolotl-flamingo-coloring.png",
    23: "axolotl-sloth-coloring.png",
    24: "axolotl-otter-coloring.png",
    25: "axolotl-hedgehog-coloring.png",
    26: "axolotl-lama-coloring.png",
    27: "axolotl-duck-coloring.png",
    28: "axolotl-chick-coloring.png",
    29: "axolotl-lamb-coloring.png",
    30: "axolotl-dragon-coloring.png",
    31: "axolotl-baby-coloring.png",  # ending
    32: "axolotl-goodbye-coloring.png",  # goodbye
}

# Colors
PINK_DARK = RGBColor(0xD6, 0x00, 0x6E)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
FACT_BG = RGBColor(0xFF, 0xF0, 0xF5)
FACT_BORDER = RGBColor(0xFF, 0xB6, 0xC1)
FACT_TEXT = RGBColor(0x8B, 0x00, 0x4A)
LIGHT_PINK_BG = RGBColor(0xFF, 0xFA, 0xFB)
PLACEHOLDER_BG = RGBColor(0xF5, 0xF5, 0xF5)
PLACEHOLDER_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

FONT_NAME = "Baloo 2"

# Page size: 8.5 x 11 inches (US Letter)
PAGE_W = Inches(8.5)
PAGE_H = Inches(11)

# Copyright page text per language (detected from intro text)
COPYRIGHT_TEXTS = {
    'en': (
        "Whose Egg Is This?\n"
        "A Coco the Axolotl Adventure\n\n"
        "Written and illustrated by Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. All rights reserved.\n\n"
        "No part of this publication may be reproduced, stored in a\n"
        "retrieval system, or transmitted in any form or by any means,\n"
        "electronic, mechanical, photocopying, recording, or otherwise,\n"
        "without the prior written permission of the copyright owner.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "First edition, 2025"
    ),
    'fr': (
        "\u00c0 qui est cet \u0153uf ?\n"
        "Une aventure de Coco l'Axolotl\n\n"
        "\u00c9crit et illustr\u00e9 par Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. Tous droits r\u00e9serv\u00e9s.\n\n"
        "Aucune partie de cette publication ne peut \u00eatre reproduite,\n"
        "stock\u00e9e dans un syst\u00e8me de r\u00e9cup\u00e9ration, ou transmise\n"
        "sous quelque forme ou par quelque moyen que ce soit,\n"
        "\u00e9lectronique, m\u00e9canique, photocopie, enregistrement ou autre,\n"
        "sans l'autorisation \u00e9crite pr\u00e9alable du titulaire des droits.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "Premi\u00e8re \u00e9dition, 2025"
    ),
    'es': (
        "\u00bfDe qui\u00e9n es este huevo?\n"
        "Una aventura de Coco el Ajolote\n\n"
        "Escrito e ilustrado por Coco the Axolotl\n\n"
        "\u00a9 2025 Coco the Axolotl. Todos los derechos reservados.\n\n"
        "Ninguna parte de esta publicaci\u00f3n puede ser reproducida,\n"
        "almacenada en un sistema de recuperaci\u00f3n, o transmitida\n"
        "de ninguna forma ni por ning\u00fan medio, electr\u00f3nico,\n"
        "mec\u00e1nico, fotocopia, grabaci\u00f3n u otro,\n"
        "sin el permiso previo por escrito del titular de los derechos.\n\n"
        "www.cocotheaxolotl.org\n\n"
        "Primera edici\u00f3n, 2025"
    ),
}


def detect_language(slides_data):
    """Detect language from intro text content."""
    intro_text = str(slides_data[1]["texts"]) if len(slides_data) > 1 else ""
    if "coquille" in intro_text or "matin" in intro_text:
        return 'fr'
    elif "mañana" in intro_text or "cascarón" in intro_text:
        return 'es'
    return 'en'

# Transparency for coloring watermark: 10% opacity = 90% transparent
WATERMARK_OPACITY_PCT = 10  # in percent


def _extract_slide_texts(slide):
    """Extract all text shapes from a single slide."""
    texts = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = shape.name
            full_text = ""
            for pi, para in enumerate(shape.text_frame.paragraphs):
                if pi > 0:
                    full_text += "\n"
                full_text += para.text
            full_text = full_text.replace("\x0b", "\n")
            full_text = re.sub(r'[\x00-\x08\x0c\x0e-\x1f]', '', full_text)
            texts[name] = full_text.strip()
    return texts


def extract_texts(pptx_path):
    """Extract text from a PPTX, always returning 33-entry list (original format).
    Handles 33-slide (original), 64-slide (spread), and 66-slide (spread+copyright) formats."""
    prs = Presentation(str(pptx_path))
    num_slides = len(prs.slides)

    if num_slides == 33:
        # Original format: direct extraction
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slides_data.append({"index": i, "texts": _extract_slide_texts(slide)})
        return slides_data

    elif num_slides == 64:
        # Spread format (no copyright page): remap to original 33-slide structure
        # 0=cover, 1=intro text, 2=intro image,
        # 3,5,...,59 = animal text pages (29),
        # 4,6,...,60 = animal image pages (29),
        # 61=ending text, 62=ending image, 63=back cover
        slides_data = []
        slides_data.append({"index": 0, "texts": _extract_slide_texts(prs.slides[0])})
        slides_data.append({"index": 1, "texts": _extract_slide_texts(prs.slides[1])})
        for animal_num in range(29):
            spread_idx = 3 + animal_num * 2
            slides_data.append({"index": animal_num + 2,
                                "texts": _extract_slide_texts(prs.slides[spread_idx])})
        slides_data.append({"index": 31, "texts": _extract_slide_texts(prs.slides[61])})
        slides_data.append({"index": 32, "texts": _extract_slide_texts(prs.slides[63])})
        return slides_data

    elif num_slides == 66:
        # Spread format with copyright page: remap to original 33-slide structure
        # 0=cover, 1=copyright, 2=intro text, 3=intro image,
        # 4,6,...,60 = animal text pages (29),
        # 5,7,...,61 = animal image pages (29),
        # 62=ending text, 63=ending image, 64=goodbye, 65=back cover
        slides_data = []
        slides_data.append({"index": 0, "texts": _extract_slide_texts(prs.slides[0])})
        # Skip copyright (slide 1), read intro text from slide 2
        slides_data.append({"index": 1, "texts": _extract_slide_texts(prs.slides[2])})
        for animal_num in range(29):
            spread_idx = 4 + animal_num * 2
            slides_data.append({"index": animal_num + 2,
                                "texts": _extract_slide_texts(prs.slides[spread_idx])})
        slides_data.append({"index": 31, "texts": _extract_slide_texts(prs.slides[62])})
        slides_data.append({"index": 32, "texts": _extract_slide_texts(prs.slides[65])})
        return slides_data

    else:
        raise ValueError(f"Unexpected slide count: {num_slides} (expected 33, 64, or 66)")


def set_default_rpr(paragraph, size_pt, bold, color, font_name=FONT_NAME):
    """Set default run properties on a paragraph via XML."""
    pPr = paragraph._element.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
    if pPr is None:
        pPr = etree.SubElement(
            paragraph._element,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        paragraph._element.insert(0, pPr)

    defRPr = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr")
    defRPr.set("sz", str(int(size_pt * 100)))
    defRPr.set("b", "1" if bold else "0")

    solidFill = etree.SubElement(defRPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    srgbClr = etree.SubElement(solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    srgbClr.set("val", str(color))

    latin = etree.SubElement(defRPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
    latin.set("typeface", font_name)


def add_text_box(slide, left, top, width, height, text, size_pt, bold, color,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_NAME, vcenter=False):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    bodyPr = tf._txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
    bodyPr.set("wrap", "square")
    if vcenter:
        bodyPr.set("anchor", "ctr")
    else:
        etree.SubElement(bodyPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}spAutoFit")

    p = tf.paragraphs[0]
    p.alignment = alignment
    set_default_rpr(p, size_pt, bold, color, font_name)

    lines = text.split("\n")
    for li, line in enumerate(lines):
        if li > 0:
            etree.SubElement(p._element, "{http://schemas.openxmlformats.org/drawingml/2006/main}br")
        r = etree.SubElement(p._element, "{http://schemas.openxmlformats.org/drawingml/2006/main}r")
        t = etree.SubElement(r, "{http://schemas.openxmlformats.org/drawingml/2006/main}t")
        t.text = line

    return txBox


def add_fun_fact_box(slide, left, top, width, height, text):
    """Add a rounded rectangle fun fact box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = FACT_BG
    shape.line.color.rgb = FACT_BORDER
    shape.line.width = Pt(2)

    # Make the fun fact box semi-transparent so watermark shows through slightly
    fill_elem = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if fill_elem is not None:
        color_elem = fill_elem.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        if color_elem is not None:
            alpha = etree.SubElement(color_elem, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
            alpha.set("val", "85000")  # 85% opacity for the fun fact box

    tf = shape.text_frame
    tf.word_wrap = True

    bodyPr = tf._txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
    bodyPr.set("wrap", "square")
    bodyPr.set("anchor", "ctr")

    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_default_rpr(p, 22, False, FACT_TEXT)

    lines = text.split("\n")
    for li, line in enumerate(lines):
        if li > 0:
            etree.SubElement(p._element, f"{a}br")
        r = etree.SubElement(p._element, f"{a}r")
        # First line (prefix) is bold
        rPr = etree.SubElement(r, f"{a}rPr")
        rPr.set("sz", str(int(22 * 100)))
        rPr.set("b", "1" if li == 0 else "0")
        solidFill = etree.SubElement(rPr, f"{a}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{a}srgbClr")
        srgbClr.set("val", str(FACT_TEXT))
        latin = etree.SubElement(rPr, f"{a}latin")
        latin.set("typeface", FONT_NAME)
        t = etree.SubElement(r, f"{a}t")
        t.text = line

    return shape


def make_transparent_image(image_path, opacity_pct):
    """Create a transparent version of an image using Pillow. Returns temp file path."""
    img = Image.open(image_path).convert("RGBA")
    # Set alpha channel to opacity_pct% of 255
    alpha_value = int(255 * opacity_pct / 100)
    # Create new alpha: for each pixel, set alpha to opacity_pct% (only where not already transparent)
    r, g, b, a = img.split()
    # Scale existing alpha by our target opacity
    a = a.point(lambda x: int(x * opacity_pct / 100))
    img = Image.merge("RGBA", (r, g, b, a))

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name, "PNG")
    tmp.close()
    return tmp.name


def add_watermark_image(slide, image_path, opacity_pct=WATERMARK_OPACITY_PCT):
    """Add a full-page coloring image as a transparent watermark background.
    Uses Pillow to bake transparency into the image for reliable rendering.
    Must be called BEFORE adding text elements so it stays behind."""
    if not image_path or not Path(image_path).exists():
        return None

    # Create transparent version with Pillow
    transparent_path = make_transparent_image(image_path, opacity_pct)

    # Fit image to page
    img = Image.open(image_path)
    img_w, img_h = img.size
    img_ratio = img_w / img_h
    page_ratio = PAGE_W / PAGE_H

    if img_ratio < page_ratio:
        new_width = PAGE_W
        new_height = int(PAGE_W / img_ratio)
        left = 0
        top = (PAGE_H - new_height) // 2
    else:
        new_height = PAGE_H
        new_width = int(PAGE_H * img_ratio)
        left = (PAGE_W - new_width) // 2
        top = 0

    pic = slide.shapes.add_picture(transparent_path, left, top, new_width, new_height)

    # Clean up temp file
    import os
    os.unlink(transparent_path)

    return pic


def add_image_slide(prs, image_path):
    """Add a full-page image slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    if image_path and Path(image_path).exists():
        img_ratio = 1024 / 1536
        page_ratio = PAGE_W / PAGE_H

        if img_ratio < page_ratio:
            new_width = PAGE_W
            new_height = int(PAGE_W / img_ratio)
            left = 0
            top = (PAGE_H - new_height) // 2
        else:
            new_height = PAGE_H
            new_width = int(PAGE_H * img_ratio)
            left = (PAGE_W - new_width) // 2
            top = 0

        slide.shapes.add_picture(str(image_path), left, top, new_width, new_height)
    else:
        margin = Inches(1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            margin, margin, PAGE_W - 2 * margin, PAGE_H - 2 * margin
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PLACEHOLDER_BG
        shape.line.color.rgb = PLACEHOLDER_BORDER
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_default_rpr(p, 16, False, RGBColor(0x99, 0x99, 0x99))
        r = etree.SubElement(p._element, "{http://schemas.openxmlformats.org/drawingml/2006/main}r")
        t = etree.SubElement(r, "{http://schemas.openxmlformats.org/drawingml/2006/main}t")
        t.text = "[ Illustration ]"

    return slide


def add_light_bg(slide):
    """Add a very light pink background to a text slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_PINK_BG


def get_coloring_path(old_slide_idx):
    """Get the coloring image path for a given old slide index."""
    coloring_name = SLIDE_TO_COLORING.get(old_slide_idx)
    if coloring_name:
        path = COLORING_DIR / coloring_name
        if path.exists():
            return path
    return None


def build_spread_pptx(slides_data, output_path):
    """Build a new PPTX with spread layout from extracted text data."""
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    margin_x = Inches(0.8)
    text_width = PAGE_W - 2 * margin_x

    # --- SLIDE 0: COVER (full page, with title overlay) ---
    cover_data = slides_data[0]["texts"]
    cover_img = SLIDE_TO_IMAGE.get(0)
    if cover_img:
        cover_img = COLORED_DIR / cover_img
    slide = add_image_slide(prs, cover_img)

    title_text = cover_data.get("TextBox 2", "Whose Egg Is This?")
    subtitle_text = cover_data.get("TextBox 3", "")
    author_text = cover_data.get("TextBox 4", "")

    add_text_box(slide, margin_x, Inches(6.5), text_width, Inches(1.5),
                 title_text, 36, True, PINK_DARK)
    if subtitle_text:
        add_text_box(slide, margin_x, Inches(8.0), text_width, Inches(0.7),
                     subtitle_text, 18, False, TEXT_DARK)
    if author_text:
        add_text_box(slide, margin_x, Inches(8.8), text_width, Inches(0.5),
                     author_text, 14, False, TEXT_DARK)

    # --- COPYRIGHT PAGE ---
    lang = detect_language(slides_data)
    copyright_text = COPYRIGHT_TEXTS.get(lang, COPYRIGHT_TEXTS['en'])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(3.5), text_width, Inches(5),
                 copyright_text, 12, False, RGBColor(0xAA, 0xAA, 0xAA))

    # --- SLIDE 1-2: INTRO (text left, image right) ---
    intro_data = slides_data[1]["texts"]
    intro_text = intro_data.get("TextBox 2", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    # Watermark on intro text page
    add_watermark_image(slide, get_coloring_path(1))
    add_text_box(slide, margin_x, Inches(1.0), text_width, Inches(8),
                 intro_text, 24, True, TEXT_DARK, vcenter=True)

    intro_img = SLIDE_TO_IMAGE.get(1)
    if intro_img:
        intro_img = COLORED_DIR / intro_img
    add_image_slide(prs, intro_img)

    # --- ANIMALS (29 spreads) ---
    for old_idx in range(2, 31):
        animal_data = slides_data[old_idx]["texts"]

        animal_name = animal_data.get("TextBox 3", f"Animal {old_idx - 1}")
        story_text = animal_data.get("TextBox 4", "")
        fun_fact = animal_data.get("Rounded Rectangle 5", "")
        # Put "Did you know?" prefix on its own line
        for prefix in ["Did you know? ", "Le savais-tu ? ", "¿Sabías que...? "]:
            if fun_fact.startswith(prefix):
                fun_fact = prefix.strip() + "\n" + fun_fact[len(prefix):]
                break
        page_num = animal_data.get("TextBox 1", "")

        # TEXT PAGE (left)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_light_bg(slide)

        # Coloring watermark (added first so it's behind everything)
        add_watermark_image(slide, get_coloring_path(old_idx))

        # Page number top-right
        if page_num:
            add_text_box(slide, PAGE_W - Inches(1.5), Inches(0.3), Inches(1.2), Inches(0.3),
                         page_num, 10, False, RGBColor(0xAA, 0xAA, 0xAA),
                         alignment=PP_ALIGN.RIGHT)

        # Animal name (big: 36pt)
        add_text_box(slide, margin_x, Inches(1.0), text_width, Inches(1.0),
                     animal_name, 36, True, PINK_DARK)

        # Story text (big: 26pt, bold, vertically centered)
        add_text_box(slide, margin_x, Inches(1.8), text_width, Inches(5.5),
                     story_text, 26, True, TEXT_DARK, vcenter=True)

        # Fun fact box at bottom (16pt)
        if fun_fact:
            add_fun_fact_box(slide, Inches(0.5), Inches(7.8), PAGE_W - Inches(1.0), Inches(2.5),
                             fun_fact)

        # IMAGE PAGE (right, full bleed)
        img_name = SLIDE_TO_IMAGE.get(old_idx)
        img_path = (COLORED_DIR / img_name) if img_name else None
        add_image_slide(prs, img_path)

    # --- ENDING (text + image spread) ---
    ending_data = slides_data[31]["texts"]
    ending_text = ending_data.get("TextBox 2", "")
    ending_end = ending_data.get("TextBox 3", "The End!")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    # Watermark on ending text page
    add_watermark_image(slide, get_coloring_path(31))
    add_text_box(slide, margin_x, Inches(0.5), text_width, Inches(8.5),
                 ending_text, 18, True, TEXT_DARK, vcenter=True)
    add_text_box(slide, margin_x, Inches(9.5), text_width, Inches(1.0),
                 ending_end, 36, True, PINK_DARK)

    ending_img = SLIDE_TO_IMAGE.get(31)
    if ending_img:
        ending_img = COLORED_DIR / ending_img
    add_image_slide(prs, ending_img)

    # --- GOODBYE PAGE (full-page illustration) ---
    goodbye_img = SLIDE_TO_IMAGE.get(32)
    if goodbye_img:
        goodbye_path = COLORED_DIR / goodbye_img
        add_image_slide(prs, goodbye_path)

    # --- BACK COVER ---
    back_data = slides_data[32]["texts"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)

    # Add baby colored image at top of back cover
    baby_img_name = SLIDE_TO_IMAGE.get(31)
    if baby_img_name:
        baby_path = COLORED_DIR / baby_img_name
        if baby_path.exists():
            img = Image.open(baby_path)
            img_w, img_h = img.size
            img_ratio = img_w / img_h
            target_h = Inches(4.0)
            target_w = int(target_h * img_ratio)
            left = (PAGE_W - target_w) // 2
            slide.shapes.add_picture(str(baby_path), left, Inches(0.5), target_w, target_h)

    enjoy_text = back_data.get("TextBox 1", "")
    visit_text = back_data.get("TextBox 2", "")
    copyright_text = back_data.get("TextBox 3", "")

    if enjoy_text:
        add_text_box(slide, margin_x, Inches(5.0), text_width, Inches(1.0),
                     enjoy_text, 22, True, PINK_DARK)
    if visit_text:
        add_text_box(slide, margin_x, Inches(6.5), text_width, Inches(2.5),
                     visit_text, 16, False, TEXT_DARK)
    if copyright_text:
        add_text_box(slide, margin_x, Inches(9.5), text_width, Inches(0.5),
                     copyright_text, 10, False, RGBColor(0xAA, 0xAA, 0xAA))

    prs.save(str(output_path))
    return len(prs.slides)


def add_bw_fun_fact_box(slide, left, top, width, height, text):
    """Add a B&W fun fact box (no pink, just black border and gray fill)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)  # very light gray
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)  # gray border
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True

    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    bodyPr = tf._txBody.find(f"{a}bodyPr")
    bodyPr.set("wrap", "square")
    bodyPr.set("anchor", "ctr")

    BW_TEXT = RGBColor(0x33, 0x33, 0x33)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_default_rpr(p, 22, False, BW_TEXT)

    lines = text.split("\n")
    for li, line in enumerate(lines):
        if li > 0:
            etree.SubElement(p._element, f"{a}br")
        r = etree.SubElement(p._element, f"{a}r")
        rPr = etree.SubElement(r, f"{a}rPr")
        rPr.set("sz", str(int(22 * 100)))
        rPr.set("b", "1" if li == 0 else "0")
        solidFill = etree.SubElement(rPr, f"{a}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{a}srgbClr")
        srgbClr.set("val", str(BW_TEXT))
        latin = etree.SubElement(rPr, f"{a}latin")
        latin.set("typeface", FONT_NAME)
        t = etree.SubElement(r, f"{a}t")
        t.text = line

    return shape


def build_coloring_pptx(slides_data, output_path):
    """Build a coloring book PPTX: same spread layout but with B&W coloring images,
    no watermark, no pink backgrounds — strictly black & white."""
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    BW_DARK = RGBColor(0x33, 0x33, 0x33)
    BW_LIGHT = RGBColor(0xAA, 0xAA, 0xAA)

    margin_x = Inches(0.8)
    text_width = PAGE_W - 2 * margin_x

    # --- COVER ---
    cover_data = slides_data[0]["texts"]
    title_text = cover_data.get("TextBox 2", "Whose Egg Is This?")
    subtitle_text = cover_data.get("TextBox 3", "")
    author_text = cover_data.get("TextBox 4", "")

    intro_text_raw = str(slides_data[1]["texts"]) if len(slides_data) > 1 else ""
    coloring_label = "Coloring Book"
    if "coquille" in intro_text_raw or "matin" in intro_text_raw:
        coloring_label = "Livre de Coloriage"
    elif "mañana" in intro_text_raw or "cascarón" in intro_text_raw:
        coloring_label = "Libro para Colorear"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(2.5), text_width, Inches(2.0),
                 title_text, 44, True, BW_DARK)
    add_text_box(slide, margin_x, Inches(5.0), text_width, Inches(1.0),
                 coloring_label, 30, True, BW_DARK)
    if subtitle_text:
        add_text_box(slide, margin_x, Inches(6.5), text_width, Inches(0.7),
                     subtitle_text, 18, False, BW_DARK)
    if author_text:
        add_text_box(slide, margin_x, Inches(7.5), text_width, Inches(0.5),
                     author_text, 14, False, BW_LIGHT)

    # --- COPYRIGHT PAGE ---
    lang = detect_language(slides_data)
    copyright_text = COPYRIGHT_TEXTS.get(lang, COPYRIGHT_TEXTS['en'])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(3.5), text_width, Inches(5),
                 copyright_text, 12, False, BW_LIGHT)

    # --- INTRO ---
    intro_data = slides_data[1]["texts"]
    intro_text = intro_data.get("TextBox 2", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(1.0), text_width, Inches(8),
                 intro_text, 24, True, BW_DARK, vcenter=True)

    intro_coloring = SLIDE_TO_COLORING.get(1)
    intro_img = (COLORING_DIR / intro_coloring) if intro_coloring else None
    add_image_slide(prs, intro_img)

    # --- ANIMALS (29 spreads) ---
    for old_idx in range(2, 31):
        animal_data = slides_data[old_idx]["texts"]

        animal_name = animal_data.get("TextBox 3", f"Animal {old_idx - 1}")
        story_text = animal_data.get("TextBox 4", "")
        fun_fact = animal_data.get("Rounded Rectangle 5", "")
        for prefix in ["Did you know? ", "Le savais-tu ? ", "¿Sabías que...? "]:
            if fun_fact.startswith(prefix):
                fun_fact = prefix.strip() + "\n" + fun_fact[len(prefix):]
                break
        page_num = animal_data.get("TextBox 1", "")

        # TEXT PAGE (left) — white background, no pink
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if page_num:
            add_text_box(slide, PAGE_W - Inches(1.5), Inches(0.3), Inches(1.2), Inches(0.3),
                         page_num, 10, False, BW_LIGHT,
                         alignment=PP_ALIGN.RIGHT)

        add_text_box(slide, margin_x, Inches(1.0), text_width, Inches(1.0),
                     animal_name, 36, True, BW_DARK)
        add_text_box(slide, margin_x, Inches(1.8), text_width, Inches(5.5),
                     story_text, 26, True, BW_DARK, vcenter=True)

        if fun_fact:
            add_bw_fun_fact_box(slide, Inches(0.5), Inches(7.8), PAGE_W - Inches(1.0), Inches(2.5),
                                fun_fact)

        # IMAGE PAGE (right) — coloring (B&W) image
        coloring_name = SLIDE_TO_COLORING.get(old_idx)
        img_path = (COLORING_DIR / coloring_name) if coloring_name else None
        add_image_slide(prs, img_path)

    # --- ENDING ---
    ending_data = slides_data[31]["texts"]
    ending_text = ending_data.get("TextBox 2", "")
    ending_end = ending_data.get("TextBox 3", "The End!")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, margin_x, Inches(0.5), text_width, Inches(8.5),
                 ending_text, 18, True, BW_DARK, vcenter=True)
    add_text_box(slide, margin_x, Inches(9.5), text_width, Inches(1.0),
                 ending_end, 36, True, BW_DARK)

    ending_coloring = SLIDE_TO_COLORING.get(31)
    ending_img = (COLORING_DIR / ending_coloring) if ending_coloring else None
    add_image_slide(prs, ending_img)

    # --- GOODBYE PAGE (full-page coloring illustration) ---
    goodbye_coloring = SLIDE_TO_COLORING.get(32)
    if goodbye_coloring:
        goodbye_path = COLORING_DIR / goodbye_coloring
        add_image_slide(prs, goodbye_path)

    # --- BACK COVER ---
    back_data = slides_data[32]["texts"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    enjoy_text = back_data.get("TextBox 1", "")
    visit_text = back_data.get("TextBox 2", "")
    copyright_text = back_data.get("TextBox 3", "")

    if enjoy_text:
        add_text_box(slide, margin_x, Inches(5.0), text_width, Inches(1.0),
                     enjoy_text, 22, True, BW_DARK)
    if visit_text:
        add_text_box(slide, margin_x, Inches(6.5), text_width, Inches(2.5),
                     visit_text, 16, False, BW_DARK)
    if copyright_text:
        add_text_box(slide, margin_x, Inches(9.5), text_width, Inches(0.5),
                     copyright_text, 10, False, BW_LIGHT)

    prs.save(str(output_path))
    return len(prs.slides)


def main():
    # Sources to read text from (use 8.5x11-spreads which have correct text)
    sources = {
        "EN": "Whose-Egg-Is-This-EN-8.5x11-spreads.pptx",
        "FR": "Whose-Egg-Is-This-FR-8.5x11-spreads.pptx",
        "ES": "Whose-Egg-Is-This-ES-8.5x11-spreads.pptx",
    }

    # Step 1: Extract text from all sources first
    all_data = {}
    for lang, source_name in sources.items():
        source_path = SCRIPT_DIR / source_name
        if not source_path.exists():
            print(f"SKIP: {source_name} not found")
            continue
        print(f"Extracting text from {source_name}...")
        all_data[lang] = extract_texts(source_path)
        print(f"  Found {len(all_data[lang])} slides")

    # Step 2: Build all outputs (safe to overwrite sources now)
    outputs = [
        ("EN", "Whose-Egg-Is-This-EN-spreads.pptx"),
        ("FR", "Whose-Egg-Is-This-FR-spreads.pptx"),
        ("ES", "Whose-Egg-Is-This-ES-spreads.pptx"),
        ("EN", "Whose-Egg-Is-This-EN-8.5x11.pptx"),
        ("FR", "Whose-Egg-Is-This-FR-8.5x11.pptx"),
        ("ES", "Whose-Egg-Is-This-ES-8.5x11.pptx"),
        ("EN", "Whose-Egg-Is-This-EN.pptx"),
        ("FR", "Whose-Egg-Is-This-FR.pptx"),
        ("ES", "Whose-Egg-Is-This-ES.pptx"),
    ]

    for lang, output_name in outputs:
        if lang not in all_data:
            continue
        output_path = SCRIPT_DIR / output_name
        print(f"Building {output_name}...")
        num_slides = build_spread_pptx(all_data[lang], output_path)
        print(f"  OK: {num_slides} slides created\n")

    # Step 3: Build coloring book versions
    coloring_outputs = [
        ("EN", "Whose-Egg-Is-This-EN-coloring-8.5x11.pptx"),
        ("FR", "Whose-Egg-Is-This-FR-coloring-8.5x11.pptx"),
        ("ES", "Whose-Egg-Is-This-ES-coloring-8.5x11.pptx"),
    ]

    for lang, output_name in coloring_outputs:
        if lang not in all_data:
            continue
        output_path = SCRIPT_DIR / output_name
        print(f"Building {output_name}...")
        num_slides = build_coloring_pptx(all_data[lang], output_path)
        print(f"  OK: {num_slides} slides created\n")

    print("Done!")


if __name__ == "__main__":
    main()
