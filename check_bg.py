import sys, io, os, hashlib
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

path = os.path.join("C:\\Users\\33612\\Documents\\AXOLOTL", "Coco Can't Sleep Tonight!", "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads.pptx")
prs = Presentation(path)

# Check the background image (first image) on text slides
text_slides = [2, 3, 5, 7, 9, 11, 13, 15, 17, 19, 21, 23, 25, 27, 29]
illust_slides = [4, 6, 8, 10, 12, 14, 16, 18, 20, 22, 24, 26, 28, 30]

print("=== TEXT SLIDES (background images) ===")
bg_hashes = {}
for si in text_slides:
    slide = prs.slides[si - 1]
    for shape in slide.shapes:
        try:
            if shape.image:
                h = hashlib.md5(shape.image.blob).hexdigest()
                sz = len(shape.image.blob) // 1024
                print(f"  Slide {si}: {h} ({sz}KB)")
                bg_hashes[h] = bg_hashes.get(h, 0) + 1
                break  # first image only
        except:
            pass

print()
print("=== Background image reuse ===")
for h, count in sorted(bg_hashes.items(), key=lambda x: -x[1]):
    print(f"  {h}: used {count} times")

print()
print("=== ILLUSTRATION SLIDES ===")
for si in illust_slides:
    slide = prs.slides[si - 1]
    for shape in slide.shapes:
        try:
            if shape.image:
                h = hashlib.md5(shape.image.blob).hexdigest()[:10]
                sz = len(shape.image.blob) // 1024
                print(f"  Slide {si}: {h} ({sz}KB)")
                break
        except:
            pass
