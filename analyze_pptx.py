import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu

path = os.path.join("C:\\Users\\33612\\Documents\\AXOLOTL", "Coco Can't Sleep Tonight!", "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads.pptx")
prs = Presentation(path)

slide_w = prs.slide_width
slide_h = prs.slide_height
mid_x = slide_w // 2

print(f"Slide size: {slide_w/914400:.2f} x {slide_h/914400:.2f} inches")
print(f"Midpoint X: {mid_x/914400:.2f} inches")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    for shape in slide.shapes:
        left_in = shape.left / 914400
        top_in = shape.top / 914400
        w_in = shape.width / 914400
        h_in = shape.height / 914400
        right_in = left_in + w_in

        # Determine if left half, right half, or spanning
        if right_in <= mid_x / 914400 + 0.1:
            side = "LEFT"
        elif left_in >= mid_x / 914400 - 0.1:
            side = "RIGHT"
        else:
            side = "FULL-SPAN"

        shape_type = "TEXT" if shape.has_text_frame else "OTHER"
        try:
            if shape.image:
                shape_type = f"IMAGE ({len(shape.image.blob)//1024}KB)"
        except:
            pass

        text_preview = ""
        if shape.has_text_frame:
            texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                text_preview = f' -> "{texts[0][:60]}..."' if len(texts[0]) > 60 else f' -> "{texts[0]}"'

        print(f"  [{side}] {shape_type} at ({left_in:.1f}, {top_in:.1f}) size ({w_in:.1f} x {h_in:.1f}){text_preview}")
    print()
