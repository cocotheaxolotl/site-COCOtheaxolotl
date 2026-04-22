"""
Convert 'Coco ne dort pas ce soir!' PPTX to EPUB
- Starry night background behind text pages
- Illustrations as separate full pages
- No emojis (they render as rectangles in most ebook readers)
"""
import sys, io, os, zipfile, uuid, hashlib
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# --- CONFIG ---
PPTX_PATH = os.path.join(
    "C:\\Users\\33612\\Documents\\AXOLOTL",
    "Coco Can't Sleep Tonight!",
    "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads.pptx"
)
OUTPUT_DIR = os.path.join(
    "C:\\Users\\33612\\Documents\\AXOLOTL",
    "Coco Can't Sleep Tonight!"
)
EPUB_FILENAME = "Coco-ne-dort-pas-ce-soir.epub"

TITLE = "Coco ne dort pas ce soir !"
SUBTITLE = "Une aventure de Coco l'Axolotl"
AUTHOR = "Dr. Anita NIRVENA"
ILLUSTRATOR = "Loopinky"
LANGUAGE = "fr"
PUBLISHER = "cocotheaxolotl.org"
DATE = "2026"
DESCRIPTION = "Un livre pour enfants sur le sommeil, avec Coco l'Axolotl qui decouvre comment dorment les differents animaux."

# --- LOAD PPTX ---
print(f"Loading: {PPTX_PATH}")
prs = Presentation(PPTX_PATH)
print(f"Slides: {len(prs.slides)}")

# --- HELPERS ---
all_images = {}

def image_ext(content_type):
    return {'image/png': '.png', 'image/jpeg': '.jpg', 'image/gif': '.gif'}.get(content_type, '.png')

def add_image(content_type, blob):
    h = hashlib.md5(blob).hexdigest()[:10]
    ext = image_ext(content_type)
    fname = f"img_{h}{ext}"
    all_images[fname] = blob
    return fname

def get_slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts

def get_slide_images(slide):
    images = []
    for shape in slide.shapes:
        try:
            if shape.image:
                images.append((shape.image.content_type, shape.image.blob))
        except:
            pass
    return images

def escape_html(text):
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

def strip_emoji(text):
    """Remove leading emoji + whitespace from text."""
    import re
    # Remove common emoji at start of string
    return re.sub(r'^[\U0001f300-\U0001f9ff\U00002702-\U000027b0\U0000fe0f\u200d]+\s*', '', text)

def format_story_text(raw_text):
    lines = raw_text.replace('\x0b', '\n').split('\n')
    result = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        escaped = escape_html(line)
        if escaped.startswith('\u00ab') or escaped.startswith('&laquo;'):
            result.append(f'<p class="dialogue">{escaped}</p>')
        else:
            result.append(f'<p>{escaped}</p>')
    return '\n'.join(result)

def format_fun_fact(text):
    escaped = escape_html(text)
    if 'Le savais-tu' in escaped:
        parts = escaped.split('?', 1)
        if len(parts) == 2:
            return f'''<div class="fun-fact">
<p class="fun-fact-title">Le savais-tu ?</p>
<p>{parts[1].strip()}</p>
</div>'''
    return f'<div class="fun-fact"><p>{escaped}</p></div>'

# --- EXTRACT DATA ---
slides_data = []
for i, slide in enumerate(prs.slides):
    imgs = get_slide_images(slide)
    img_files = [add_image(ct, blob) for ct, blob in imgs]
    slides_data.append({
        'index': i + 1,
        'texts': get_slide_texts(slide),
        'images': imgs,
        'image_files': img_files,
    })

# The background image is the same on all text slides - identify it
bg_image_file = slides_data[2]['image_files'][0]  # slide 3's first image
print(f"Background image: {bg_image_file}")

# Cover image (slide 1, first image)
cover_image_file = slides_data[0]['image_files'][0] if slides_data[0]['image_files'] else None
# Second cover image (Coco)
cover_coco_file = slides_data[0]['image_files'][1] if len(slides_data[0]['image_files']) > 1 else None

print(f"Total unique images: {len(all_images)}")

# --- CHAPTER DEFINITIONS ---
# Structure: text page (odd slide with bg) + illustration page (even slide)
# No emojis in titles
chapter_defs = [
    {"type": "cover", "title": "Couverture", "slides": [1]},
    {"type": "title_page", "title": "Page de titre", "slides": [2]},
    {"type": "spread", "title": "L'heure de dormir", "text_slide": 3, "illust_slide": 4},
    {"type": "spread", "title": "L'aventure commence", "text_slide": 5, "illust_slide": 6},
    {"type": "spread", "title": "Le Hibou", "text_slide": 7, "illust_slide": 8},
    {"type": "spread", "title": "Le Dauphin", "text_slide": 9, "illust_slide": 10},
    {"type": "spread", "title": "La Loutre", "text_slide": 11, "illust_slide": 12},
    {"type": "spread", "title": "Le Koala", "text_slide": 13, "illust_slide": 14},
    {"type": "spread", "title": "La Chauve-Souris", "text_slide": 15, "illust_slide": 16},
    {"type": "spread", "title": "Le Chat", "text_slide": 17, "illust_slide": 18},
    {"type": "spread", "title": "Le Cheval", "text_slide": 19, "illust_slide": 20},
    {"type": "spread", "title": "Le Flamant Rose", "text_slide": 21, "illust_slide": 22},
    {"type": "spread", "title": "Le Paresseux", "text_slide": 23, "illust_slide": 24},
    {"type": "spread", "title": "Maman Axolotl", "text_slide": 25, "illust_slide": 26},
    {"type": "spread", "title": "Le super-pouvoir", "text_slide": 27, "illust_slide": 28},
    {"type": "spread", "title": "Bonne nuit, Coco", "text_slide": 29, "illust_slide": 30},
]

# --- BUILD XHTML ---

def build_cover_xhtml():
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Couverture</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body>
<div class="cover-page">
<img src="images/{cover_image_file}" alt="Couverture" class="cover-img"/>
</div>
</body>
</html>'''

def build_title_page_xhtml():
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Page de titre</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body>
<div class="text-page">
<img src="images/{bg_image_file}" alt="" class="bg-img"/>
<div class="text-content title-page">
<h1>{escape_html(TITLE)}</h1>
<p class="subtitle">{escape_html(SUBTITLE)}</p>
<p class="author">Par : {escape_html(AUTHOR)}</p>
<p class="illustrator">Illustrations : {escape_html(ILLUSTRATOR)}</p>
<p class="copyright">\u00a9 {DATE} {escape_html(AUTHOR)}<br/>Tous droits r\u00e9serv\u00e9s.<br/>{escape_html(PUBLISHER)}<br/>Premi\u00e8re \u00e9dition, {DATE}</p>
</div>
</div>
</body>
</html>'''

def build_text_page_xhtml(ch_def):
    """Build the text page with starry night background."""
    sd = slides_data[ch_def['text_slide'] - 1]
    title = ch_def['title']

    body_parts = []
    body_parts.append(f'<h2 class="chapter-title">{escape_html(title)}</h2>')

    for text in sd['texts']:
        cleaned = strip_emoji(text)
        if not cleaned:
            continue
        # Skip if it's just the animal name (already in chapter title)
        if cleaned == title or cleaned.strip() == title.strip():
            continue
        if 'Le savais-tu' in cleaned:
            body_parts.append(format_fun_fact(cleaned))
        else:
            body_parts.append(format_story_text(cleaned))

    content = '\n'.join(body_parts)
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>{escape_html(title)}</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body>
<div class="text-page">
<img src="images/{bg_image_file}" alt="" class="bg-img"/>
<div class="text-content">
{content}
</div>
</div>
</body>
</html>'''

def build_illust_page_xhtml(ch_def):
    """Build the illustration-only page."""
    sd = slides_data[ch_def['illust_slide'] - 1]
    title = ch_def['title']
    # Get the illustration image (not the background)
    img_file = None
    for fname in sd['image_files']:
        if fname != bg_image_file:
            img_file = fname
            break
    if not img_file and sd['image_files']:
        img_file = sd['image_files'][0]

    return f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>{escape_html(title)} - Illustration</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body>
<div class="illust-page">
<img src="images/{img_file}" alt="{escape_html(title)}"/>
</div>
</body>
</html>'''

# --- CSS ---
CSS = '''@charset "UTF-8";

body {
    margin: 0;
    padding: 0;
    font-family: "Georgia", "Palatino", serif;
    line-height: 1.6;
    color: #2c2c2c;
    background-color: #fff;
}

/* Text page */
.text-page {
    page-break-after: always;
}

.bg-img {
    width: 100%;
    height: auto;
    display: block;
}

.text-content {
    text-align: center;
    color: #2c2c2c;
    padding: 1.5em 1.5em;
}

.text-content p {
    color: #2c2c2c;
    font-size: 1.05em;
    margin: 0.6em 0;
    text-align: center;
}

.text-content .dialogue {
    font-style: italic;
    color: #5b3a7a;
    text-align: center;
}

h1 {
    font-size: 1.8em;
    text-align: center;
    color: #e75480;
    margin-top: 0.5em;
    margin-bottom: 0.3em;
}

h2.chapter-title {
    font-size: 1.5em;
    text-align: center;
    color: #e75480;
    margin-top: 0.5em;
    margin-bottom: 0.8em;
}

.subtitle {
    text-align: center;
    font-style: italic;
    font-size: 1.1em;
    color: #7b68ae;
}

.author, .illustrator {
    text-align: center;
    font-size: 1em;
    color: #555;
}

.copyright {
    text-align: center;
    font-size: 0.85em;
    color: #888;
    margin-top: 2em;
}

.title-page {
    text-align: center;
    padding-top: 3em;
}

/* Cover page */
.cover-page {
    text-align: center;
    margin: 0;
    padding: 0;
    page-break-after: always;
}

.cover-img {
    max-width: 100%;
    height: auto;
    display: block;
    margin: 0 auto;
}

/* Illustration page */
.illust-page {
    text-align: center;
    margin: 0;
    padding: 0;
    page-break-after: always;
}

.illust-page img {
    max-width: 100%;
    height: auto;
    display: block;
    margin: 0 auto;
}

/* Fun fact box */
.fun-fact {
    background-color: #fef9e7;
    border: 2px solid #f7c948;
    border-radius: 10px;
    padding: 0.8em 1em;
    margin: 1em 0;
    page-break-inside: avoid;
}

.fun-fact p {
    color: #333333;
    font-size: 0.95em;
    text-align: center;
}

.fun-fact-title {
    font-weight: bold;
    color: #e67e22;
    font-size: 1.1em;
    margin-bottom: 0.3em;
}
'''

# --- BUILD EPUB ---
book_uuid = str(uuid.uuid4())

# Generate all XHTML files
xhtml_files = []  # (filename, content, title)

for i, ch_def in enumerate(chapter_defs):
    if ch_def['type'] == 'cover':
        xhtml_files.append((f"page_{i:02d}_cover.xhtml", build_cover_xhtml(), ch_def['title']))
    elif ch_def['type'] == 'title_page':
        xhtml_files.append((f"page_{i:02d}_title.xhtml", build_title_page_xhtml(), ch_def['title']))
    elif ch_def['type'] == 'spread':
        # Text page first (left), then illustration (right)
        xhtml_files.append((f"page_{i:02d}_text.xhtml", build_text_page_xhtml(ch_def), ch_def['title']))
        xhtml_files.append((f"page_{i:02d}_illust.xhtml", build_illust_page_xhtml(ch_def), f"{ch_def['title']} - Illustration"))

print(f"Generated {len(xhtml_files)} XHTML pages")

# content.opf
def build_opf():
    manifest_items = []
    spine_items = []

    manifest_items.append('<item id="style" href="style.css" media-type="text/css"/>')

    for i, (fname, _, _) in enumerate(xhtml_files):
        item_id = f"page_{i:02d}"
        manifest_items.append(f'<item id="{item_id}" href="{fname}" media-type="application/xhtml+xml"/>')
        spine_items.append(f'<itemref idref="{item_id}"/>')

    for img_fname in sorted(all_images.keys()):
        img_id = img_fname.replace('.', '_').replace('-', '_')
        mt = 'image/png' if img_fname.endswith('.png') else 'image/jpeg'
        manifest_items.append(f'<item id="{img_id}" href="images/{img_fname}" media-type="{mt}"/>')

    cover_meta = ''
    if cover_image_file:
        cover_id = cover_image_file.replace('.', '_').replace('-', '_')
        cover_meta = f'<meta name="cover" content="{cover_id}"/>'

    manifest_items.append('<item id="ncx" href="toc.ncx" media-type="application/x-dtbncx+xml"/>')
    manifest_items.append('<item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>')

    return f'''<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid" version="3.0">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="bookid">urn:uuid:{book_uuid}</dc:identifier>
    <dc:title>{escape_html(TITLE)}</dc:title>
    <dc:creator>{escape_html(AUTHOR)}</dc:creator>
    <dc:language>{LANGUAGE}</dc:language>
    <dc:publisher>{escape_html(PUBLISHER)}</dc:publisher>
    <dc:date>{DATE}</dc:date>
    <dc:description>{escape_html(DESCRIPTION)}</dc:description>
    <meta property="dcterms:modified">2026-02-19T00:00:00Z</meta>
    {cover_meta}
  </metadata>
  <manifest>
    {chr(10).join("    " + m for m in manifest_items)}
  </manifest>
  <spine toc="ncx">
    {chr(10).join("    " + s for s in spine_items)}
  </spine>
</package>'''

def build_ncx():
    nav_points = []
    # Only add text pages and cover to TOC (not illustration pages)
    play_order = 0
    for i, (fname, _, title) in enumerate(xhtml_files):
        if 'illust' not in fname:
            play_order += 1
            nav_points.append(f'''    <navPoint id="navpoint-{play_order}" playOrder="{play_order}">
      <navLabel><text>{escape_html(title)}</text></navLabel>
      <content src="{fname}"/>
    </navPoint>''')
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">
  <head>
    <meta name="dtb:uid" content="urn:uuid:{book_uuid}"/>
    <meta name="dtb:depth" content="1"/>
    <meta name="dtb:totalPageCount" content="0"/>
    <meta name="dtb:maxPageNumber" content="0"/>
  </head>
  <docTitle><text>{escape_html(TITLE)}</text></docTitle>
  <navMap>
{chr(10).join(nav_points)}
  </navMap>
</ncx>'''

def build_nav():
    items = []
    for fname, _, title in xhtml_files:
        if 'illust' not in fname:
            items.append(f'      <li><a href="{fname}">{escape_html(title)}</a></li>')
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Table des matieres</title></head>
<body>
<nav epub:type="toc" id="toc">
  <h1>Table des matieres</h1>
  <ol>
{chr(10).join(items)}
  </ol>
</nav>
</body>
</html>'''

# --- WRITE EPUB ---
epub_path = os.path.join(OUTPUT_DIR, EPUB_FILENAME)
print(f"Writing EPUB to: {epub_path}")

with zipfile.ZipFile(epub_path, 'w', zipfile.ZIP_DEFLATED) as zf:
    zf.writestr('mimetype', 'application/epub+zip', compress_type=zipfile.ZIP_STORED)

    zf.writestr('META-INF/container.xml', '''<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>''')

    zf.writestr('OEBPS/content.opf', build_opf())
    zf.writestr('OEBPS/toc.ncx', build_ncx())
    zf.writestr('OEBPS/nav.xhtml', build_nav())
    zf.writestr('OEBPS/style.css', CSS)

    for fname, content, _ in xhtml_files:
        zf.writestr(f'OEBPS/{fname}', content)

    for img_fname, img_blob in all_images.items():
        zf.writestr(f'OEBPS/images/{img_fname}', img_blob)

print(f"\nDone! EPUB created: {epub_path}")
print(f"  - {len(xhtml_files)} pages")
print(f"  - {len(all_images)} unique images")
print(f"  - File size: {os.path.getsize(epub_path) / 1024 / 1024:.1f} MB")
