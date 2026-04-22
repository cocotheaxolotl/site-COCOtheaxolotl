import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

path = os.path.join("C:\\Users\\33612\\Documents\\AXOLOTL", "Coco Can't Sleep Tonight!", "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads.pptx")
print("Path:", path)
print("Exists:", os.path.exists(path))

prs = Presentation(path)
print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width/914400:.1f} x {prs.slide_height/914400:.1f} inches")
print()

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    images = []
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
        try:
            if shape.image:
                images.append(f"{shape.image.content_type}, {len(shape.image.blob)//1024}KB")
        except:
            pass
    for t in texts:
        print(f"  TEXT: {t}")
    for img in images:
        print(f"  IMG: {img}")
    print()
