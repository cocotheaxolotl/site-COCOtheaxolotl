import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu

path = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
prs = Presentation(path)

for si in [10, 11, 12, 39, 40]:
    slide = prs.slides[si - 1]
    print(f"=== Slide {si} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            print(f"  TextFrame: pos=({shape.left/914400:.2f}, {shape.top/914400:.2f}) size=({shape.width/914400:.2f} x {shape.height/914400:.2f})")
            print(f"  Word wrap: {tf.word_wrap}")
            for pi, para in enumerate(tf.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                print(f"  Para {pi}: '{text}'")
                print(f"    Alignment: {para.alignment}")
                for ri, run in enumerate(para.runs):
                    f = run.font
                    try:
                        color_val = f.color.rgb
                    except:
                        try:
                            color_val = f"type={f.color.type}"
                        except:
                            color_val = "inherited"
                    print(f"    Run {ri}: '{run.text}' font={f.name} size={f.size} bold={f.bold} italic={f.italic} color={color_val}")
    print()

# Check slides 33, 69 for text shape existence
print("=== Slides that need text ===")
for si in [33, 34, 69, 70, 84]:
    slide = prs.slides[si - 1]
    print(f"Slide {si}: {len(slide.shapes)} shapes")
    for shape in slide.shapes:
        stype = "IMG" if not shape.has_text_frame else "TXT"
        try:
            if shape.image:
                stype = f"IMG({len(shape.image.blob)//1024}KB)"
        except:
            pass
        print(f"  {stype}: name='{shape.name}' pos=({shape.left/914400:.2f},{shape.top/914400:.2f}) size=({shape.width/914400:.2f}x{shape.height/914400:.2f})")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    print(f"    Text: '{para.text.strip()}'")
                    for run in para.runs:
                        try:
                            color_val = run.font.color.rgb
                        except:
                            color_val = "inherited"
                        print(f"    font={run.font.name} size={run.font.size} color={color_val}")
