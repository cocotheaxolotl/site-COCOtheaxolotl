"""
Update animal names in the wild portraits coloring book PPTX.
- Slides 33-38, 84: ADD text box with name (none exists)
- Slides 69-83: REPLACE "lapin" with correct name
- Slide 47: fix typo "chimpanfé" -> "chimpanzé"
- Font: Baguet Script, 18pt, white, centered
"""
import sys, io, os, copy, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

SRC = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
DST = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"

# Backup first
BACKUP = SRC.replace('.pptx', '_backup.pptx')
if not os.path.exists(BACKUP):
    shutil.copy2(SRC, BACKUP)
    print(f"Backup saved: {BACKUP}")
else:
    print(f"Backup already exists: {BACKUP}")

prs = Presentation(SRC)

# Text box position and size (matching existing slides)
TX_LEFT = Inches(3.25)
TX_TOP = Inches(9.83)
TX_WIDTH = Inches(2.15)
TX_HEIGHT = Inches(0.40)

# --- NAMES TO UPDATE ---

# Slides that need a NEW text box (no "ZoneTexte 3" exists)
add_names = {
    33: "girafe",
    34: "éléphant",
    35: "éléphanteau",
    36: "axolotl",
    37: "buffle",
    38: "boeuf musqué",
    84: "béluga",
}

# Slides where "lapin" needs to be REPLACED
replace_names = {
    69: "toucan",
    70: "flamant rose",
    71: "paon",
    72: "macareux",
    73: "pélican",
    74: "cigogne",
    75: "grue couronnée",
    76: "colibri",
    77: "canard mandarin",
    78: "faisan doré",
    79: "cygne",
    80: "martin-pêcheur",
    81: "pic-vert",
    82: "épervier",
    83: "ara",
}

# Typo fixes
typo_fixes = {
    47: ("chimpanfé", "chimpanzé"),
}

def set_text_style(run, text):
    """Apply Baguet Script 18pt white to a run."""
    run.text = text
    run.font.name = "Baguet Script"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = False
    run.font.italic = False

def add_text_box(slide, text):
    """Add a new text box with animal name."""
    txBox = slide.shapes.add_textbox(TX_LEFT, TX_TOP, TX_WIDTH, TX_HEIGHT)
    txBox.name = "ZoneTexte 3"
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    set_text_style(run, text)
    return txBox

# --- PROCESS ---

changes = 0

# 1. Add new text boxes
for slide_num, name in add_names.items():
    slide = prs.slides[slide_num - 1]
    add_text_box(slide, name)
    print(f"  Slide {slide_num}: ADDED '{name}'")
    changes += 1

# 2. Replace "lapin" text
for slide_num, name in replace_names.items():
    slide = prs.slides[slide_num - 1]
    found = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "ZoneTexte 3":
            # Clear existing runs and set new text
            para = shape.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            # Remove existing runs
            for run in para.runs:
                run.text = ""
            # Set first run or add new one
            if para.runs:
                set_text_style(para.runs[0], name)
            else:
                run = para.add_run()
                set_text_style(run, name)
            found = True
            print(f"  Slide {slide_num}: 'lapin' -> '{name}'")
            changes += 1
            break
    if not found:
        # No ZoneTexte 3, add one
        add_text_box(slide, name)
        print(f"  Slide {slide_num}: ADDED '{name}' (no existing text box)")
        changes += 1

# 3. Fix typos
for slide_num, (old_text, new_text) in typo_fixes.items():
    slide = prs.slides[slide_num - 1]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if old_text in para.text:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            print(f"  Slide {slide_num}: '{old_text}' -> '{new_text}' (typo fix)")
                            changes += 1

# --- SAVE ---
prs.save(DST)
print(f"\nDone! {changes} changes saved to: {DST}")
