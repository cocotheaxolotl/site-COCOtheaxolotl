"""
Convert 'Coco ne dort pas ce soir!' PPTX to EPUB
Text pages = composite images (text rendered on starry night with Pillow)
Illustration pages = standalone images
Works on ALL ebook readers.
"""
import sys, io, os, zipfile, uuid, hashlib, textwrap, re
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# --- CONFIG ---
PPTX_PATH = os.path.join(
    "C:\\Users\\33612\\Documents\\AXOLOTL",
    "Coco Can't Sleep Tonight!",
    "Coco-ne-dort-pas-ce-soir-FR-8.5x11-spreads.pptx"
)
OUTPUT_DIR = os.path.join(
    "C:\\Users\\33612\\Documents\\AXOLOTL",
    "Coco Can't Sleep Tonight!"
)
EPUB_FILENAME = "Coco-ne-dort-pas-ce-soir.epub"

TITLE = "Coco ne dort pas ce soir !"
SUBTITLE = "Une aventure de Coco l'Axolotl"
AUTHOR = "Dr. Anita NIRVENA"
ILLUSTRATOR = "Loopinky"
LANGUAGE = "fr"
PUBLISHER = "cocotheaxolotl.org"
DATE = "2026"
DESCRIPTION = "Un livre pour enfants sur le sommeil, avec Coco l'Axolotl."

# --- FONTS ---
FONT_DIR = "C:/Windows/Fonts/"
FONT_REGULAR = ImageFont.truetype(os.path.join(FONT_DIR, "georgia.ttf"), 52)
FONT_BOLD = ImageFont.truetype(os.path.join(FONT_DIR, "georgiab.ttf"), 52)
FONT_ITALIC = ImageFont.truetype(os.path.join(FONT_DIR, "georgiai.ttf"), 50)
FONT_TITLE = ImageFont.truetype(os.path.join(FONT_DIR, "georgiab.ttf"), 68)
FONT_H1 = ImageFont.truetype(os.path.join(FONT_DIR, "georgiab.ttf"), 78)
FONT_SMALL = ImageFont.truetype(os.path.join(FONT_DIR, "georgia.ttf"), 42)
FONT_FACT_TITLE = ImageFont.truetype(os.path.join(FONT_DIR, "georgiab.ttf"), 46)
FONT_FACT = ImageFont.truetype(os.path.join(FONT_DIR, "georgia.ttf"), 44)
FONT_COPYRIGHT = ImageFont.truetype(os.path.join(FONT_DIR, "georgia.ttf"), 34)

# Colors
BLACK = (0, 0, 0)
YELLOW = (0, 0, 0)
ORANGE_LIGHT = (0, 0, 0)
FACT_BG = (254, 249, 231, 220)
FACT_BORDER = (247, 201, 72)
FACT_TITLE_COLOR = (230, 126, 34)
FACT_TEXT_COLOR = (51, 51, 51)

# Page size (8.5x11 at 150 DPI)
PAGE_W = 1275
PAGE_H = 1650

# --- LOAD PPTX ---
print(f"Loading: {PPTX_PATH}")
prs = Presentation(PPTX_PATH)
print(f"Slides: {len(prs.slides)}")

def get_slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts

def get_slide_images(slide):
    images = []
    for shape in slide.shapes:
        try:
            if shape.image:
                images.append((shape.image.content_type, shape.image.blob))
        except:
            pass
    return images

def strip_emoji(text):
    return re.sub(r'^[\U0001f300-\U0001f9ff\U00002702-\U000027b0\U0000fe0f\u200d]+\s*', '', text)

# Extract all slide data
slides_data = []
for i, slide in enumerate(prs.slides):
    imgs = get_slide_images(slide)
    slides_data.append({
        'index': i + 1,
        'texts': get_slide_texts(slide),
        'images': imgs,
    })

# Get background image blob (same on all text slides)
bg_blob = slides_data[2]['images'][0][1]  # slide 3, first image
bg_img = Image.open(io.BytesIO(bg_blob)).convert('RGBA')
print(f"Background image: {bg_img.size}")

# --- IMAGE RENDERING HELPERS ---

def wrap_text(text, font, max_width, draw):
    """Word-wrap text to fit max_width."""
    # Coller la ponctuation française au mot précédent (évite "!" seul sur une ligne)
    text = re.sub(r'\s+([!?;:])', '\u00a0\\1', text)
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines

def draw_centered_text(draw, text, font, color, y, page_w, max_width=None):
    """Draw centered text, return new y position."""
    if max_width is None:
        max_width = page_w - 120
    lines = wrap_text(text, font, max_width, draw)
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        x = (page_w - tw) // 2
        draw.text((x, y), line, font=font, fill=color)
        y += th + 10
    return y + 5

def draw_fun_fact_box(draw, img, title_text, fact_text, y, page_w, max_width):
    """Draw a fun fact box with background."""
    margin = 60
    box_x = margin
    box_w = page_w - 2 * margin
    padding = 20

    # Pre-calculate text height
    title_lines = wrap_text(title_text, FONT_FACT_TITLE, box_w - 2 * padding, draw)
    fact_lines = wrap_text(fact_text, FONT_FACT, box_w - 2 * padding, draw)

    title_h = sum(draw.textbbox((0,0), l, font=FONT_FACT_TITLE)[3] - draw.textbbox((0,0), l, font=FONT_FACT_TITLE)[1] + 8 for l in title_lines)
    fact_h = sum(draw.textbbox((0,0), l, font=FONT_FACT)[3] - draw.textbbox((0,0), l, font=FONT_FACT)[1] + 8 for l in fact_lines)

    box_h = padding * 2 + title_h + fact_h + 10

    # Draw rounded rectangle background
    overlay = Image.new('RGBA', img.size, (0, 0, 0, 0))
    overlay_draw = ImageDraw.Draw(overlay)
    overlay_draw.rounded_rectangle(
        [box_x, y, box_x + box_w, y + box_h],
        radius=15,
        fill=FACT_BG,
        outline=FACT_BORDER,
        width=3
    )
    img.alpha_composite(overlay)

    # Re-get draw on the composite
    draw = ImageDraw.Draw(img)

    # Draw title
    ty = y + padding
    for line in title_lines:
        bbox = draw.textbbox((0, 0), line, font=FONT_FACT_TITLE)
        tw = bbox[2] - bbox[0]
        x = (page_w - tw) // 2
        draw.text((x, ty), line, font=FONT_FACT_TITLE, fill=FACT_TITLE_COLOR)
        ty += bbox[3] - bbox[1] + 8

    ty += 5
    # Draw fact text
    for line in fact_lines:
        bbox = draw.textbbox((0, 0), line, font=FONT_FACT)
        tw = bbox[2] - bbox[0]
        x = (page_w - tw) // 2
        draw.text((x, ty), line, font=FONT_FACT, fill=FACT_TEXT_COLOR)
        ty += bbox[3] - bbox[1] + 8

    return y + box_h + 20, draw

def _draw_page_content(draw, img, ch_title, texts, is_title_page, y_start, max_w):
    """Draw all text content starting at y_start. Returns final y."""
    y = y_start
    if is_title_page:
        y = draw_centered_text(draw, TITLE, FONT_H1, YELLOW, y, PAGE_W, max_w)
        y += 10
        y = draw_centered_text(draw, SUBTITLE, FONT_ITALIC, ORANGE_LIGHT, y, PAGE_W, max_w)
        y += 30
        y = draw_centered_text(draw, f"Par : {AUTHOR}", FONT_SMALL, BLACK, y, PAGE_W, max_w)
        y = draw_centered_text(draw, f"Illustrations : {ILLUSTRATOR}", FONT_SMALL, BLACK, y, PAGE_W, max_w)
        y += 40
        y = draw_centered_text(draw, f"\u00a9 {DATE} {AUTHOR}", FONT_COPYRIGHT, (180, 180, 180), y, PAGE_W, max_w)
        y = draw_centered_text(draw, "Tous droits r\u00e9serv\u00e9s.", FONT_COPYRIGHT, (180, 180, 180), y, PAGE_W, max_w)
        y = draw_centered_text(draw, PUBLISHER, FONT_COPYRIGHT, (180, 180, 180), y, PAGE_W, max_w)
    else:
        y = draw_centered_text(draw, ch_title, FONT_TITLE, YELLOW, y, PAGE_W, max_w)
        y += 20
        for text in texts:
            cleaned = strip_emoji(text)
            if not cleaned:
                continue
            if cleaned == ch_title or cleaned.strip() == ch_title.strip():
                continue
            if 'Le savais-tu' in cleaned:
                y += 80  # 2 lignes d'espace avant la boîte
                parts = cleaned.split('?', 1)
                fact_title = "Le savais-tu ?"
                fact_body = parts[1].strip() if len(parts) == 2 else cleaned
                y, draw = draw_fun_fact_box(draw, img, fact_title, fact_body, y, PAGE_W, max_w)
            else:
                for line in cleaned.replace('\x0b', '\n').split('\n'):
                    line = line.strip()
                    if not line:
                        continue
                    if line.startswith('\u00ab'):
                        y = draw_centered_text(draw, line, FONT_ITALIC, ORANGE_LIGHT, y, PAGE_W, max_w)
                    else:
                        y = draw_centered_text(draw, line, FONT_REGULAR, BLACK, y, PAGE_W, max_w)
    return y

def render_text_page(ch_title, texts, is_title_page=False):
    """Render a text page: starry background + text centered vertically."""
    max_w = PAGE_W - 120

    # First pass: measure total content height on a temp image
    tmp = Image.new('RGBA', (PAGE_W, PAGE_H * 2), (0, 0, 0, 0))
    tmp_draw = ImageDraw.Draw(tmp)
    end_y = _draw_page_content(tmp_draw, tmp, ch_title, texts, is_title_page, 0, max_w)
    content_h = end_y

    # Calculate starting y to center vertically
    y_start = max(40, (PAGE_H - content_h) // 2)

    # Second pass: draw for real on the starry background
    img = bg_img.copy()
    img = img.resize((PAGE_W, PAGE_H), Image.LANCZOS)
    draw = ImageDraw.Draw(img)
    _draw_page_content(draw, img, ch_title, texts, is_title_page, y_start, max_w)

    return img.convert('RGB')

# --- CHAPTER DEFINITIONS ---
chapter_defs = [
    {"type": "cover", "title": "Couverture", "slides": [1]},
    {"type": "title_page", "title": "Page de titre", "slides": [2]},
    {"type": "spread", "title": "L'heure de dormir", "text_slide": 3, "illust_slide": 4},
    {"type": "spread", "title": "L'aventure commence", "text_slide": 5, "illust_slide": 6},
    {"type": "spread", "title": "Le Hibou", "text_slide": 7, "illust_slide": 8},
    {"type": "spread", "title": "Le Dauphin", "text_slide": 9, "illust_slide": 10},
    {"type": "spread", "title": "La Loutre", "text_slide": 11, "illust_slide": 12},
    {"type": "spread", "title": "Le Koala", "text_slide": 13, "illust_slide": 14},
    {"type": "spread", "title": "La Chauve-Souris", "text_slide": 15, "illust_slide": 16},
    {"type": "spread", "title": "Le Chat", "text_slide": 17, "illust_slide": 18},
    {"type": "spread", "title": "Le Cheval", "text_slide": 19, "illust_slide": 20},
    {"type": "spread", "title": "Le Flamant Rose", "text_slide": 21, "illust_slide": 22},
    {"type": "spread", "title": "Le Paresseux", "text_slide": 23, "illust_slide": 24},
    {"type": "spread", "title": "Maman Axolotl", "text_slide": 25, "illust_slide": 26},
    {"type": "spread", "title": "Le super-pouvoir", "text_slide": 27, "illust_slide": 28},
    {"type": "spread", "title": "Bonne nuit, Coco", "text_slide": 29, "illust_slide": 30},
]

# --- GENERATE ALL IMAGES ---
all_images = {}  # filename -> bytes

def add_pil_image(pil_img, name):
    buf = io.BytesIO()
    pil_img.save(buf, format='PNG', optimize=True)
    blob = buf.getvalue()
    all_images[name] = blob
    return name

def add_raw_image(content_type, blob):
    h = hashlib.md5(blob).hexdigest()[:10]
    ext = '.png' if 'png' in content_type else '.jpg'
    fname = f"img_{h}{ext}"
    all_images[fname] = blob
    return fname

def escape_html(text):
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

# Generate pages
xhtml_files = []  # (filename, content, title, for_toc)

print("Generating pages...")

for ch_i, ch_def in enumerate(chapter_defs):
    if ch_def['type'] == 'cover':
        # Cover: use first image from slide 1
        ct, blob = slides_data[0]['images'][0]
        img_name = add_raw_image(ct, blob)
        xhtml = f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Couverture</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body><div class="full-page"><img src="images/{img_name}" alt="Couverture"/></div></body>
</html>'''
        xhtml_files.append((f"p{ch_i:02d}_cover.xhtml", xhtml, "Couverture", True))
        print(f"  Cover: {img_name}")

    elif ch_def['type'] == 'title_page':
        # Title page: render text on starry background
        rendered = render_text_page("", [], is_title_page=True)
        img_name = add_pil_image(rendered, "page_title.png")
        xhtml = f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Page de titre</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body><div class="full-page"><img src="images/{img_name}" alt="Page de titre"/></div></body>
</html>'''
        xhtml_files.append((f"p{ch_i:02d}_title.xhtml", xhtml, "Page de titre", True))
        print(f"  Title page: {img_name}")

    elif ch_def['type'] == 'spread':
        title = ch_def['title']
        # Text page: render composite image
        sd = slides_data[ch_def['text_slide'] - 1]
        rendered = render_text_page(title, sd['texts'])
        img_name = add_pil_image(rendered, f"page_{ch_i:02d}_text.png")
        xhtml = f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>{escape_html(title)}</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body><div class="full-page"><img src="images/{img_name}" alt="{escape_html(title)}"/></div></body>
</html>'''
        xhtml_files.append((f"p{ch_i:02d}_text.xhtml", xhtml, title, True))

        # Illustration page
        illust_sd = slides_data[ch_def['illust_slide'] - 1]
        if illust_sd['images']:
            ct, blob = illust_sd['images'][0]
            ill_name = add_raw_image(ct, blob)
            xhtml2 = f'''<?xml version="1.0" encoding="UTF-8"?>
<!DOCTYPE html>
<html xmlns="http://www.w3.org/1999/xhtml" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>{escape_html(title)} - Illustration</title><link rel="stylesheet" type="text/css" href="style.css"/></head>
<body><div class="full-page"><img src="images/{ill_name}" alt="{escape_html(title)}"/></div></body>
</html>'''
            xhtml_files.append((f"p{ch_i:02d}_illust.xhtml", xhtml2, f"{title} - Illustration", False))

        print(f"  {title}: text={img_name}, illust={'yes' if illust_sd['images'] else 'no'}")

# --- CSS (minimal - just full page images) ---
CSS = '''@charset "UTF-8";
body { margin: 0; padding: 0; }
.full-page { text-align: center; page-break-after: always; }
.full-page img { max-width: 100%; height: auto; display: block; margin: 0 auto; }
'''

# --- BUILD EPUB ---
book_uuid = str(uuid.uuid4())

def build_opf():
    manifest = ['<item id="style" href="style.css" media-type="text/css"/>']
    spine = []
    for i, (fname, _, _, _) in enumerate(xhtml_files):
        mid = f"p{i:02d}"
        manifest.append(f'<item id="{mid}" href="{fname}" media-type="application/xhtml+xml"/>')
        spine.append(f'<itemref idref="{mid}"/>')
    for img_name in sorted(all_images.keys()):
        iid = img_name.replace('.', '_').replace('-', '_')
        mt = 'image/png' if img_name.endswith('.png') else 'image/jpeg'
        manifest.append(f'<item id="{iid}" href="images/{img_name}" media-type="{mt}"/>')
    manifest.append('<item id="ncx" href="toc.ncx" media-type="application/x-dtbncx+xml"/>')
    manifest.append('<item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>')
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="bookid" version="3.0">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:identifier id="bookid">urn:uuid:{book_uuid}</dc:identifier>
    <dc:title>{escape_html(TITLE)}</dc:title>
    <dc:creator>{escape_html(AUTHOR)}</dc:creator>
    <dc:language>{LANGUAGE}</dc:language>
    <dc:publisher>{escape_html(PUBLISHER)}</dc:publisher>
    <dc:date>{DATE}</dc:date>
    <dc:description>{escape_html(DESCRIPTION)}</dc:description>
    <meta property="dcterms:modified">2026-02-19T00:00:00Z</meta>
  </metadata>
  <manifest>
    {chr(10).join("    " + m for m in manifest)}
  </manifest>
  <spine toc="ncx">
    {chr(10).join("    " + s for s in spine)}
  </spine>
</package>'''

def build_ncx():
    pts = []
    n = 0
    for fname, _, title, for_toc in xhtml_files:
        if for_toc:
            n += 1
            pts.append(f'    <navPoint id="np{n}" playOrder="{n}"><navLabel><text>{escape_html(title)}</text></navLabel><content src="{fname}"/></navPoint>')
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">
  <head><meta name="dtb:uid" content="urn:uuid:{book_uuid}"/></head>
  <docTitle><text>{escape_html(TITLE)}</text></docTitle>
  <navMap>
{chr(10).join(pts)}
  </navMap>
</ncx>'''

def build_nav():
    items = []
    for fname, _, title, for_toc in xhtml_files:
        if for_toc:
            items.append(f'      <li><a href="{fname}">{escape_html(title)}</a></li>')
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops" xml:lang="fr" lang="fr">
<head><meta charset="UTF-8"/><title>Table des matieres</title></head>
<body><nav epub:type="toc" id="toc"><h1>Table des matieres</h1><ol>
{chr(10).join(items)}
</ol></nav></body></html>'''

# --- WRITE EPUB ---
epub_path = os.path.join(OUTPUT_DIR, EPUB_FILENAME)
print(f"Writing EPUB to: {epub_path}")

with zipfile.ZipFile(epub_path, 'w', zipfile.ZIP_DEFLATED) as zf:
    zf.writestr('mimetype', 'application/epub+zip', compress_type=zipfile.ZIP_STORED)
    zf.writestr('META-INF/container.xml', '''<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles>
</container>''')
    zf.writestr('OEBPS/content.opf', build_opf())
    zf.writestr('OEBPS/toc.ncx', build_ncx())
    zf.writestr('OEBPS/nav.xhtml', build_nav())
    zf.writestr('OEBPS/style.css', CSS)
    for fname, content, _, _ in xhtml_files:
        zf.writestr(f'OEBPS/{fname}', content)
    for img_name, img_blob in all_images.items():
        zf.writestr(f'OEBPS/images/{img_name}', img_blob)

print(f"\nDone!")
print(f"  - {len(xhtml_files)} pages")
print(f"  - {len(all_images)} images")
print(f"  - File size: {os.path.getsize(epub_path) / 1024 / 1024:.1f} MB")
