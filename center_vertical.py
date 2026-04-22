"""
Center ALL elements as a group vertically on each slide.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu

path = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
prs = Presentation(path)
slide_h = prs.slide_height

fixes = 0
for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    if slide_num < 4 or slide_num > 84:
        continue

    shapes = list(slide.shapes)
    if not shapes:
        continue

    min_top = None
    max_bottom = None
    for shape in shapes:
        top = shape.top
        bottom = shape.top + shape.height
        if min_top is None or top < min_top:
            min_top = top
        if max_bottom is None or bottom > max_bottom:
            max_bottom = bottom

    if min_top is None:
        continue

    group_height = max_bottom - min_top
    current_center = min_top + group_height // 2
    target_center = slide_h // 2
    offset = target_center - current_center

    if abs(offset) < Emu(10000):
        continue

    for shape in shapes:
        shape.top = shape.top + offset

    fixes += 1
    offset_in = offset / 914400
    print(f"  Slide {slide_num}: shifted {offset_in:+.2f}\" ({len(shapes)} elements)")

prs.save(path)
print(f"\nDone! Centered {fixes} slides vertically.")
