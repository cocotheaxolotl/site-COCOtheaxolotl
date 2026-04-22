import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

path = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
out_dir = r"C:\Users\33612\Documents\coloring book for adults\extracted_images"
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(path)

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    # Skip intro slides (1-3) and end promo slides (108+)
    if slide_num < 4 or slide_num > 107:
        continue

    for shape in slide.shapes:
        try:
            if shape.image:
                blob = shape.image.blob
                size_kb = len(blob) // 1024
                # Skip the small 33KB watermark/placeholder
                if size_kb < 100:
                    continue
                ext = '.png' if shape.image.content_type == 'image/png' else '.jpg'
                fname = f"slide_{slide_num:03d}{ext}"
                fpath = os.path.join(out_dir, fname)
                with open(fpath, 'wb') as f:
                    f.write(blob)
                print(f"Slide {slide_num:3d}: saved {fname} ({size_kb}KB)")
        except:
            pass

print(f"\nDone! Images saved to: {out_dir}")
print(f"Total files: {len(os.listdir(out_dir))}")
