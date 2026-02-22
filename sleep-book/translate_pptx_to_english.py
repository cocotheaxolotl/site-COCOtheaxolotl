"""
Translate the French PPTX "Coco ne dort pas ce soir" to English.
Creates a new PPTX with translated text (preserving all formatting),
then exports slides as JPG via PowerPoint COM.

Usage:
  python translate_pptx_to_english.py              # Full pipeline: translate + export JPG
  python translate_pptx_to_english.py --translate   # Only create translated PPTX
  python translate_pptx_to_english.py --export      # Only export existing EN PPTX to JPG
  python translate_pptx_to_english.py --extract      # Only print extracted text (dry run)
"""

import sys, os, shutil, argparse
from copy import deepcopy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FR_PPTX = os.path.join(SCRIPT_DIR, 'Coco-ne-dort-pas-ce-soir-FR-sleep-book.pptx')
EN_PPTX = os.path.join(SCRIPT_DIR, 'Coco-doesnt-sleep-tonight-EN-sleep-book.pptx')
EN_JPG_DIR = os.path.join(SCRIPT_DIR, 'Coco-doesnt-sleep-tonight-EN-sleep-book')

# ── Translation Dictionary ──
# Keys = normalized French paragraph text, Values = English translation
TRANSLATIONS = {
    # === SLIDE 1: Cover ===
    "histoire du soir":
        "bedtime story",
    "COCO":
        "COCO",
    "NE DORT PAS":
        "DOESN'T SLEEP",
    "CE SOIR":
        "TONIGHT",

    # === SLIDE 2: Bedtime ===
    "Ce soir, c\u2019\u00e9tait l\u2019heure de dormir. Mais Coco n\u2019avait pas du tout sommeil.\u2014 Maman, je n\u2019arrive pas \u00e0 dormir !\u2014 Compte les moutons et...":
        "That night, it was time for bed. But Coco wasn\u2019t sleepy at all.\n\u2014 Mom, I can\u2019t fall asleep!\n\u2014 Count some sheep and...",
    "\u2014 Mais Maman, les moutons ne savent pas nager !Maman ria.":
        "\u2014 But Mom, sheep can\u2019t swim!\nMom laughed.",
    "\u2014 C\u2019est vrai ! Alors reste tranquille et le sommeil viendra.":
        "\u2014 That\u2019s true! Just stay still and sleep will come.",
    "\U0001f319 L\u2019heure de dormir":
        "\U0001f319 Bedtime",

    # === SLIDE 4: The Adventure Begins ===
    "Mais Coco ne pouvait pas rester tranquille. Pas ce soir.Quand Maman et Papa se furent endormis, Coco se faufila doucement hors de la rivi\u00e8re.\u2014 Si je ne dors pas, je vais d\u00e9couvrir ce que font les autres animaux la nuit !":
        "But Coco couldn\u2019t stay still. Not tonight.\nWhen Mom and Dad fell asleep, Coco quietly slipped out of the river.\n\u2014 If I can\u2019t sleep, I\u2019ll find out what other animals do at night!",
    "\u2728 L\u2019aventure commence":
        "\u2728 The Adventure Begins",

    # === SLIDE 6: The Owl ===
    "\U0001f989 Le Hibou":
        "\U0001f989 The Owl",
    "Coco rencontra un hibou perch\u00e9 sur une branche.":
        "Coco met an owl perched on a branch.",
    "\u2014 Bonsoir, Hibou ! Pourquoi tu ne dors pas ?":
        "\u2014 Good evening, Owl! Why aren\u2019t you sleeping?",
    "\u2014 Hou hou ! Je suis nocturne ! Je travaille la nuit et je dors le jour. Mais m\u00eame moi, j\u2019ai besoin de dormir. Tout le monde a besoin de dormir !":
        "\u2014 Hoo hoo! I\u2019m nocturnal! I work at night and sleep during the day. But even I need to sleep. Everyone needs to sleep!",
    "Les hiboux dorment le jour, cach\u00e9s dans les arbres. Ils peuvent tourner la t\u00eate presque tout autour, comme pour regarder derri\u00e8re eux !":
        "Owls sleep during the day, hidden in trees. They can turn their heads almost all the way around, as if looking behind them!",

    # === SLIDE 8: The Dolphin ===
    "\U0001f42c Le Dauphin":
        "\U0001f42c The Dolphin",
    "Au bord de l\u2019oc\u00e9an, un dauphin nageait tranquillement.\u2014 Tu ne dors pas non plus ?\u2014 Oh si ! Je n\u2019endors qu\u2019une moiti\u00e9 de mon cerveau \u00e0 la fois. L\u2019autre reste \u00e9veill\u00e9e pour nager et respirer !\u2014 Tu dors ET tu nages en m\u00eame temps ?!\u2014 Le sommeil est si important que j\u2019ai trouv\u00e9 un moyen de ne jamais m\u2019en passer !":
        "By the ocean, a dolphin was swimming peacefully.\n\u2014 You\u2019re not sleeping either?\n\u2014 Oh, I am! I only put half my brain to sleep at a time. The other half stays awake to swim and breathe!\n\u2014 You sleep AND swim at the same time?!\n\u2014 Sleep is so important that I found a way to never go without it!",
    "Les dauphins dorment avec un seul \u0153il ferm\u00e9 ! Le c\u00f4t\u00e9 du cerveau qui dort change toutes les deux heures.":
        "Dolphins sleep with only one eye closed! The side of the brain that sleeps switches every two hours.",

    # === SLIDE 10: The Otter ===
    "\U0001f9a6 La Loutre":
        "\U0001f9a6 The Otter",
    "Deux loutres flottaient sur le dos, patte dans la patte.\u2014 Pourquoi vous vous tenez la main ?\u2014 Pour ne pas s\u2019\u00e9loigner en dormant ! Se sentir en s\u00e9curit\u00e9, c\u2019est tr\u00e8s important pour bien dormir.Coco sourit.":
        "Two otters were floating on their backs, paw in paw.\n\u2014 Why are you holding hands?\n\u2014 So we don\u2019t drift apart while sleeping! Feeling safe is very important for a good night\u2019s sleep.\nCoco smiled.",
    "\u2014 C\u2019est comme quand je me blottis contre ma peluche !":
        "\u2014 It\u2019s like when I snuggle up with my stuffed animal!",
    "Les loutres de mer se tiennent vraiment par la main en dormant ! Elles s\u2019enroulent aussi dans des algues pour ne pas d\u00e9river.":
        "Sea otters really do hold hands while sleeping! They also wrap themselves in seaweed so they don\u2019t drift away.",

    # === SLIDE 12: The Koala ===
    "\U0001f428 Le Koala":
        "\U0001f428 The Koala",
    "Dans un eucalyptus, un koala dormait profond\u00e9ment.\u2014 Koala ! Tu dors ?\u2014 Zzz... oui... 22 heures par jour...\u2014 22 HEURES ?! Mais pourquoi autant ?\u2014 Mon corps a besoin d\u2019\u00e9nergie pour dig\u00e9rer mes feuilles... Le sommeil, \u00e7a recharge ton corps en \u00e9nergie !":
        "In a eucalyptus tree, a koala was sleeping soundly.\n\u2014 Koala! Are you sleeping?\n\u2014 Zzz... yes... 22 hours a day...\n\u2014 22 HOURS?! But why so much?\n\u2014 My body needs energy to digest my leaves... Sleep recharges your body with energy!",
    "Et il se rendormit aussit\u00f4t.":
        "And he fell right back asleep.",
    "Les koalas dorment jusqu\u2019\u00e0 22 heures par jour ! Les feuilles d\u2019eucalyptus sont tr\u00e8s dures \u00e0 dig\u00e9rer et donnent tr\u00e8s peu d\u2019\u00e9nergie.":
        "Koalas sleep up to 22 hours a day! Eucalyptus leaves are very hard to digest and provide very little energy.",

    # === SLIDE 14: The Bat ===
    "\U0001f987 La Chauve-Souris":
        "\U0001f987 The Bat",
    "Sous un pont, des chauves-souris \u00e9taient suspendues la t\u00eate en bas.\u2014 Vous dormez \u00e0 l\u2019envers ?!\u2014 Bien s\u00fbr ! Et pendant qu\u2019on dort, notre cerveau trie tout ce qu\u2019on a appris \u2014 o\u00f9 sont les insectes, quels chemins \u00e9viter...\u2014 C\u2019est pour \u00e7a qu\u2019on se sent plus malin apr\u00e8s une bonne nuit !":
        "Under a bridge, bats were hanging upside down.\n\u2014 You sleep upside down?!\n\u2014 Of course! And while we sleep, our brain sorts everything we\u2019ve learned \u2014 where the insects are, which paths to avoid...\n\u2014 That\u2019s why you feel smarter after a good night\u2019s sleep!",
    "Pendant le sommeil, le cerveau trie et range les souvenirs. C\u2019est pour \u00e7a qu\u2019on retient mieux ses le\u00e7ons apr\u00e8s avoir bien dormi !":
        "During sleep, the brain sorts and organizes memories. That\u2019s why you remember your lessons better after a good night\u2019s sleep!",

    # === SLIDE 16: The Cat ===
    "\U0001f431 Le Chat":
        "\U0001f431 The Cat",
    "Pr\u00e8s d\u2019une maison, un chat \u00e9tait roul\u00e9 en boule sur un muret.\u2014 Tu fais un gros dodo ?\u2014 Je fais plein de petits dodos. On appelle \u00e7a des micro-siestes !\u2014 Pour les petits axolotls, c\u2019est mieux de faire un long sommeil la nuit. Pendant le sommeil profond, ton corps se r\u00e9pare !":
        "Near a house, a cat was curled up on a low wall.\n\u2014 Are you having a big nap?\n\u2014 I take lots of little naps. They\u2019re called catnaps!\n\u2014 For little axolotls, it\u2019s better to have one long sleep at night. During deep sleep, your body repairs itself!",
    "Les chats dorment entre 12 et 16 heures par jour ! C\u2019est de l\u00e0 que vient l\u2019expression micro-sieste.":
        "Cats sleep between 12 and 16 hours a day! That\u2019s where the expression \u201ccatnap\u201d comes from.",

    # === SLIDE 18: The Horse ===
    "\U0001f40e Le Cheval":
        "\U0001f40e The Horse",
    "Dans un pr\u00e9, un cheval se tenait immobile sous la lune.\u2014 Tu... tu dors debout ?\u2014 Oui ! Comme \u00e7a je suis pr\u00eat \u00e0 courir. Mais pour r\u00eaver, je dois m\u2019allonger.\u2014 Le sommeil a diff\u00e9rentes \u00e9tapes : d\u2019abord l\u00e9ger, puis profond, puis les r\u00eaves ! Et \u00e7a recommence toute la nuit.":
        "In a meadow, a horse stood still under the moon.\n\u2014 You... you sleep standing up?\n\u2014 Yes! That way I\u2019m ready to run. But to dream, I have to lie down.\n\u2014 Sleep has different stages: first light, then deep, then dreams! And it starts all over again throughout the night.",
    "Les chevaux ne dorment profond\u00e9ment que 2 \u00e0 3 heures par jour \u2014 et seulement allong\u00e9s !":
        "Horses only sleep deeply for 2 to 3 hours a day \u2014 and only when lying down!",

    # === SLIDE 20: The Flamingo ===
    "\U0001f9a9 Le Flamant Rose":
        "\U0001f9a9 The Flamingo",
    "Au bord d\u2019un \u00e9tang, un flamant rose dormait sur une seule patte !\u2014 Comment tu ne tombes pas ?!\u2014 L\u2019\u00e9quilibre, \u00e7a demande un corps bien repos\u00e9 ! Quand je suis fatigu\u00e9...Le flamant vacilla.":
        "By a pond, a flamingo was sleeping on one leg!\n\u2014 How don\u2019t you fall over?!\n\u2014 Balance requires a well-rested body! When I\u2019m tired...\nThe flamingo wobbled.",
    "\u2014 ...Voil\u00e0 ce qui arrive !Coco rigola, mais commen\u00e7ait \u00e0 comprendre.":
        "\u2014 ...That\u2019s what happens!\nCoco laughed, but was starting to understand.",
    "Les flamants dorment souvent sur une seule patte ! Les scientifiques pensent que c\u2019est pour garder l\u2019autre au chaud.":
        "Flamingos often sleep on one leg! Scientists think it\u2019s to keep the other one warm.",

    # === SLIDE 22: The Sloth ===
    "\U0001f9a5 Le Paresseux":
        "\U0001f9a5 The Sloth",
    "Coco commen\u00e7ait \u00e0 b\u00e2iller. Un paresseux pendait d\u2019une branche.\u2014 Je commence... \u00e0 \u00eatre... fatigu\u00e9...\u2014 Bienvenuuuue... au cluuuub... Quand on ne dort pas assez... on devient lent... grognon... et tout est plus difficile...\u2014 Allez... va te coucher... petit axolotl... zzz...":
        "Coco was starting to yawn. A sloth was hanging from a branch.\n\u2014 I\u2019m starting... to feel... tired...\n\u2014 Welcomeee... to the cluuub... When you don\u2019t sleep enough... you become slow... grumpy... and everything is harder...\n\u2014 Go on... go to bed... little axolotl... zzz...",
    "Les paresseux dorment environ 15 heures par jour ! Sans assez de sommeil, notre corps se sent lourd et notre cerveau a du mal \u00e0 se concentrer.":
        "Sloths sleep about 15 hours a day! Without enough sleep, our body feels heavy and our brain has trouble concentrating.",

    # === SLIDE 24: Mom Axolotl ===
    "\U0001f495 Maman Axolotl":
        "\U0001f495 Mama Axolotl",
    "\u2014 Coco ?C\u2019\u00e9tait Maman, au bord de la rivi\u00e8re.\u2014 Maman ! J\u2019ai d\u00e9couvert comment tous les animaux dorment ! Tout le monde en a besoin. Moi aussi !\u2014 Et sais-tu comment dorment les axolotls ? On se pose au fond de l\u2019eau, bien tranquilles, berc\u00e9s par le courant...":
        "\u2014 Coco?\nIt was Mom, by the river.\n\u2014 Mom! I discovered how all the animals sleep! Everyone needs it. Me too!\n\u2014 And do you know how axolotls sleep? We settle at the bottom of the water, nice and quiet, rocked by the current...",
    "Les axolotls vivent dans l\u2019eau douce et respirent gr\u00e2ce \u00e0 leurs jolies branchies en forme de plume sur leur t\u00eate.":
        "Axolotls live in fresh water and breathe through their beautiful feather-shaped gills on their heads.",

    # === SLIDE 26: The Superpower ===
    "\u26a1Le super-pouvoir":
        "\u26a1The Superpower",
    "Maman porta Coco jusqu\u2019\u00e0 son coin douillet au fond de la rivi\u00e8re.\u2014 Le sommeil, \u00e7a recharge l\u2019\u00e9nergie, \u00e7a range le cerveau, \u00e7a r\u00e9pare le corps... pas vrai ?\u2014 Tout \u00e7a \u00e0 la fois, mon Coco. Et nous, les axolotls, on peut m\u00eame faire repousser nos pattes ! Mais pour \u00e7a, il faut bien dormir.":
        "Mom carried Coco to his cozy spot at the bottom of the river.\n\u2014 Sleep recharges your energy, organizes your brain, repairs your body... right?\n\u2014 All of that at once, my Coco. And we axolotls can even regrow our legs! But for that, you need to sleep well.",
    "Les axolotls ont un super-pouvoir : ils peuvent r\u00e9g\u00e9n\u00e9rer leurs pattes, leur queue et m\u00eame des parties de leur cerveau !":
        "Axolotls have a superpower: they can regenerate their legs, their tail, and even parts of their brain!",

    # === SLIDE 28: Goodnight ===
    "\U0001f4a4 Bonne nuit Coco !":
        "\U0001f4a4 Good Night Coco!",
    "Coco se blottit dans son petit nid d\u2019eau douce avec sa peluche.Il pensa \u00e0 tous les animaux qui dormaient en ce moment \u2014 chacun \u00e0 sa fa\u00e7on, mais tous parce que le sommeil est magique.Ses branchies redevinrent toutes roses et touffues.Et Coco s\u2019endormit.":
        "Coco snuggled into his little freshwater nest with his stuffed animal.\nHe thought of all the animals sleeping right now \u2014 each in their own way, but all because sleep is magical.\nHis gills turned rosy and fluffy again.\nAnd Coco fell asleep.",
    "Bonne nuit !":
        "Good Night!",
    "Coco":
        "Coco",

    # === SLIDE 30: Website ===
    "Retrouve mon monde sur mon site :":
        "Discover my world on my website:",
    "cocotheaxolotl.org":
        "cocotheaxolotl.org",
    "\u00c0 bient\u00f4t !":
        "See You Soon!",

    # === SLIDE 31: Dear Parent ===
    "CHER PARENT":
        "DEAR PARENT",
    "Le sommeil est un ph\u00e9nom\u00e8ne universel pour tous les humains et les animaux. Cependant, expliquer ce qu\u2019est le sommeil et pourquoi il est essentiel, m\u00eame pour les adultes, n\u2019est pas simple. Ce livre vise \u00e0 aider les enfants et les parents \u00e0 en apprendre davantage sur le sommeil de mani\u00e8re ludique et engageante, \u00e0 travers l\u2019histoire de Coco, un petit axolotl curieux qui part \u00e0 la d\u00e9couverte du monde nocturne.":
        "Sleep is a universal phenomenon for all humans and animals. However, explaining what sleep is and why it is essential, even for adults, is not simple. This book aims to help children and parents learn more about sleep in a fun and engaging way, through the story of Coco, a curious little axolotl who sets out to discover the nocturnal world.",
    "A travers ses rencontres avec diff\u00e9rents animaux, Coco d\u00e9couvre que chacun dort d\u2019une fa\u00e7on unique et fascinante \u2014 et comprend peu \u00e0 peu pourquoi le sommeil est si important.":
        "Through his encounters with different animals, Coco discovers that each one sleeps in a unique and fascinating way \u2014 and gradually understands why sleep is so important.",
    "Les principes d\u00e9crits dans ce livre sont bas\u00e9s sur des ann\u00e9es de recherche dans le domaine du sommeil chez l\u2019enfant.":
        "The principles described in this book are based on years of research in the field of children\u2019s sleep.",
    "Dr. Anita NIRVENA":
        "Dr. Anita NIRVENA",
    "Coco ne dort pas ce soir !Une aventure de Coco l\u2019AxolotlPar : Dr. Anita NIRVENAIllustrations : LOOPINKY\u00a9 2026 Dr. Anita NIRVENATous droits r\u00e9serv\u00e9s.www.cocotheaxolotl.orgPremi\u00e8re \u00e9dition, 2026":
        "Coco Doesn\u2019t Sleep Tonight!\nA Coco the Axolotl Adventure\nBy: Dr. Anita NIRVENA\nIllustrations: LOOPINKY\n\u00a9 2026 Dr. Anita NIRVENA\nAll rights reserved.\nwww.cocotheaxolotl.org\nFirst edition, 2026",

    # === Common "Le savais-tu ?" header (appears on many slides) ===
    "Le savais-tu ?":
        "Did You Know?",
}


def normalize(text):
    """Normalize text for dictionary lookup: strip, remove vertical tabs, collapse whitespace."""
    # PPTX uses \x0b (vertical tab) for soft line breaks within a paragraph
    text = text.replace('\x0b', '')
    # Normalize curly quotes to straight for consistent matching
    text = text.replace('\u2019', "'").replace('\u2018', "'")
    text = text.replace('\u201c', '"').replace('\u201d', '"')
    return ' '.join(text.split()).strip()


def find_translation(fr_text):
    """Look up the French text in the translation dictionary."""
    key = normalize(fr_text)
    if key in TRANSLATIONS:
        return TRANSLATIONS[key]
    # Try without leading/trailing punctuation differences
    for k, v in TRANSLATIONS.items():
        if normalize(k) == key:
            return v
    return None


def replace_paragraph_text(para, new_text):
    """Replace all text in a paragraph with new_text, preserving first run's formatting."""
    if not para.runs:
        return

    # Save formatting from the first run
    first_run = para.runs[0]
    rPr = first_run._r.find(qn('a:rPr'))
    saved_rPr = deepcopy(rPr) if rPr is not None else None

    # Remove ALL content elements: runs (a:r), line breaks (a:br), AND endParaRPr
    # We'll re-add endParaRPr at the very end so runs come before it
    saved_endParaRPr = None
    for child in list(para._p):
        tag = child.tag
        if tag in (qn('a:r'), qn('a:br')):
            para._p.remove(child)
        elif tag == qn('a:endParaRPr'):
            saved_endParaRPr = deepcopy(child)
            para._p.remove(child)

    # Re-create content with <a:br/> for line breaks
    parts = new_text.split('\n')
    for idx, part in enumerate(parts):
        if idx > 0:
            # Insert a line break element <a:br/>
            br = etree.SubElement(para._p, qn('a:br'))
            if saved_rPr is not None:
                br_rPr = deepcopy(saved_rPr)
                br.insert(0, br_rPr)
        new_r = etree.SubElement(para._p, qn('a:r'))
        if saved_rPr is not None:
            new_r.insert(0, deepcopy(saved_rPr))
        new_t = etree.SubElement(new_r, qn('a:t'))
        new_t.text = part
        new_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

    # Re-add endParaRPr at the very end (must come after all runs)
    if saved_endParaRPr is not None:
        para._p.append(saved_endParaRPr)


def translate_pptx():
    """Copy FR PPTX and replace all text with English translations."""
    print(f"Copying {FR_PPTX} -> {EN_PPTX}")
    shutil.copy2(FR_PPTX, EN_PPTX)

    prs = Presentation(EN_PPTX)
    translated = 0
    not_found = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                fr_text = para.text.strip()
                if not fr_text:
                    continue

                en_text = find_translation(fr_text)
                if en_text is not None:
                    replace_paragraph_text(para, en_text)
                    translated += 1
                else:
                    not_found.append((slide_idx, fr_text[:80]))

    prs.save(EN_PPTX)
    print(f"\nTranslated {translated} paragraphs.")
    if not_found:
        print(f"\nWARNING: {len(not_found)} paragraphs not found in dictionary:")
        for slide_num, text in not_found:
            print(f"  Slide {slide_num}: {text}")
    else:
        print("All paragraphs translated successfully!")
    print(f"\nSaved: {EN_PPTX}")


def export_to_jpg():
    """Export EN PPTX slides as JPG using PowerPoint COM automation."""
    import comtypes.client

    pptx_path = os.path.abspath(EN_PPTX)
    output_dir = os.path.abspath(EN_JPG_DIR)

    if not os.path.exists(pptx_path):
        print(f"ERROR: {pptx_path} not found. Run --translate first.")
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    print(f"Opening PowerPoint...")
    pp = comtypes.client.CreateObject('PowerPoint.Application')
    pp.Visible = 1

    print(f"Opening {pptx_path}...")
    pres = pp.Presentations.Open(pptx_path, WithWindow=False)

    slide_count = pres.Slides.Count
    print(f"Exporting {slide_count} slides as JPG (840x1080)...")

    for i in range(1, slide_count + 1):
        slide = pres.Slides(i)
        jpg_path = os.path.join(output_dir, f'Diapositive{i}.JPG')
        slide.Export(jpg_path, 'JPG', 840, 1080)
        print(f"  Slide {i}/{slide_count}: {jpg_path}")

    pres.Close()
    pp.Quit()
    print(f"\nDone! JPGs saved in: {output_dir}")


def extract_text():
    """Print all text from the FR PPTX (dry run)."""
    prs = Presentation(FR_PPTX)
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            print(f"=== SLIDE {i} ===")
            for t in texts:
                en = find_translation(t)
                status = "OK" if en else "MISSING"
                print(f"  [{status}] {t[:100]}")
                if en:
                    print(f"     -> {en[:100]}")
            print()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Translate Coco PPTX to English')
    parser.add_argument('--translate', action='store_true', help='Only create translated PPTX')
    parser.add_argument('--export', action='store_true', help='Only export EN PPTX to JPG')
    parser.add_argument('--extract', action='store_true', help='Only print extracted text (dry run)')
    args = parser.parse_args()

    if args.extract:
        extract_text()
    elif args.translate:
        translate_pptx()
    elif args.export:
        export_to_jpg()
    else:
        # Full pipeline
        translate_pptx()
        export_to_jpg()
