"""
Fix z-order in wild portraits PPTX:
Move illustration images BEHIND text elements.
"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from lxml import etree

path = r"C:\Users\33612\Documents\coloring book for adults\wild portraits coloring book for adults.pptx"
prs = Presentation(path)

fixes = 0
for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    # Skip intro/outro slides
    if slide_num < 4 or slide_num > 84:
        continue

    # The spTree contains all shapes - order = z-order (first = back)
    spTree = slide.shapes._spTree

    # Find the large illustration image element and move it to the front (first child after background)
    big_img_elem = None
    for sp in list(spTree):
        # Check if this is a picture element with a large image
        if sp.tag.endswith('}pic'):
            # It's a picture shape - check size
            for shape in slide.shapes:
                if hasattr(shape, '_element') and shape._element is sp:
                    try:
                        if shape.image and len(shape.image.blob) > 100 * 1024:
                            big_img_elem = sp
                    except:
                        pass

    if big_img_elem is not None:
        # Remove and re-insert at the beginning (after the spTree properties)
        spTree.remove(big_img_elem)
        # Insert after the first child (which is usually grpSpPr or nvGrpSpPr)
        # Find the right position - after all the group properties
        insert_pos = 0
        for idx, child in enumerate(spTree):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('nvGrpSpPr', 'grpSpPr'):
                insert_pos = idx + 1
        spTree.insert(insert_pos, big_img_elem)
        fixes += 1
        if slide_num <= 40 or slide_num >= 69:  # Only print some for brevity
            print(f"  Slide {slide_num}: moved illustration to back")

prs.save(path)
print(f"\nDone! Fixed z-order on {fixes} slides.")
